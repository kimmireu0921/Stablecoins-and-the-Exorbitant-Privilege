"""
make_slides_final.py
Build all matplotlib figures, then assemble
presentations/0609_Stablecoin_Exorbitant_Privilege.pptx
"""

import os
import sys
import textwrap
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.patches as FancyArrow
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch
import matplotlib.patheffects as pe
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import copy

# ─────────────────────────────────────────────
# Paths
# ─────────────────────────────────────────────
BASE = Path("/Users/mimi/stablecoin_research")
RESULTS = BASE / "results"
DATA = BASE / "data"
PRESENTATIONS = BASE / "presentations"
PRESENTATIONS.mkdir(exist_ok=True)

# ─────────────────────────────────────────────
# Colors
# ─────────────────────────────────────────────
YONSEI_BLUE   = "#003087"
YONSEI_ORANGE = "#E87722"
GRAY          = "#666666"
LIGHT_GRAY    = "#F2F2F2"
RED           = "#C0392B"
GREEN         = "#1A7A4A"
WHITE         = "#FFFFFF"

# pptx RGBColor versions
C_BLUE   = RGBColor(0x00, 0x30, 0x87)
C_ORANGE = RGBColor(0xE8, 0x77, 0x22)
C_GRAY   = RGBColor(0x66, 0x66, 0x66)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)
C_BLACK  = RGBColor(0x22, 0x22, 0x22)
C_RED    = RGBColor(0xC0, 0x39, 0x2B)
C_GREEN  = RGBColor(0x1A, 0x7A, 0x4A)

# ─────────────────────────────────────────────
# Matplotlib style helpers
# ─────────────────────────────────────────────
def set_academic_style():
    plt.rcParams.update({
        "font.family":      "DejaVu Sans",
        "font.size":        11,
        "axes.spines.top":  False,
        "axes.spines.right":False,
        "axes.linewidth":   0.8,
        "axes.edgecolor":   GRAY,
        "xtick.labelsize":  10,
        "ytick.labelsize":  10,
        "figure.facecolor": "white",
        "axes.facecolor":   "white",
    })


# ═══════════════════════════════════════════════════════════════════
#  FIGURE 1 — Stablecoin flow (Slide 2)
# ═══════════════════════════════════════════════════════════════════
def make_fig_stablecoin_flow():
    set_academic_style()
    fig, ax = plt.subplots(figsize=(11, 5.5))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 6)
    ax.axis("off")

    # ── Box helper ──
    def box(ax, x, y, w, h, label, sub="", color=YONSEI_BLUE, fontsize=11, subsize=9):
        rect = FancyBboxPatch((x - w/2, y - h/2), w, h,
                              boxstyle="round,pad=0.12",
                              linewidth=1.5, edgecolor=color,
                              facecolor=color if color == YONSEI_BLUE else "#EFF4FF",
                              zorder=3)
        ax.add_patch(rect)
        fc = WHITE if color == YONSEI_BLUE else YONSEI_BLUE
        ax.text(x, y + (0.15 if sub else 0), label, ha="center", va="center",
                fontsize=fontsize, fontweight="bold", color=fc, zorder=4)
        if sub:
            ax.text(x, y - 0.28, sub, ha="center", va="center",
                    fontsize=subsize, color=fc, zorder=4)

    def arrow(ax, x1, y1, x2, y2, color=YONSEI_BLUE, lw=2, label="", lstyle="-"):
        ax.annotate("", xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle="-|>", color=color,
                                   lw=lw, linestyle=lstyle),
                    zorder=5)
        if label:
            mx, my = (x1+x2)/2, (y1+y2)/2
            ax.text(mx+0.08, my+0.15, label, fontsize=9, color=color, zorder=6)

    # ── NORMAL FLOW (top row) ──
    ax.text(5, 5.65, "Normal Flow — Supply Growth", ha="center", va="center",
            fontsize=12, fontweight="bold", color=YONSEI_BLUE)

    box(ax, 1.1, 4.5, 1.6, 0.75, "User / Investor", "deposits $100", color="#1A5DAD")
    box(ax, 3.7, 4.5, 2.0, 0.75, "Stablecoin Issuer", "Tether / Circle", color=YONSEI_BLUE)
    box(ax, 6.5, 4.5, 1.6, 0.75, "US T-bills", "Short-term UST", color="#1A5DAD")
    box(ax, 9.0, 4.5, 1.6, 0.75, "US Treasury", "receives funding", color="#1A5DAD")

    arrow(ax, 1.9,  4.5, 2.7,  4.5, color=YONSEI_BLUE, label="$100 cash")
    arrow(ax, 4.7,  4.5, 5.7,  4.5, color=YONSEI_BLUE, label="buys T-bills")
    arrow(ax, 7.3,  4.5, 8.2,  4.5, color=YONSEI_BLUE, label="~$150B total")

    # Mint tokens (down from issuer)
    box(ax, 3.7, 3.0, 2.0, 0.70, "100 USDT tokens", "$1 peg each", color="#1A5DAD")
    arrow(ax, 3.7, 4.12, 3.7, 3.35, color=YONSEI_BLUE, label="mints")

    # Scale note
    ax.text(5, 2.35,
            "Scale: combined $300B supply  →  ~$150B in T-bills  →  comparable to mid-sized sovereign reserve manager",
            ha="center", va="center", fontsize=9.5, color=GRAY,
            style="italic",
            bbox=dict(boxstyle="round,pad=0.25", facecolor="#F7F7F7", edgecolor=GRAY, lw=0.8))

    # ── STRESS FLOW (bottom row) ──
    ax.text(5, 1.90, "Stress Scenario — Run / Redemption", ha="center", va="center",
            fontsize=12, fontweight="bold", color=RED)

    box(ax, 1.1, 1.0, 1.6, 0.65, "User redeems", "wants cash back", color=RED)
    box(ax, 3.7, 1.0, 2.0, 0.65, "Issuer must", "sell T-bills", color=RED)
    box(ax, 6.5, 1.0, 1.8, 0.65, "T-bill yields", "RISE", color=RED)
    box(ax, 9.0, 1.0, 1.6, 0.65, "Spread widens", "Privilege lost", color=RED)

    arrow(ax, 1.9,  1.0, 2.7,  1.0, color=RED, lw=2)
    arrow(ax, 4.7,  1.0, 5.6,  1.0, color=RED, lw=2)
    arrow(ax, 7.4,  1.0, 8.2,  1.0, color=RED, lw=2)

    fig.tight_layout(pad=0.4)
    path = RESULTS / "fig_stablecoin_flow.png"
    fig.savefig(path, dpi=150, bbox_inches="tight")
    plt.close(fig)
    print(f"  Saved {path}")


# ═══════════════════════════════════════════════════════════════════
#  FIGURE 2 — Triffin dilemma two-panel (Slide 4)
# ═══════════════════════════════════════════════════════════════════
def make_fig_triffin_dilemma():
    set_academic_style()
    fig, axes = plt.subplots(1, 2, figsize=(11, 5.5),
                             gridspec_kw={"wspace": 0.35})

    def panel(ax, title, steps, color, title_color):
        ax.set_xlim(0, 4)
        ax.set_ylim(0, len(steps) + 1.2)
        ax.axis("off")

        # Panel background
        bg = FancyBboxPatch((0.05, 0.1), 3.9, len(steps) + 0.95,
                            boxstyle="round,pad=0.15",
                            facecolor=color + "18", edgecolor=color,
                            linewidth=2, zorder=0)
        ax.add_patch(bg)

        ax.text(2, len(steps) + 0.85, title,
                ha="center", va="center", fontsize=13,
                fontweight="bold", color=title_color)

        for i, (step, detail) in enumerate(reversed(steps)):
            y = i + 0.65
            rect = FancyBboxPatch((0.3, y - 0.27), 3.4, 0.54,
                                  boxstyle="round,pad=0.08",
                                  facecolor=color, edgecolor=color,
                                  alpha=0.85, zorder=2)
            ax.add_patch(rect)
            ax.text(2.0, y + 0.03, step, ha="center", va="center",
                    fontsize=10, fontweight="bold", color=WHITE, zorder=3)
            if detail:
                ax.text(2.0, y - 0.15, detail, ha="center", va="center",
                        fontsize=8.5, color=WHITE, alpha=0.9, zorder=3)

            if i < len(steps) - 1:
                ax.annotate("", xy=(2, y + 0.28), xytext=(2, y - 0.27 + 0.54 + 0.05),
                            arrowprops=dict(arrowstyle="-|>", color=color,
                                           lw=1.8, mutation_scale=14),
                            zorder=4)

    # Left — Normal times
    left_steps = [
        ("Supply Grows (ΔS > 0)", "new USDT/USDC issued"),
        ("Issuer Buys T-bills", "algorithmic, rule-bound"),
        ("T-bill Yields Fall", "excess demand → cheap UST"),
        ("Spread Compresses", "OIS-Treasury spread ↓"),
        ("PRIVILEGE AMPLIFIED ✓", "β₁ < 0 confirmed"),
    ]
    panel(axes[0], "Normal Times", left_steps, YONSEI_BLUE, YONSEI_BLUE)

    # Right — Stress
    right_steps = [
        ("Run / Redemptions", "users demand cash"),
        ("Issuer Sells T-bills", "forced, rapid"),
        ("T-bill Yields Rise", "supply shock"),
        ("Spread Widens", "OIS-Treasury spread ↑"),
        ("PRIVILEGE REVERSED ✗", "low liquid buffer L < 13%"),
    ]
    panel(axes[1], "Stress / Run", right_steps, RED, RED)

    # Center connector text (between subplots — use figure text)
    fig.text(0.5, 0.52,
             "Same Mechanism\nTwo Faces",
             ha="center", va="center", fontsize=12,
             fontweight="bold", color=GRAY,
             bbox=dict(boxstyle="round,pad=0.4", facecolor="#F5F5F5",
                       edgecolor=GRAY, lw=1.2))

    fig.text(0.5, 0.13,
             "Critical threshold at  L* = 13%  — liquid buffer separates regimes",
             ha="center", va="center", fontsize=11,
             fontweight="bold", color=YONSEI_ORANGE,
             bbox=dict(boxstyle="round,pad=0.3", facecolor="#FFF3E0",
                       edgecolor=YONSEI_ORANGE, lw=1.5))

    fig.tight_layout(rect=[0, 0.08, 1, 1])
    path = RESULTS / "fig_triffin_dilemma.png"
    fig.savefig(path, dpi=150, bbox_inches="tight")
    plt.close(fig)
    print(f"  Saved {path}")


# ═══════════════════════════════════════════════════════════════════
#  FIGURE 3 — Beta intuition causal chain (Slide 7)
# ═══════════════════════════════════════════════════════════════════
def make_fig_beta_intuition():
    set_academic_style()
    fig, ax = plt.subplots(figsize=(11, 4.5))
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 3.5)
    ax.axis("off")

    steps = [
        ("Stablecoin supply\ngrows (ΔS > 0)", YONSEI_BLUE),
        ("Issuer must buy\nT-bills to back\nnew supply",    YONSEI_BLUE),
        ("T-bill demand\nincreases",                        "#1A5DAD"),
        ("T-bill yields\nfall",                             "#1A5DAD"),
        ("Spread compresses\n(more negative)",              "#1A7A4A"),
        ("β₁ < 0 :\nPrivilege deepens",                    GREEN),
    ]

    xs = [1.0, 3.0, 5.0, 7.0, 9.0, 11.0]
    y0 = 2.1
    bw, bh = 1.6, 1.1

    for i, ((label, col), x) in enumerate(zip(steps, xs)):
        rect = FancyBboxPatch((x - bw/2, y0 - bh/2), bw, bh,
                              boxstyle="round,pad=0.12",
                              facecolor=col, edgecolor=col,
                              linewidth=1.5, alpha=0.92, zorder=3)
        ax.add_patch(rect)
        ax.text(x, y0, label, ha="center", va="center",
                fontsize=9.5, fontweight="bold", color=WHITE, zorder=4)

        if i < len(steps) - 1:
            ax.annotate("", xy=(xs[i+1] - bw/2 - 0.07, y0),
                        xytext=(x + bw/2 + 0.07, y0),
                        arrowprops=dict(arrowstyle="-|>",
                                        color=GRAY, lw=1.8,
                                        mutation_scale=14),
                        zorder=5)

    # Bottom note
    ax.text(6, 0.55,
            'Intuition:  β₁  should be NEGATIVE and SIGNIFICANT   ✓',
            ha="center", va="center", fontsize=12,
            fontweight="bold", color=GREEN,
            bbox=dict(boxstyle="round,pad=0.35",
                      facecolor="#E8F5E9", edgecolor=GREEN, lw=1.5))

    fig.tight_layout(pad=0.3)
    path = RESULTS / "fig_beta_intuition.png"
    fig.savefig(path, dpi=150, bbox_inches="tight")
    plt.close(fig)
    print(f"  Saved {path}")


# ═══════════════════════════════════════════════════════════════════
#  FIGURE 4 — Threshold intuition (Slide 9)
# ═══════════════════════════════════════════════════════════════════
def make_fig_threshold_intuition():
    set_academic_style()
    fig, ax = plt.subplots(figsize=(9, 5))

    q_thresh = 13.0
    q_lo = np.linspace(0, q_thresh, 80)
    q_hi = np.linspace(q_thresh, 30, 80)

    beta_lo = -6.97 - 0.05 * q_lo          # negative slope (privilege regime)
    beta_hi =  1.26 + 0.04 * (q_hi - q_thresh)   # positive/flat (stress regime)

    ax.fill_betweenx([min(beta_lo) - 1, 4],  0, q_thresh,
                     alpha=0.12, color=YONSEI_BLUE, zorder=0)
    ax.fill_betweenx([min(beta_lo) - 1, 4], q_thresh, 30,
                     alpha=0.10, color=RED, zorder=0)

    ax.plot(q_lo, beta_lo, color=YONSEI_BLUE, lw=2.8, label="β = −6.97  (low-buffer regime)")
    ax.plot(q_hi, beta_hi, color=RED, lw=2.8, label="β = +1.26  (high-buffer regime)")

    ax.axvline(q_thresh, color=YONSEI_ORANGE, lw=2.2, linestyle="--", zorder=5)
    ax.axhline(0, color=GRAY, lw=0.8, linestyle="-", alpha=0.5)

    ax.annotate("q* = 13%\n(threshold)", xy=(q_thresh, -1.5),
                xytext=(q_thresh + 4.5, -1.5),
                fontsize=11, fontweight="bold", color=YONSEI_ORANGE,
                arrowprops=dict(arrowstyle="-|>", color=YONSEI_ORANGE, lw=1.6),
                va="center")

    # regime labels
    ax.text(5.5, 3.2,
            "Below 13%:\nsupply growth\ncompresses spreads\n(privilege amplified)",
            ha="center", va="top", fontsize=9.5, color=YONSEI_BLUE,
            bbox=dict(boxstyle="round,pad=0.3", facecolor="#EFF4FF",
                      edgecolor=YONSEI_BLUE, lw=1))

    ax.text(21.5, 2.5,
            "Above 13%:\neffect reverses sign\n(forced-liquidation risk)",
            ha="center", va="top", fontsize=9.5, color=RED,
            bbox=dict(boxstyle="round,pad=0.3", facecolor="#FFEDED",
                      edgecolor=RED, lw=1))

    # Arrow: intuition note
    ax.annotate("Low buffer → forced\nliquidation risk → sign flip",
                xy=(13, 1.26), xytext=(18, 3.5),
                fontsize=9, color=GRAY,
                arrowprops=dict(arrowstyle="-|>", color=GRAY, lw=1.2),
                ha="center")

    ax.set_xlabel("Liquid Buffer L (%)", fontsize=12, color=GRAY)
    ax.set_ylabel("Effect of ΔlnS on Spread (bps)", fontsize=12, color=GRAY)
    ax.set_title("Threshold Model: Regime-Dependent Effect of Supply Growth",
                 fontsize=13, fontweight="bold", color=YONSEI_BLUE, pad=10)
    ax.set_xlim(0, 30)
    ax.set_ylim(min(beta_lo) - 1.5, 5)
    ax.legend(fontsize=10, framealpha=0.9, loc="lower right")

    fig.tight_layout(pad=0.6)
    path = RESULTS / "fig_threshold_intuition.png"
    fig.savefig(path, dpi=150, bbox_inches="tight")
    plt.close(fig)
    print(f"  Saved {path}")


# ═══════════════════════════════════════════════════════════════════
#  FIGURE 5 & 6 — Raw event study plots (Slides 11)
# ═══════════════════════════════════════════════════════════════════
def make_fig_event_raw():
    df = pd.read_csv(DATA / "daily_panel.csv", parse_dates=["date"])
    df["spread_bps"] = df["spread"] * 100

    # ── LUNA / UST ──
    mask = (df["date"] >= "2022-04-25") & (df["date"] <= "2022-05-25")
    sub = df[mask].copy()

    fig, ax = plt.subplots(figsize=(9, 4.5))
    set_academic_style()

    ax.plot(sub["date"], sub["spread_bps"],
            color=YONSEI_BLUE, lw=2.2, zorder=3)
    ax.fill_between(sub["date"], sub["spread_bps"],
                    sub["spread_bps"].min() - 2,
                    alpha=0.08, color=YONSEI_BLUE)

    # Event lines
    fomc_may4  = pd.Timestamp("2022-05-04")
    luna_may9  = pd.Timestamp("2022-05-09")
    usdt_may12 = pd.Timestamp("2022-05-12")

    ax.axvline(fomc_may4,  color=RED,         lw=2.2, linestyle="--", zorder=5)
    ax.axvline(luna_may9,  color=YONSEI_BLUE, lw=2.2, linestyle="-",  zorder=5)
    ax.axvline(usdt_may12, color=YONSEI_BLUE, lw=2.0, linestyle=":",  zorder=5)

    ymax = sub["spread_bps"].max()
    ymin = sub["spread_bps"].min()
    yr   = ymax - ymin

    ax.text(fomc_may4,  ymax - yr*0.04, " FOMC\n May 4\n +50bps",
            color=RED, fontsize=9, fontweight="bold", va="top")
    ax.text(luna_may9,  ymax - yr*0.04, " LUNA\n May 9",
            color=YONSEI_BLUE, fontsize=9, fontweight="bold", va="top")
    ax.text(usdt_may12, ymin + yr*0.06, " USDT\n depeg\n May 12",
            color=YONSEI_BLUE, fontsize=9, va="bottom")

    # contamination label
    ax.annotate("",
                xy=(luna_may9, (ymax+ymin)/2),
                xytext=(fomc_may4, (ymax+ymin)/2),
                arrowprops=dict(arrowstyle="<->",
                                color=YONSEI_ORANGE, lw=1.8))
    ax.text(pd.Timestamp("2022-05-06"), (ymax+ymin)/2 + yr*0.05,
            "5 days", ha="center", fontsize=9,
            color=YONSEI_ORANGE, fontweight="bold")

    ax.set_title("Raw OIS-Treasury Spread: LUNA/UST Event (Apr 25 – May 25, 2022)",
                 fontsize=12, fontweight="bold", color=YONSEI_BLUE)
    ax.set_ylabel("Spread (bps)", fontsize=11)
    ax.set_xlabel("")
    ax.tick_params(axis="x", rotation=20)

    # Contamination note box
    ax.text(0.5, -0.18,
            "FOMC hike just 5 days before LUNA — contamination prevents clean causal inference",
            transform=ax.transAxes, ha="center", fontsize=9.5,
            color=RED, fontweight="bold",
            bbox=dict(boxstyle="round,pad=0.25", facecolor="#FFEDED",
                      edgecolor=RED, lw=1))

    fig.tight_layout(pad=0.6, rect=[0, 0.07, 1, 1])
    path = RESULTS / "fig_event_raw_luna.png"
    fig.savefig(path, dpi=150, bbox_inches="tight")
    plt.close(fig)
    print(f"  Saved {path}")

    # ── SVB / USDC ──
    mask2 = (df["date"] >= "2023-02-25") & (df["date"] <= "2023-04-01")
    sub2 = df[mask2].copy()

    fig2, ax2 = plt.subplots(figsize=(9, 4.5))
    set_academic_style()

    ax2.plot(sub2["date"], sub2["spread_bps"],
             color=YONSEI_BLUE, lw=2.2, zorder=3)
    ax2.fill_between(sub2["date"], sub2["spread_bps"], -20,
                     alpha=0.08, color=YONSEI_BLUE)
    ax2.axhline(0, color=GRAY, lw=0.9, linestyle="-", alpha=0.6, label="0 bps")

    svb_mar11   = pd.Timestamp("2023-03-11")
    fomc_mar22  = pd.Timestamp("2023-03-22")

    ax2.axvline(svb_mar11,  color=YONSEI_BLUE, lw=2.2, linestyle="-",  zorder=5)
    ax2.axvline(fomc_mar22, color=RED,         lw=2.2, linestyle="--", zorder=5)

    ymax2 = sub2["spread_bps"].max()
    ymin2 = sub2["spread_bps"].min()

    ax2.text(svb_mar11, ymax2 * 0.95, " SVB\n Mar 11",
             color=YONSEI_BLUE, fontsize=9, fontweight="bold", va="top")
    ax2.text(fomc_mar22, ymax2 * 0.6, " FOMC\n Mar 22\n +25bps",
             color=RED, fontsize=9, fontweight="bold", va="top")

    # Peak-to-trough annotation
    ax2.annotate("",
                 xy=(svb_mar11 + pd.Timedelta("6d"), -8.6),
                 xytext=(svb_mar11 - pd.Timedelta("4d"), 43.2),
                 arrowprops=dict(arrowstyle="-|>", color=YONSEI_ORANGE,
                                 lw=2.0, mutation_scale=14))
    ax2.text(pd.Timestamp("2023-03-13"), 22,
             "−51.8 bps\nin 6 days",
             ha="left", fontsize=9.5, color=YONSEI_ORANGE,
             fontweight="bold")

    ax2.set_title("Raw OIS-Treasury Spread: SVB / USDC Event (Feb 25 – Apr 1, 2023)",
                  fontsize=12, fontweight="bold", color=YONSEI_BLUE)
    ax2.set_ylabel("Spread (bps)", fontsize=11)
    ax2.tick_params(axis="x", rotation=20)

    ax2.text(0.5, -0.18,
             "Flight-to-safety: T-bills became MORE expensive than risk-free — government guarantee\nprevented forced Tether liquidation; this is NOT a stablecoin-driven event",
             transform=ax2.transAxes, ha="center", fontsize=9.2,
             color=RED, fontweight="bold",
             bbox=dict(boxstyle="round,pad=0.25", facecolor="#FFEDED",
                       edgecolor=RED, lw=1))

    fig2.tight_layout(pad=0.6, rect=[0, 0.08, 1, 1])
    path2 = RESULTS / "fig_event_raw_svb.png"
    fig2.savefig(path2, dpi=150, bbox_inches="tight")
    plt.close(fig2)
    print(f"  Saved {path2}")


# ═══════════════════════════════════════════════════════════════════
#  Build all figures
# ═══════════════════════════════════════════════════════════════════
print("Building matplotlib figures…")
make_fig_stablecoin_flow()
make_fig_triffin_dilemma()
make_fig_beta_intuition()
make_fig_threshold_intuition()
make_fig_event_raw()
print("All figures saved.\n")


# ═══════════════════════════════════════════════════════════════════
#  PPTX helpers
# ═══════════════════════════════════════════════════════════════════

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_layout(prs):
    """Return the blank slide layout."""
    return prs.slide_layouts[6]

# ── color setters ──
def rgb(r, g, b):
    return RGBColor(r, g, b)

def set_fill(shape, r, g, b):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb(r, g, b)

def set_border_none(shape):
    try:
        shape.line.fill.background()
    except Exception:
        pass

# ── text box ──
def add_textbox(slide, left, top, width, height,
                text, font_size=14, bold=False, italic=False,
                color=C_BLACK, align=PP_ALIGN.LEFT,
                font_name="Calibri", wrap=True, line_spacing=None):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name  = font_name
    f.size  = Pt(font_size)
    f.bold  = bold
    f.italic= italic
    f.color.rgb = color
    if line_spacing:
        from pptx.util import Pt as _Pt
        p.line_spacing = line_spacing
    return txb

def add_paragraph(tf, text, font_size=11, bold=False, italic=False,
                  color=C_BLACK, align=PP_ALIGN.LEFT, font_name="Calibri",
                  bullet=False, indent=0):
    p = tf.add_paragraph()
    p.alignment = align
    if indent:
        p.level = indent
    run = p.add_run()
    run.text = text
    f = run.font
    f.name  = font_name
    f.size  = Pt(font_size)
    f.bold  = bold
    f.italic= italic
    f.color.rgb = color
    return p

# ── colored rectangle ──
def add_rect(slide, left, top, width, height, r, g, b):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height)
    set_fill(shape, r, g, b)
    shape.line.fill.background()
    return shape

# ── image ──
def add_image(slide, path, left, top, width, height=None):
    if height:
        return slide.shapes.add_picture(str(path), left, top, width, height)
    return slide.shapes.add_picture(str(path), left, top, width)

# ── slide header bar ──
def add_header(slide, title, subtitle="", section=""):
    # Dark blue header bar
    hdr = add_rect(slide, 0, 0, SLIDE_W, Inches(0.9), 0, 0x30, 0x87)
    # title text
    add_textbox(slide, Inches(0.3), Inches(0.07), Inches(11.5), Inches(0.5),
                title, font_size=20, bold=True, color=C_WHITE,
                align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, Inches(0.3), Inches(0.55), Inches(9), Inches(0.32),
                    subtitle, font_size=11, italic=True, color=C_WHITE,
                    align=PP_ALIGN.LEFT)
    if section:
        add_textbox(slide, Inches(10.5), Inches(0.05), Inches(2.6), Inches(0.35),
                    section, font_size=9, color=C_ORANGE,
                    align=PP_ALIGN.RIGHT, bold=True)

# ── sidebar (slides 10–16) ──
def add_sidebar(slide):
    sb_left  = Inches(10.83)
    sb_width = Inches(2.50)
    # background
    add_rect(slide, sb_left, Inches(0.9), sb_width, Inches(6.6), 0, 0x30, 0x87)
    # top orange accent
    add_rect(slide, sb_left, Inches(0.9), sb_width, Inches(0.04), 0xE8, 0x77, 0x22)

    add_textbox(slide, sb_left + Inches(0.1), Inches(1.0), Inches(2.3), Inches(1.2),
                "Key Results", font_size=11, bold=True, color=C_ORANGE,
                align=PP_ALIGN.CENTER)

    add_textbox(slide, sb_left + Inches(0.08), Inches(2.0), Inches(2.3), Inches(0.6),
                "β₁ = −6.02***", font_size=15, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, sb_left + Inches(0.08), Inches(2.55), Inches(2.3), Inches(0.4),
                "p = 0.006", font_size=10, italic=True, color=C_ORANGE,
                align=PP_ALIGN.CENTER)

    add_rect(slide, sb_left + Inches(0.15), Inches(3.05),
             Inches(2.2), Inches(0.025), 0xE8, 0x77, 0x22)

    add_textbox(slide, sb_left + Inches(0.08), Inches(3.15), Inches(2.3), Inches(0.6),
                "q* = 13%", font_size=15, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, sb_left + Inches(0.08), Inches(3.70), Inches(2.3), Inches(0.5),
                "Bootstrap CI\n[3.2%, 14.5%]", font_size=9, italic=True,
                color=C_ORANGE, align=PP_ALIGN.CENTER)

    add_rect(slide, sb_left + Inches(0.15), Inches(4.3),
             Inches(2.2), Inches(0.025), 0xE8, 0x77, 0x22)

    add_textbox(slide, sb_left + Inches(0.08), Inches(4.40), Inches(2.3), Inches(1.6),
                "3 convergent\nmethods:\n\nβ₁ ✓\nq* ✓\nEvent study\n(directional)",
                font_size=9.5, color=C_WHITE, align=PP_ALIGN.CENTER)

# ── table helper ──
def add_table_slide(slide, left, top, width, height,
                    rows, headers, col_widths,
                    header_bg=(0, 0x30, 0x87),
                    row_bg1=(0xF2, 0xF2, 0xF2),
                    row_bg2=(0xFF, 0xFF, 0xFF),
                    font_size=10):
    """Add a formatted table to the slide."""
    nrows = len(rows) + 1  # +1 for header
    ncols = len(headers)
    tbl = slide.shapes.add_table(nrows, ncols,
                                  left, top, width, height).table
    # Column widths
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = cw

    # Header row
    for j, hdr in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(*header_bg)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.bold  = True
        run.font.color.rgb = C_WHITE
        run.font.size  = Pt(font_size)
        run.font.name  = "Calibri"

    # Data rows
    for i, row in enumerate(rows):
        bg = row_bg1 if i % 2 == 0 else row_bg2
        for j, val in enumerate(row):
            cell = tbl.cell(i+1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(*bg)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size  = Pt(font_size)
            run.font.name  = "Calibri"
            run.font.color.rgb = C_BLACK
    return tbl


# ═══════════════════════════════════════════════════════════════════
#  BUILD PRESENTATION
# ═══════════════════════════════════════════════════════════════════
print("Building PPTX…")
prs = new_prs()
blank = blank_layout(prs)

# ── content width helpers ──
MAIN_LEFT    = Inches(0.3)
MAIN_TOP     = Inches(1.0)
MAIN_WIDTH   = Inches(10.4)     # without sidebar
MAIN_WIDTH_S = Inches(10.15)    # with sidebar (slides 10+)
MAIN_HEIGHT  = Inches(6.3)


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(blank)

# Full dark blue background
add_rect(slide1, 0, 0, SLIDE_W, SLIDE_H, 0, 0x30, 0x87)

# Orange accent bar
add_rect(slide1, 0, Inches(3.6), SLIDE_W, Inches(0.06), 0xE8, 0x77, 0x22)

# Title
add_textbox(slide1, Inches(0.8), Inches(1.2), Inches(11.7), Inches(1.0),
            "Stablecoins and the Exorbitant Privilege",
            font_size=34, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# Subtitle
add_textbox(slide1, Inches(0.8), Inches(2.4), Inches(11.7), Inches(0.7),
            "Safe-Asset Demand and Its Systemic Fragility",
            font_size=21, italic=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

# Authors
add_textbox(slide1, Inches(0.8), Inches(3.85), Inches(11.7), Inches(0.7),
            "Mireu Kim   Sara Chekroune   Oybek Ibragimov   Jade Zhu",
            font_size=13, color=C_WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide1, Inches(0.8), Inches(4.35), Inches(11.7), Inches(0.55),
            "Alexandre Godefroy   Baptiste Degand   Minjin Kim",
            font_size=13, color=C_WHITE, align=PP_ALIGN.CENTER)

# Institution & course
add_textbox(slide1, Inches(0.8), Inches(5.1), Inches(11.7), Inches(0.5),
            "Yonsei GSIS  |  Topics in International Finance  |  June 2026",
            font_size=12, italic=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

# Key metrics at bottom
add_rect(slide1, Inches(2.5), Inches(6.0), Inches(3.5), Inches(0.8), 0, 0x20, 0x60)
add_textbox(slide1, Inches(2.5), Inches(6.0), Inches(3.5), Inches(0.8),
            "β₁ = −6.02***  (p = 0.006)", font_size=14, bold=True,
            color=C_ORANGE, align=PP_ALIGN.CENTER)

add_rect(slide1, Inches(7.3), Inches(6.0), Inches(3.5), Inches(0.8), 0, 0x20, 0x60)
add_textbox(slide1, Inches(7.3), Inches(6.0), Inches(3.5), Inches(0.8),
            "q* = 13.0%  liquid buffer threshold",
            font_size=14, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 2 — What is a stablecoin?
# ══════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(blank)
add_rect(slide2, 0, 0, SLIDE_W, SLIDE_H, 0xF8, 0xF9, 0xFA)
add_header(slide2, "What Is a Stablecoin?",
           "A programmatic, rule-bound channel for US dollar safe-asset demand",
           "Slide 2 / 16")

add_image(slide2, RESULTS / "fig_stablecoin_flow.png",
          Inches(0.3), Inches(1.0), Inches(12.7))


# ══════════════════════════════════════════════════════════════
# SLIDE 3 — The Exorbitant Privilege
# ══════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(blank)
add_rect(slide3, 0, 0, SLIDE_W, SLIDE_H, 0xF8, 0xF9, 0xFA)
add_header(slide3, "The Exorbitant Privilege — What Finance Already Knows",
           "Stablecoins as a new, algorithmic channel for this classic mechanism",
           "Slide 3 / 16")

# Left column
lc_left  = Inches(0.3)
lc_top   = Inches(1.05)
lc_width = Inches(6.2)

# Gourinchas & Rey box
gr_box = add_rect(slide3, lc_left, lc_top, lc_width, Inches(1.3),
                  0xEF, 0xF4, 0xFF)
add_textbox(slide3, lc_left + Inches(0.1), lc_top + Inches(0.05),
            lc_width - Inches(0.2), Inches(0.35),
            "Gourinchas & Rey (2007)", font_size=12, bold=True,
            color=C_BLUE)
add_textbox(slide3, lc_left + Inches(0.1), lc_top + Inches(0.38),
            lc_width - Inches(0.2), Inches(0.8),
            "US earns more on foreign assets than it pays on liabilities → structural "
            "borrowing advantage. Global demand for dollar safe assets keeps T-bill "
            "yields low → US borrows cheap.",
            font_size=10.5, color=C_BLACK)

# Mechanism box
mech_top = lc_top + Inches(1.45)
add_rect(slide3, lc_left, mech_top, lc_width, Inches(1.25),
         0xEF, 0xF4, 0xFF)
add_textbox(slide3, lc_left + Inches(0.1), mech_top + Inches(0.05),
            lc_width - Inches(0.2), Inches(0.35),
            "Mechanism: OIS-Treasury Spread as Our Measure",
            font_size=12, bold=True, color=C_BLUE)
add_textbox(slide3, lc_left + Inches(0.1), mech_top + Inches(0.38),
            lc_width - Inches(0.2), Inches(0.75),
            "T-bill yield minus OIS ≈ 0 in normal times.\n"
            "NEGATIVE spread = T-bills MORE expensive than risk-free = extreme privilege.\n"
            "Normal range: 10–50 bps.   10 bps IS economically large.",
            font_size=10.5, color=C_BLACK)

# Our new channel box
nc_top = mech_top + Inches(1.4)
add_rect(slide3, lc_left, nc_top, lc_width, Inches(1.35),
         0xFF, 0xF3, 0xE0)
add_textbox(slide3, lc_left + Inches(0.1), nc_top + Inches(0.05),
            lc_width - Inches(0.2), Inches(0.35),
            "The New Channel: Stablecoins", font_size=12, bold=True,
            color=C_ORANGE)
add_textbox(slide3, lc_left + Inches(0.1), nc_top + Inches(0.38),
            lc_width - Inches(0.2), Inches(0.85),
            "Not a central bank, not a fund manager — an ALGORITHM that must buy T-bills "
            "as supply grows. Combined supply ~$300B → ~$150B in T-bills.\n"
            "Programmatic, rule-bound, and growing fast.",
            font_size=10.5, color=C_BLACK)

# Right column — visual timeline/quote
rc_left  = Inches(6.7)
rc_width = Inches(6.3)

add_rect(slide3, rc_left, lc_top, rc_width, Inches(4.2), 0, 0x30, 0x87)

add_textbox(slide3, rc_left + Inches(0.2), lc_top + Inches(0.15),
            rc_width - Inches(0.4), Inches(0.5),
            "How the Privilege Works", font_size=13, bold=True, color=C_ORANGE,
            align=PP_ALIGN.CENTER)

privilege_steps = [
    ("1", "Global investors want safe dollar assets"),
    ("2", "They buy US T-bills → demand pushes yields down"),
    ("3", "OIS-Treasury spread compresses (or goes negative)"),
    ("4", "US Treasury borrows at below-market rates"),
    ("5", "Stablecoins: algorithmic, always-on version of step 2"),
]
for idx, (num, txt) in enumerate(privilege_steps):
    y_pos = lc_top + Inches(0.75) + idx * Inches(0.68)
    add_rect(slide3, rc_left + Inches(0.15), y_pos,
             Inches(0.45), Inches(0.45), 0xE8, 0x77, 0x22)
    add_textbox(slide3, rc_left + Inches(0.15), y_pos,
                Inches(0.45), Inches(0.45),
                num, font_size=13, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(slide3, rc_left + Inches(0.7), y_pos,
                rc_width - Inches(0.9), Inches(0.45),
                txt, font_size=10.5, color=C_WHITE)

# Bottom callout
add_rect(slide3, Inches(0.3), Inches(5.6), Inches(12.73), Inches(0.65),
         0x00, 0x30, 0x87)
add_textbox(slide3, Inches(0.5), Inches(5.65), Inches(12.33), Inches(0.55),
            "Key insight: Stablecoins create a NEW demand channel for T-bills — "
            "one that is larger, faster, and more rule-bound than sovereign reserve managers.",
            font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 4 — New Triffin Dilemma
# ══════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(blank)
add_rect(slide4, 0, 0, SLIDE_W, SLIDE_H, 0xF8, 0xF9, 0xFA)
add_header(slide4, "The New Triffin Dilemma — Our Core Insight",
           "Same mechanism, two faces: supply growth amplifies OR undermines the privilege",
           "Slide 4 / 16")

add_image(slide4, RESULTS / "fig_triffin_dilemma.png",
          Inches(0.3), Inches(1.0), Inches(12.7))


# ══════════════════════════════════════════════════════════════
# SLIDE 5 — Literature & Our Contribution
# ══════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(blank)
add_rect(slide5, 0, 0, SLIDE_W, SLIDE_H, 0xF8, 0xF9, 0xFA)
add_header(slide5, "Literature & Our Contribution",
           "We fill a gap no existing paper has addressed empirically",
           "Slide 5 / 16")

lit_headers = ["What we build on", "Paper", "Key insight"]
lit_rows = [
    ["Safe-asset demand & privilege",  "Maggiori (2017)",       "Exorbitant privilege from risk-sharing equilibrium"],
    ["Crisis zone mechanism",          "Cole & Kehoe (2000)",   "Run dynamics when fundamentals enter crisis zone"],
    ["Stablecoin banking",             "Gorton & Zhang (2021)", "Stablecoins = wildcat banking, uninsured deposits"],
    ["Reserve adequacy",               "Jeanne & Rancière (2011)", "Optimal buffers against sudden stops"],
    ["Stablecoin regulation",          "Duffie (2022)",         "Reserve composition is first-order concern"],
]
col_widths = [Inches(2.8), Inches(2.4), Inches(5.3)]
add_table_slide(slide5,
                Inches(0.3), Inches(1.05),
                Inches(10.5), Inches(2.9),
                lit_rows, lit_headers, col_widths,
                font_size=10)

# OUR GAP box
gap_top = Inches(4.15)
add_rect(slide5, Inches(0.3), gap_top, Inches(5.9), Inches(1.1),
         0xFF, 0xED, 0xED)
add_textbox(slide5, Inches(0.4), gap_top + Inches(0.08),
            Inches(5.7), Inches(0.4),
            "OUR GAP", font_size=12, bold=True, color=C_RED)
add_textbox(slide5, Inches(0.4), gap_top + Inches(0.44),
            Inches(5.7), Inches(0.55),
            "No paper has connected reserve adequacy → sovereign bond spreads empirically. We do.",
            font_size=10.5, color=C_BLACK)

# OUR CONTRIBUTION box
add_rect(slide5, Inches(6.4), gap_top, Inches(6.9), Inches(1.1),
         0xE8, 0xF5, 0xE9)
add_textbox(slide5, Inches(6.5), gap_top + Inches(0.08),
            Inches(6.7), Inches(0.4),
            "OUR CONTRIBUTION", font_size=12, bold=True, color=C_GREEN)
add_textbox(slide5, Inches(6.5), gap_top + Inches(0.44),
            Inches(6.7), Inches(0.55),
            "First empirical threshold for stablecoin liquid reserves that separates "
            "privilege-amplification from forced-liquidation regimes.",
            font_size=10.5, color=C_BLACK)

# Bottom bar
add_rect(slide5, Inches(0.3), Inches(5.45), Inches(12.73), Inches(0.7),
         0, 0x30, 0x87)
add_textbox(slide5, Inches(0.5), Inches(5.5), Inches(12.33), Inches(0.6),
            "β₁ = −6.02***  |  q* = 13%  |  Three convergent methods: "
            "OLS regression  +  Hansen threshold  +  LSTAR nonlinear model",
            font_size=12, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 6 — Data & The Spread
# ══════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(blank)
add_rect(slide6, 0, 0, SLIDE_W, SLIDE_H, 0xF8, 0xF9, 0xFA)
add_header(slide6, "Data & The OIS-Treasury Spread",
           "Monthly panel 2020–2024; spread from FRED; stablecoin supply from DeFiLlama",
           "Slide 6 / 16")

# timeseries figure
add_image(slide6, RESULTS / "fig_timeseries.png",
          Inches(0.25), Inches(1.0), Inches(8.3))

# Right annotation panel
rp = Inches(8.7)
add_rect(slide6, rp, Inches(1.05), Inches(4.35), Inches(6.2), 0, 0x30, 0x87)

add_textbox(slide6, rp + Inches(0.15), Inches(1.15),
            Inches(4.05), Inches(0.45),
            "Reading the Spread", font_size=12, bold=True,
            color=C_ORANGE)

annotations = [
    ("Normal range:", "10–50 bps", "10 bps IS economically large"),
    ("March 2023:", "−8.6 bps", "T-bills MORE expensive than risk-free\n= extreme flight-to-safety"),
    ("LUNA May 2022:", "+76 bps peak", "Spread WIDENED during crypto stress"),
    ("Our measure:", "OIS minus T-bill yield", "≈ 0 in normal times;\nnegative = extreme privilege"),
]
y_ann = Inches(1.7)
for label, val, detail in annotations:
    add_rect(slide6, rp + Inches(0.1), y_ann,
             Inches(4.15), Inches(1.05), 0, 0x20, 0x65)
    add_textbox(slide6, rp + Inches(0.2), y_ann + Inches(0.04),
                Inches(3.9), Inches(0.32),
                f"{label}  {val}", font_size=10.5, bold=True,
                color=C_ORANGE)
    add_textbox(slide6, rp + Inches(0.2), y_ann + Inches(0.35),
                Inches(3.9), Inches(0.62),
                detail, font_size=9.5, color=C_WHITE)
    y_ann += Inches(1.15)

# Stress events note
add_rect(slide6, rp + Inches(0.1), y_ann, Inches(4.15), Inches(0.65),
         0xC0, 0x39, 0x2B)
add_textbox(slide6, rp + Inches(0.2), y_ann + Inches(0.08),
            Inches(3.9), Inches(0.5),
            "3 stress events marked:\nLUNA/UST   USDC/SVB   USDT Mar 2023",
            font_size=9.5, bold=True, color=C_WHITE)


# ══════════════════════════════════════════════════════════════
# SLIDE 7 — Intuition for β₁
# ══════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(blank)
add_rect(slide7, 0, 0, SLIDE_W, SLIDE_H, 0xF8, 0xF9, 0xFA)
add_header(slide7, "Intuition for β₁ — The Privilege Channel",
           "Before seeing the result: what should happen mechanically?",
           "Slide 7 / 16")

add_image(slide7, RESULTS / "fig_beta_intuition.png",
          Inches(0.3), Inches(1.05), Inches(12.7))

# Bottom context
add_rect(slide7, Inches(0.3), Inches(5.65), Inches(12.73), Inches(0.7),
         0, 0x30, 0x87)
add_textbox(slide7, Inches(0.5), Inches(5.70), Inches(12.33), Inches(0.6),
            "If stablecoins truly act as a safe-asset demand channel, β₁ must be NEGATIVE — "
            "supply growth compresses the spread.   We test this with N = 51 monthly observations.",
            font_size=11.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 8 — Main Result: β₁ = −6.02
# ══════════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(blank)
add_rect(slide8, 0, 0, SLIDE_W, SLIDE_H, 0xF8, 0xF9, 0xFA)
add_header(slide8, "Main Result: β₁ = −6.02***   (THE KEY FINDING)",
           "Table 2 — OLS with HAC standard errors, N = 51 monthly observations",
           "Slide 8 / 16")

# ── Regression table ──
reg_headers = ["Variable", "Coeff.", "Std. Err.", "p-value", "Interpretation"]
reg_rows = [
    ["Δln(Supply)  β₁",  "−6.02", "2.07", "0.006**",  "Supply growth → spread compression"],
    ["VIX  β₂",          "+0.043","0.019","0.032*",   "Global risk-off widens spread"],
    ["ΔFed Funds  β₃",   "+0.18", "0.14", "0.213",    "Rate changes (partially absorbed)"],
    ["Constant",         "+12.3", "8.4",  "0.148",    "Baseline spread level"],
    ["R²  =  0.502",     "—",     "—",    "—",         "Regression explains 50% of variance"],
]
col_widths_r = [Inches(2.3), Inches(1.2), Inches(1.2), Inches(1.4), Inches(4.1)]
add_table_slide(slide8,
                Inches(0.3), Inches(1.05),
                Inches(10.2), Inches(2.5),
                reg_rows, reg_headers, col_widths_r,
                font_size=10.5)

# ── Big callout box ──
cb_top = Inches(3.7)
add_rect(slide8, Inches(0.3), cb_top, Inches(7.5), Inches(1.6),
         0, 0x30, 0x87)
add_textbox(slide8, Inches(0.45), cb_top + Inches(0.08),
            Inches(7.2), Inches(0.45),
            "Economic Magnitude", font_size=13, bold=True, color=C_ORANGE)
add_textbox(slide8, Inches(0.45), cb_top + Inches(0.52),
            Inches(7.2), Inches(0.95),
            "1-SD supply growth (4 pp)  ×  β₁  =  −6.02  →  24 bps spread compression\n"
            "This is larger than a typical FOMC surprise effect on T-bill spreads.",
            font_size=11.5, color=C_WHITE)

# ── Intuition check ──
add_rect(slide8, Inches(8.0), cb_top, Inches(5.05), Inches(1.6),
         0x1A, 0x7A, 0x4A)
add_textbox(slide8, Inches(8.15), cb_top + Inches(0.08),
            Inches(4.75), Inches(0.45),
            "Intuition Check", font_size=13, bold=True, color=C_WHITE)
add_textbox(slide8, Inches(8.15), cb_top + Inches(0.52),
            Inches(4.75), Inches(0.95),
            "β₁ < 0  ✓  (supply growth compresses spread)\n"
            "VIX > 0  ✓  (risk-off widens spread)\n"
            "MATCHES theoretical prediction",
            font_size=11.5, color=C_WHITE)

# ── Bottom summary ──
add_rect(slide8, Inches(0.3), Inches(5.5), Inches(12.73), Inches(0.8),
         0, 0x30, 0x87)
add_textbox(slide8, Inches(0.5), Inches(5.55), Inches(12.33), Inches(0.7),
            "\"A one-standard-deviation increase in stablecoin supply growth compresses the "
            "OIS-Treasury spread by 24 basis points — confirming H1\"",
            font_size=12.5, bold=True, italic=True, color=C_ORANGE,
            align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 9 — Threshold Intuition
# ══════════════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(blank)
add_rect(slide9, 0, 0, SLIDE_W, SLIDE_H, 0xF8, 0xF9, 0xFA)
add_header(slide9, "Intuition for the Threshold — Why Should There Be a Regime Break?",
           "Low buffer = forced liquidation risk; high buffer = orderly redemption",
           "Slide 9 / 16")

add_image(slide9, RESULTS / "fig_threshold_intuition.png",
          Inches(0.4), Inches(1.0), Inches(8.5))

# Right column explanation
rp9 = Inches(9.1)
add_rect(slide9, rp9, Inches(1.05), Inches(4.0), Inches(6.2), 0, 0x30, 0x87)

add_textbox(slide9, rp9 + Inches(0.15), Inches(1.15),
            Inches(3.7), Inches(0.45),
            "The Logic", font_size=12, bold=True, color=C_ORANGE)

logic_pts = [
    "When liquid buffer L is LOW:",
    "→ A redemption wave forces rapid T-bill sales",
    "→ Spread dynamics flip sign",
    "",
    "When liquid buffer L is HIGH:",
    "→ Issuer can absorb redemptions from cash",
    "→ T-bill holdings untouched",
    "→ Privilege channel continues",
    "",
    "Critical threshold: L* ≈ 13%",
    "Hansen grid search confirms this",
    "LSTAR c* = 13.1% (convergent)",
]
y_lp = Inches(1.7)
for pt in logic_pts:
    if pt == "":
        y_lp += Inches(0.12)
        continue
    fc = C_ORANGE if pt.startswith("Critical") or pt.startswith("Hansen") or pt.startswith("LSTAR") else C_WHITE
    fs = 10.5 if pt.endswith(":") else 9.5
    bld = pt.endswith(":")
    add_textbox(slide9, rp9 + Inches(0.15), y_lp,
                Inches(3.7), Inches(0.35),
                pt, font_size=fs, bold=bld, color=fc)
    y_lp += Inches(0.36)


# ══════════════════════════════════════════════════════════════
# SLIDE 10 — Threshold Result  [SIDEBAR STARTS]
# ══════════════════════════════════════════════════════════════
slide10 = prs.slides.add_slide(blank)
add_rect(slide10, 0, 0, SLIDE_W, SLIDE_H, 0xF8, 0xF9, 0xFA)
add_header(slide10, "Threshold Result: q* = 13%",
           "Hansen (2000) grid search + LSTAR convergent validity",
           "Slide 10 / 16")
add_sidebar(slide10)

# SSR plot
add_image(slide10, RESULTS / "threshold_ssr.png",
          Inches(0.3), Inches(1.05), Inches(5.8))

# Results table
th_headers = ["", "Low-buffer  (N=38)", "High-buffer  (N=13)"]
th_rows = [
    ["β_ΔlnS",    "−6.97",  "+1.26  ← SIGN FLIP"],
    ["q* (Hansen)", "0.1301 (13.0%)", "—"],
    ["LSTAR c*",  "0.1310 (13.1%)", "convergent validity ✓"],
    ["Bootstrap p", "0.260 (suggestive)", "N=51 limits power"],
    ["TRIM stability", "13% at TRIM 15/20/25%", "stable across trimming"],
]
col_w_th = [Inches(2.2), Inches(3.3), Inches(3.0)]
add_table_slide(slide10,
                Inches(6.2), Inches(1.05),
                Inches(8.5), Inches(2.65),
                th_rows, th_headers, col_w_th,
                font_size=10)

# Note box
add_rect(slide10, Inches(6.2), Inches(3.9), Inches(4.45), Inches(1.1),
         0xFF, 0xF3, 0xE0)
add_textbox(slide10, Inches(6.3), Inches(3.95),
            Inches(4.25), Inches(1.0),
            "Why p = 0.260?   N = 51 monthly obs.   Bootstrap\n"
            "distribution is wide — 90% CI [3.2%, 14.5%].\n"
            "Direction confirmed; magnitude suggestive.",
            font_size=10, color=C_BLACK)

add_rect(slide10, Inches(0.3), Inches(5.5), Inches(10.35), Inches(0.75),
         0, 0x30, 0x87)
add_textbox(slide10, Inches(0.5), Inches(5.55),
            Inches(10.0), Inches(0.65),
            "Hansen threshold at 13% replicated by independent LSTAR estimation (c* = 13.1%).\n"
            "Convergent validity across two methods strengthens our confidence.",
            font_size=11, bold=True, color=C_WHITE)


# ══════════════════════════════════════════════════════════════
# SLIDE 11 — Event Study
# ══════════════════════════════════════════════════════════════
slide11 = prs.slides.add_slide(blank)
add_rect(slide11, 0, 0, SLIDE_W, SLIDE_H, 0xF8, 0xF9, 0xFA)
add_header(slide11, "Event Study: Raw Spread Around Stress Events",
           "Showing raw daily OIS-Treasury spread — NOT CARs (corrected approach)",
           "Slide 11 / 16")
add_sidebar(slide11)

# LUNA figure (left)
add_textbox(slide11, Inches(0.3), Inches(1.0), Inches(5.0), Inches(0.35),
            "Panel A: LUNA/UST — May 2022",
            font_size=11, bold=True, color=C_BLUE)
add_image(slide11, RESULTS / "fig_event_raw_luna.png",
          Inches(0.2), Inches(1.38), Inches(5.3))

# SVB figure (right)
add_textbox(slide11, Inches(5.5), Inches(1.0), Inches(5.1), Inches(0.35),
            "Panel B: SVB / USDC — March 2023",
            font_size=11, bold=True, color=C_BLUE)
add_image(slide11, RESULTS / "fig_event_raw_svb.png",
          Inches(5.4), Inches(1.38), Inches(5.3))

# Takeaway bar
add_rect(slide11, Inches(0.2), Inches(5.55), Inches(10.4), Inches(0.75),
         0, 0x30, 0x87)
add_textbox(slide11, Inches(0.35), Inches(5.60),
            Inches(10.1), Inches(0.65),
            "LUNA: spread barely moved — FOMC hike 5 days prior contaminates inference.\n"
            "SVB: dramatic compression (−52 bps) = flight-to-safety, NOT forced T-bill selling.",
            font_size=10.5, bold=True, color=C_WHITE)


# ══════════════════════════════════════════════════════════════
# SLIDE 12 — Where Results Differ from Intuition
# ══════════════════════════════════════════════════════════════
slide12 = prs.slides.add_slide(blank)
add_rect(slide12, 0, 0, SLIDE_W, SLIDE_H, 0xF8, 0xF9, 0xFA)
add_header(slide12, "Intuition vs. Results — Where They Converge & Diverge",
           "Honest assessment of what the data can and cannot tell us",
           "Slide 12 / 16")
add_sidebar(slide12)

cmp_headers = ["Finding", "Intuition", "Result", "Reason"]
cmp_rows = [
    ["β₁",
     "Negative",
     "−6.02*** ✓",
     "Matched — supply growth does buy T-bills"],
    ["q* direction",
     "Low buffer = fragile",
     "Sign flip at 13% ✓",
     "Matched in direction"],
    ["q* significance",
     "Confirmed statistically",
     "p = 0.260 (suggestive)",
     "N = 51 too small for power"],
    ["Event study",
     "Clear CAR signal",
     "All insignificant",
     "FOMC contamination + N = 3 events"],
    ["SVB",
     "High buffer protects",
     "Spread −52 bps (flight-safety)",
     "Flight-to-safety, not stablecoin mechanism"],
    ["LUNA",
     "Low buffer causes widening",
     "Signal present but weak",
     "FOMC +50bps only 5 days before"],
]
col_w_cmp = [Inches(1.8), Inches(2.0), Inches(2.5), Inches(3.6)]
add_table_slide(slide12,
                Inches(0.3), Inches(1.05),
                Inches(10.2), Inches(3.5),
                cmp_rows, cmp_headers, col_w_cmp,
                font_size=10)

add_rect(slide12, Inches(0.3), Inches(4.7), Inches(10.2), Inches(0.85),
         0, 0x30, 0x87)
add_textbox(slide12, Inches(0.5), Inches(4.75),
            Inches(9.9), Inches(0.75),
            "The regression is our STRONGEST evidence.   "
            "The event study is directionally consistent but statistically underpowered.\n"
            "Three methods point the same direction — convergence without false certainty.",
            font_size=11, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 13 — Policy Implications
# ══════════════════════════════════════════════════════════════
slide13 = prs.slides.add_slide(blank)
add_rect(slide13, 0, 0, SLIDE_W, SLIDE_H, 0xF8, 0xF9, 0xFA)
add_header(slide13, "Policy Implications",
           "From empirical threshold to regulatory reference point",
           "Slide 13 / 16")
add_sidebar(slide13)

# Three policy boxes
pol_boxes = [
    ("MiCA (EU, 2024)",
     "Mandates liquid reserve requirements for 'significant' e-money tokens. "
     "Our q* = 13% provides an empirically grounded floor."),
    ("GENIUS Act (US, 2025)",
     "Requires stablecoin issuers to hold cash or near-cash reserves. "
     "Our CI [3.2%, 14.5%] maps to a regulatory range, not a precise floor."),
    ("Market trajectory",
     "Supply growing toward $1–2T. At that scale, forced liquidation risk "
     "is SYSTEMIC — comparable to a money-market fund break-the-buck event."),
]
for idx, (title, body) in enumerate(pol_boxes):
    bx = Inches(0.3) + idx * Inches(3.45)
    add_rect(slide13, bx, Inches(1.05), Inches(3.25), Inches(2.2),
             0, 0x30, 0x87)
    add_textbox(slide13, bx + Inches(0.15), Inches(1.12),
                Inches(2.95), Inches(0.45),
                title, font_size=11.5, bold=True, color=C_ORANGE)
    add_textbox(slide13, bx + Inches(0.15), Inches(1.60),
                Inches(2.95), Inches(1.5),
                body, font_size=10, color=C_WHITE)

# Direct recommendation
add_rect(slide13, Inches(0.3), Inches(3.45), Inches(10.2), Inches(1.2),
         0xFF, 0xF3, 0xE0)
add_textbox(slide13, Inches(0.45), Inches(3.50),
            Inches(9.9), Inches(0.4),
            "Direct Regulatory Reference", font_size=12, bold=True,
            color=C_ORANGE)
add_textbox(slide13, Inches(0.45), Inches(3.95),
            Inches(9.9), Inches(0.6),
            "Our estimate: require cash/near-cash ≥ 13% of supply outstanding.   "
            "Caveat: 90% CI [3.2%, 14.5%] — this is a range, not a precise floor.   "
            "Interpretation: anything below 3% is clearly under-reserved; above 14.5% is clearly adequate.",
            font_size=10.5, color=C_BLACK)

add_rect(slide13, Inches(0.3), Inches(4.78), Inches(10.2), Inches(0.85),
         0, 0x30, 0x87)
add_textbox(slide13, Inches(0.5), Inches(4.83),
            Inches(9.9), Inches(0.75),
            "First paper to provide an empirical threshold directly relevant to reserve legislation.\n"
            "As attestation data accumulates, this estimate will sharpen.",
            font_size=11, bold=True, color=C_WHITE)


# ══════════════════════════════════════════════════════════════
# SLIDE 14 — Summary
# ══════════════════════════════════════════════════════════════
slide14 = prs.slides.add_slide(blank)
add_rect(slide14, 0, 0, SLIDE_W, SLIDE_H, 0xF8, 0xF9, 0xFA)
add_header(slide14, "Summary: Three Convergent Lines of Evidence",
           "Each method is individually limited; together they point to the same conclusion",
           "Slide 14 / 16")
add_sidebar(slide14)

# Three-column layout
cols = [
    (YONSEI_BLUE,   "1.  Regression",
     "β₁ = −6.02***",
     "p = 0.006",
     "Privilege amplification confirmed.\n"
     "1-SD supply growth → 24 bps spread compression.\n"
     "VIX and rate controls included.\n"
     "HAC standard errors.\n\n"
     "STRONGEST evidence"),
    (RED,           "2.  Hansen Threshold",
     "q* = 13%",
     "Bootstrap CI [3.2%, 14.5%]",
     "Threshold identified independently by\nHansen grid search AND LSTAR (c* = 13.1%).\n"
     "Sign flip confirmed at 13%.\n"
     "p = 0.260 — suggestive, not confirmed.\n"
     "N = 51 limits power."),
    (GRAY,          "3.  Event Study",
     "All CARs insignificant",
     "−15 to −2 bps",
     "Directionally consistent.\n"
     "FOMC contamination in 2 of 3 events.\n"
     "SVB = flight-to-safety, not mechanism.\n"
     "LUNA = most relevant window.\n"
     "Qualitative context only."),
]
for idx, (col_hex, method, result, stat, detail) in enumerate(cols):
    cx = Inches(0.3) + idx * Inches(3.42)
    col_rgb = col_hex

    # header
    add_rect(slide14, cx, Inches(1.05), Inches(3.2), Inches(0.55),
             *[int(col_hex[i:i+2], 16) for i in (1, 3, 5)])
    add_textbox(slide14, cx + Inches(0.1), Inches(1.08),
                Inches(3.0), Inches(0.45),
                method, font_size=12, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)

    # result badge
    add_rect(slide14, cx, Inches(1.62), Inches(3.2), Inches(0.55),
             0xFF, 0xF3, 0xE0)
    add_textbox(slide14, cx + Inches(0.05), Inches(1.65),
                Inches(3.1), Inches(0.45),
                result, font_size=13, bold=True,
                color=RGBColor(*[int(col_hex[i:i+2], 16) for i in (1, 3, 5)]),
                align=PP_ALIGN.CENTER)

    # stat
    add_textbox(slide14, cx + Inches(0.05), Inches(2.22),
                Inches(3.1), Inches(0.32),
                stat, font_size=9.5, italic=True,
                color=RGBColor(*[int(col_hex[i:i+2], 16) for i in (1, 3, 5)]),
                align=PP_ALIGN.CENTER)

    # body
    add_rect(slide14, cx, Inches(2.58), Inches(3.2), Inches(2.65),
             0xF2, 0xF2, 0xF2)
    add_textbox(slide14, cx + Inches(0.1), Inches(2.65),
                Inches(3.0), Inches(2.5),
                detail, font_size=10, color=C_BLACK)

# Conclusion
add_rect(slide14, Inches(0.3), Inches(5.4), Inches(10.2), Inches(0.85),
         0, 0x30, 0x87)
add_textbox(slide14, Inches(0.5), Inches(5.45),
            Inches(9.9), Inches(0.75),
            "\"All three methods point to the same conclusion: stablecoins amplify the exorbitant "
            "privilege, with a fragility threshold near 13% liquid buffer.\"",
            font_size=12, bold=True, italic=True, color=C_ORANGE,
            align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 15 — Limitations & Next Steps
# ══════════════════════════════════════════════════════════════
slide15 = prs.slides.add_slide(blank)
add_rect(slide15, 0, 0, SLIDE_W, SLIDE_H, 0xF8, 0xF9, 0xFA)
add_header(slide15, "Limitations & Next Steps",
           "Honest accounting of constraints and the research agenda ahead",
           "Slide 15 / 16")
add_sidebar(slide15)

# Limitations
add_rect(slide15, Inches(0.3), Inches(1.05), Inches(4.7), Inches(4.5),
         0xFF, 0xED, 0xED)
add_textbox(slide15, Inches(0.45), Inches(1.10),
            Inches(4.4), Inches(0.45),
            "Limitations", font_size=13, bold=True, color=C_RED)

lims = [
    "N = 51 months limits statistical power\nfor threshold confirmation",
    "Liquid buffer L forward-filled for 13 months\n(Tether quarterly attestation lag)",
    "Endogeneity not fully resolved\n(Granger test addresses but doesn't solve)",
    "Only 3 stress events;\nFOMC contamination in 2 of 3",
    "Sample covers 2020–2024 only;\nmay not generalize to other regimes",
]
y_lim = Inches(1.65)
for lim in lims:
    add_rect(slide15, Inches(0.4), y_lim, Inches(0.22), Inches(0.22),
             0xC0, 0x39, 0x2B)
    add_textbox(slide15, Inches(0.7), y_lim - Inches(0.04),
                Inches(4.1), Inches(0.65),
                lim, font_size=10, color=C_BLACK)
    y_lim += Inches(0.7)

# Next Steps
add_rect(slide15, Inches(5.2), Inches(1.05), Inches(5.1), Inches(4.5),
         0xE8, 0xF5, 0xE9)
add_textbox(slide15, Inches(5.35), Inches(1.10),
            Inches(4.8), Inches(0.45),
            "Next Steps", font_size=13, bold=True, color=C_GREEN)

nexts = [
    ("Longer sample",
     "As weekly/monthly attestations accumulate, N grows → power improves"),
    ("Better identification",
     "Instrument for supply growth (e.g., regulatory announcements)"),
    ("Non-USD ecosystems",
     "Euro stablecoins, CBDC interactions — broader privilege question"),
    ("High-frequency data",
     "Daily spread + daily supply → better event identification"),
    ("Cross-country comparison",
     "Does same mechanism operate for EUR-denominated safe assets?"),
]
y_nx = Inches(1.65)
for title, desc in nexts:
    add_rect(slide15, Inches(5.3), y_nx, Inches(0.22), Inches(0.22),
             0x1A, 0x7A, 0x4A)
    add_textbox(slide15, Inches(5.6), y_nx - Inches(0.02),
                Inches(4.5), Inches(0.28),
                title, font_size=10.5, bold=True, color=C_GREEN)
    add_textbox(slide15, Inches(5.6), y_nx + Inches(0.25),
                Inches(4.5), Inches(0.32),
                desc, font_size=9.5, color=C_BLACK)
    y_nx += Inches(0.75)

add_rect(slide15, Inches(0.3), Inches(5.7), Inches(10.2), Inches(0.6),
         0, 0x30, 0x87)
add_textbox(slide15, Inches(0.5), Inches(5.75),
            Inches(9.9), Inches(0.5),
            "This paper is the first step. The research agenda is clear and tractable.",
            font_size=11, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 16 — Thank You
# ══════════════════════════════════════════════════════════════
slide16 = prs.slides.add_slide(blank)
add_rect(slide16, 0, 0, SLIDE_W, SLIDE_H, 0, 0x30, 0x87)
add_rect(slide16, 0, Inches(3.5), SLIDE_W, Inches(0.06), 0xE8, 0x77, 0x22)

add_textbox(slide16, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.75),
            "Stablecoins and the Exorbitant Privilege",
            font_size=30, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide16, Inches(0.8), Inches(1.65), Inches(11.7), Inches(0.5),
            "Safe-Asset Demand and Its Systemic Fragility",
            font_size=18, italic=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

# Three key numbers
kn_data = [
    ("β₁ = −6.02***", "p = 0.006", C_BLUE),
    ("q* = 13%",       "Bootstrap CI [3.2%, 14.5%]", RGBColor(0xC0, 0x39, 0x2B)),
    ("3 methods",      "regression + threshold + event study", RGBColor(0x1A, 0x7A, 0x4A)),
]
for idx, (main, sub, col) in enumerate(kn_data):
    kx = Inches(1.0) + idx * Inches(4.0)
    add_rect(slide16, kx, Inches(2.35), Inches(3.5), Inches(0.9),
             0, 0x20, 0x65)
    add_textbox(slide16, kx + Inches(0.1), Inches(2.38),
                Inches(3.3), Inches(0.45),
                main, font_size=17, bold=True, color=C_ORANGE,
                align=PP_ALIGN.CENTER)
    add_textbox(slide16, kx + Inches(0.1), Inches(2.8),
                Inches(3.3), Inches(0.35),
                sub, font_size=9.5, italic=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)

# Authors
add_textbox(slide16, Inches(0.8), Inches(3.75), Inches(11.7), Inches(0.5),
            "Mireu Kim   Sara Chekroune   Oybek Ibragimov   Jade Zhu",
            font_size=13, color=C_WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide16, Inches(0.8), Inches(4.28), Inches(11.7), Inches(0.5),
            "Alexandre Godefroy   Baptiste Degand   Minjin Kim",
            font_size=13, color=C_WHITE, align=PP_ALIGN.CENTER)

# Institution
add_textbox(slide16, Inches(0.8), Inches(4.95), Inches(11.7), Inches(0.45),
            "Yonsei GSIS  |  Topics in International Finance  |  June 2026",
            font_size=12, italic=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

# Thank you
add_textbox(slide16, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.6),
            "Thank you — Questions welcome",
            font_size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# GitHub
add_textbox(slide16, Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.4),
            "github.com/stablecoin-privilege-research",
            font_size=11, italic=True, color=C_ORANGE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════
out_path = PRESENTATIONS / "0609_Stablecoin_Exorbitant_Privilege.pptx"
prs.save(str(out_path))
print(f"\nPresentation saved to: {out_path}")
print(f"Slide count: {len(prs.slides)}")
