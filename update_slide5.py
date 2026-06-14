"""
update_slide5.py
Updates slide 5 (index 4, 4B: "Why? Diagnosing Spurious Regression"):
  - Replaces long paragraph rows with approved short bullet text
  - Shrinks cards to make room for a conclusion bar
  - Adds a dark conclusion bar at the bottom
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

TEAL   = RGBColor(0x1B, 0x99, 0x8B)
GOLD   = RGBColor(0xE8, 0xA1, 0x2C)
NAVY   = RGBColor(0x13, 0x29, 0x4B)
NAVYDK = RGBColor(0x12, 0x28, 0x44)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LTGRAY = RGBColor(0xF2, 0xF4, 0xF8)
RED    = RGBColor(0xC1, 0x12, 0x1F)
DKGRAY = RGBColor(0x33, 0x42, 0x5C)
MGRAY  = RGBColor(0x5B, 0x6B, 0x82)
SERIF  = "Georgia"
SANS   = "Calibri"

ANS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def get_shape(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None


def clear_tf(tf):
    txBody = tf._txBody
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)
    txBody.append(parse_xml(f'<a:p xmlns:a="{ANS}"/>'))


def fp(tf, text, font, sz, bold=False, color=NAVY,
       align=PP_ALIGN.LEFT, sb=0, sa=0):
    p = tf.paragraphs[0]
    p.alignment = align
    if sb: p.space_before = Pt(sb)
    if sa: p.space_after  = Pt(sa)
    r = p.add_run()
    r.text = text; r.font.name = font; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.color.rgb = color


def ap(tf, text, font, sz, bold=False, color=NAVY,
       align=PP_ALIGN.LEFT, sb=0, sa=0):
    p = tf.add_paragraph()
    p.alignment = align
    if sb: p.space_before = Pt(sb)
    if sa: p.space_after  = Pt(sa)
    r = p.add_run()
    r.text = text; r.font.name = font; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.color.rgb = color


def add_rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def add_tb(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


# ── New card dimensions ────────────────────────────────────────────────────────
NEW_CARD_H = 3500000    # shrunk from 4100000 to make room for conclusion bar
CARD_TOP   = 2011680

# conclusion bar sits just below the cards with a small gap
BAR_TOP    = CARD_TOP + NEW_CARD_H + 90000  # = 5601680
BAR_H      = 540000
BAR_L      = 548640
BAR_W      = 11091672

# ── Approved bullet content per card ──────────────────────────────────────────
CARDS = [
    dict(
        bg="Rectangle 8", accent="Rectangle 9",
        content_tb="TextBox 11",
        rows=[
            ("What it is",   TEAL, "Does the variable drift without a stable mean? If yes → non-stationary"),
            ("Why we ran it", TEAL, "Non-stationary variables produce spurious regressions — false significance from shared drift"),
            ("Our result",   GOLD, "Spread p = 0.494  ·  Buffer L p = 0.902 → both non-stationary, I(1)"),
            ("What it means", RED, "Both variables drifted down together — coincidence, not causation"),
        ]
    ),
    dict(
        bg="Rectangle 12", accent="Rectangle 13",
        content_tb="TextBox 15",
        rows=[
            ("What it is",   TEAL, "Even if non-stationary individually, a stable long-run link (cointegration) makes regression valid"),
            ("Why we ran it", TEAL, "Standard rescue — if cointegration holds, levels regression survives"),
            ("Our result",   GOLD, "p = 0.120 → cannot reject no cointegration"),
            ("What it means", RED, "No stable long-run link — they fell together during 2022–24 hiking, then diverged"),
        ]
    ),
    dict(
        bg="Rectangle 16", accent="Rectangle 17",
        content_tb="TextBox 19",
        rows=[
            ("What it is",   TEAL, "Stronger, multivariate cointegration test — checks the full system, not one pair"),
            ("Why we ran it", TEAL, "More reliable in small samples — the definitive confirmation"),
            ("Our result",   GOLD, "Fails to reject r = 0 (zero cointegrating vectors)"),
            ("What it means", RED, "Both tests fail → levels regression is spurious"),
        ]
    ),
]

CONCLUSION = (
    "All three tests agree: the variables are non-stationary and not cointegrated. "
    "The original results were driven by shared drift during the 2022–24 Fed hiking cycle "
    "— not a genuine economic relationship."
)

# ══════════════════════════════════════════════════════════════════════════════
prs = Presentation("presentations/0616_Stablecoin_Exorbitant_Privilege.pptx")
s5  = prs.slides[4]

# Shrink each card background and accent bar to NEW_CARD_H
for card in CARDS:
    bg_sh = get_shape(s5, card["bg"])
    ac_sh = get_shape(s5, card["accent"])
    bg_sh.height = NEW_CARD_H
    ac_sh.height = NEW_CARD_H

    # Replace content text box text
    tb_sh = get_shape(s5, card["content_tb"])
    tb_sh.height = NEW_CARD_H - 640000
    tf = tb_sh.text_frame
    tf.word_wrap = True
    clear_tf(tf)

    first = True
    for label, label_color, body in card["rows"]:
        if first:
            fp(tf, label, SANS, 10.5, bold=True, color=label_color, sa=2)
            first = False
        else:
            ap(tf, label, SANS, 10.5, bold=True, color=label_color, sb=11, sa=2)
        ap(tf, body, SANS, 10, color=DKGRAY)

# Add conclusion bar
add_rect(s5, BAR_L, BAR_TOP, BAR_W, BAR_H, NAVYDK)
tf = add_tb(s5, BAR_L + 200000, BAR_TOP + 80000, BAR_W - 300000, BAR_H - 100000)
fp(tf, CONCLUSION, SANS, 11, color=WHITE, align=PP_ALIGN.LEFT)

prs.save("presentations/0616_Stablecoin_Exorbitant_Privilege.pptx")
print("Done — slide 5 updated with bullets and conclusion bar.")
