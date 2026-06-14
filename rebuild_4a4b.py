"""
rebuild_4a4b.py
Modifies presentations/0616_Stablecoin_Exorbitant_Privilege.pptx:
  - Slide 4  → 4A: "What the Re-Run Revealed"
              Left card: β₁ and β₄ bullet explanations
              Dark bar: "What we dropped" kept as-is
  - Insert new slide 4B after: "Why? Diagnosing Spurious Regression"
              Three column cards — ADF / Engle-Granger / Johansen
  - Renumber slides 5-15 → 6-16
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# ── Palette ───────────────────────────────────────────────────────────────────
TEAL   = RGBColor(0x1B, 0x99, 0x8B)
GOLD   = RGBColor(0xE8, 0xA1, 0x2C)
NAVY   = RGBColor(0x13, 0x29, 0x4B)
NAVYDK = RGBColor(0x12, 0x28, 0x44)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LTGRAY = RGBColor(0xF2, 0xF4, 0xF8)
RED    = RGBColor(0xC1, 0x12, 0x1F)
DKGRAY = RGBColor(0x33, 0x42, 0x5C)
MGRAY  = RGBColor(0x5B, 0x6B, 0x82)
SERIF  = "Georgia"
SANS   = "Calibri"

ANS = "http://schemas.openxmlformats.org/drawingml/2006/main"

# ── Helpers ───────────────────────────────────────────────────────────────────
def get_shape(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None

def clear_tf(tf):
    txBody = tf._txBody
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)
    txBody.append(parse_xml(f'<a:p xmlns:a="{ANS}"/>'))

def fp(tf, text, font, sz, bold=False, color=NAVY,
       align=PP_ALIGN.LEFT, sb=0, sa=0):
    """Set text on the FIRST paragraph (use after clear_tf or fresh textbox)."""
    p = tf.paragraphs[0]
    p.alignment = align
    if sb: p.space_before = Pt(sb)
    if sa: p.space_after  = Pt(sa)
    r = p.add_run()
    r.text = text;  r.font.name = font;  r.font.size = Pt(sz)
    r.font.bold = bold;  r.font.color.rgb = color

def ap(tf, text, font, sz, bold=False, color=NAVY,
       align=PP_ALIGN.LEFT, sb=0, sa=0):
    """Add a NEW paragraph."""
    p = tf.add_paragraph()
    p.alignment = align
    if sb: p.space_before = Pt(sb)
    if sa: p.space_after  = Pt(sa)
    r = p.add_run()
    r.text = text;  r.font.name = font;  r.font.size = Pt(sz)
    r.font.bold = bold;  r.font.color.rgb = color

def add_rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid();  s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def add_tb(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame

def build_header(slide, kicker, title, subtitle, num, H):
    add_rect(slide, 0,      0, 182880, H, NAVY)
    add_rect(slide, 182880, 0,  91440, H, TEAL)

    tf = add_tb(slide, 566928, 384048, 11064240, 292608)
    fp(tf, kicker, SANS, 12, bold=True, color=TEAL)

    tf = add_tb(slide, 548640, 676656, 11064240, 777240)
    fp(tf, title, SERIF, 27, bold=True, color=NAVY)

    tf = add_tb(slide, 566928, 1481328, 11064240, 411480)
    fp(tf, subtitle, SERIF, 15.5, color=MGRAY)

    tf = add_tb(slide, 566928, 6455664, 8229600, 274320)
    fp(tf, "Stablecoins & the Exorbitant Privilege", SANS, 9, color=MGRAY)

    tf = add_tb(slide, 10972800, 6400800, 640080, 274320)
    fp(tf, str(num), SANS, 11, color=MGRAY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
prs = Presentation("presentations/0616_Stablecoin_Exorbitant_Privilege.pptx")
W, H = prs.slide_width, prs.slide_height

# ── SLIDE 4A: modify in place ─────────────────────────────────────────────────
s4 = prs.slides[3]

# Remove the "diagnosis" card (right side)
for name in ["Rectangle 10", "Rectangle 11", "TextBox 12", "TextBox 13"]:
    sh = get_shape(s4, name)
    if sh:
        s4.shapes._spTree.remove(sh._element)

# Extend card background to full width
get_shape(s4, "Rectangle 6").width = 11091672

# Extend card title width
get_shape(s4, "TextBox 8").width = 10700000

# Replace card body with β₁ / β₄ bullets
cb = get_shape(s4, "TextBox 9")
cb.width  = 10700000
cb.height = 1600000
tf = cb.text_frame
tf.word_wrap = True
clear_tf(tf)

fp(tf,
   "β₁ — Privilege Amplification",
   SANS, 13, bold=True, color=NAVY, sa=3)
ap(tf,
   "Measures whether stablecoin supply amplifies how sensitive spreads are to reserve buffers. "
   "After the corrected re-run: β₁ flipped from −6.02*** to +2.74 (not significant). "
   "We cannot claim that stablecoin growth makes T-bill demand more fragile under stress.",
   SANS, 11.5, color=DKGRAY, sa=14)
ap(tf,
   "β₄ — Buffer Interaction",
   SANS, 13, bold=True, color=NAVY, sa=3)
ap(tf,
   "Measures the direct effect of reserve buffers on spreads. In levels: −35.89 (p = 0.032). "
   "Under first-differencing: +8.86 (p = 0.31) — sign flipped and became insignificant. "
   "A sign reversal under differencing is the textbook signal of spurious regression.",
   SANS, 11.5, color=DKGRAY)

# ── SLIDE 4B: insert after slide 4 ───────────────────────────────────────────
blank_layout = prs.slide_layouts[6]
new_slide = prs.slides.add_slide(blank_layout)

# Move new slide from end → position index 4 (right after 4A)
sld_id_lst = prs.slides._sldIdLst
items = list(sld_id_lst)
el = items[-1]
sld_id_lst.remove(el)
sld_id_lst.insert(4, el)

s5 = prs.slides[4]  # 4B is now at index 4

# Remove any placeholders that the blank layout injected
sp_tree = s5.shapes._spTree
for sp in list(sp_tree):
    tag = sp.tag.split("}")[-1] if "}" in sp.tag else sp.tag
    if tag == "sp":
        sp_tree.remove(sp)

build_header(s5,
             "AFTER THE PROFESSOR'S FEEDBACK",
             "Why? Diagnosing Spurious Regression",
             "We ran three standard tests to confirm the levels regression is not valid",
             5, H)

# ── Three test cards ──────────────────────────────────────────────────────────
CARD_TOP = 2011680
CARD_H   = 4100000        # tall — no dark bar on this slide

TOTAL_W  = 11091672       # from left=548640 to right=11640312
GAP      = 200000
CWIDTH   = (TOTAL_W - 2 * GAP) // 3
STARTS   = [548640,
            548640 + CWIDTH + GAP,
            548640 + 2 * (CWIDTH + GAP)]
ACCENT_W = 73152
TXT_OFF  = ACCENT_W + 109728    # text starts after accent bar + padding
TXT_W    = CWIDTH - TXT_OFF - 91440

TESTS = [
    dict(accent=NAVY, title="ADF Unit Root Test",
         rows=[
             ("What it is", TEAL,
              "Tests whether a variable wanders without a stable mean — this is called "
              "non-stationary. A non-stationary variable has no equilibrium; it simply drifts over time."),
             ("Why we ran it", TEAL,
              "If both the spread and reserve buffer L are non-stationary, a regression between them "
              "can look significant purely from shared drift — not because one causes the other."),
             ("Our result", GOLD,
              "Spread p = 0.494  ·  Buffer L p = 0.902\n→ Both non-stationary, I(1)"),
             ("What it means", RED,
              "Two drifting variables can produce high R² and low p-values by coincidence. "
              "The regression may be picking up shared drift, not a real economic relationship."),
         ]),
    dict(accent=TEAL, title="Engle–Granger Cointegration",
         rows=[
             ("What it is", TEAL,
              "Even if two variables are individually non-stationary, their regression is still valid "
              "if they share a stable long-run relationship — called cointegration."),
             ("Why we ran it", TEAL,
              "Cointegration is the standard rescue for two non-stationary variables. "
              "If it holds, the levels regression remains valid despite non-stationarity."),
             ("Our result", GOLD,
              "p = 0.120 — cannot reject no cointegration"),
             ("What it means", RED,
              "No stable long-run link. Both variables fell during the 2022–24 Fed hiking cycle "
              "and diverged after. Coincidental co-movement — not a structural relationship."),
         ]),
    dict(accent=GOLD, title="Johansen Test",
         rows=[
             ("What it is", TEAL,
              "A more powerful, multivariate version of the cointegration test. "
              "It examines the full system of variables together — not just one pair at a time."),
             ("Why we ran it", TEAL,
              "Robustness check: Engle–Granger can miss cointegration in small samples. "
              "Johansen is the stronger and more reliable confirmation."),
             ("Our result", GOLD,
              "Fails to reject r = 0 (zero cointegrating vectors)"),
             ("What it means", RED,
              "Confirms Engle–Granger. With both tests failing, the levels regression is spurious "
              "and cannot support any causal claim about reserve buffers and spreads."),
         ]),
]

for i, test in enumerate(TESTS):
    cx = STARTS[i]
    tx = cx + TXT_OFF

    # Card background + accent bar
    add_rect(s5, cx, CARD_TOP, CWIDTH,   CARD_H, LTGRAY)
    add_rect(s5, cx, CARD_TOP, ACCENT_W, CARD_H, test["accent"])

    # Card title
    tf = add_tb(s5, tx, CARD_TOP + 114000, TXT_W, 400000)
    fp(tf, test["title"], SERIF, 13, bold=True, color=NAVY)

    # Content rows (all in one text box)
    tf = add_tb(s5, tx, CARD_TOP + 570000, TXT_W, CARD_H - 640000)
    first = True
    for label, label_color, body in test["rows"]:
        if first:
            fp(tf, label, SANS, 10, bold=True, color=label_color, sa=2)
            first = False
        else:
            ap(tf, label, SANS, 10, bold=True, color=label_color, sb=10, sa=2)
        ap(tf, body, SANS, 9.5, color=DKGRAY)

# ── Renumber slides 5-15 → 6-16 ──────────────────────────────────────────────
for new_idx in range(5, len(prs.slides)):
    slide = prs.slides[new_idx]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # slide number boxes sit at top > 6 000 000 and left > 10 000 000
        if shape.top > 6000000 and shape.left > 10000000:
            tf = shape.text_frame
            txt = tf.paragraphs[0].text.strip()
            if txt.isdigit():
                clear_tf(tf)
                fp(tf, str(new_idx + 1), SANS, 11, color=MGRAY, align=PP_ALIGN.RIGHT)

prs.save("presentations/0616_Stablecoin_Exorbitant_Privilege.pptx")
print("Done — 16-slide deck saved.")
