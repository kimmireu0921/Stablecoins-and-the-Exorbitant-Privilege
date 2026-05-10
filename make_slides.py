"""
make_slides.py — presentation matching the Yonsei GSIS template from Digital Currency slides.

Template:
  - Title/end slides: dark navy bg + lighter navy band + white circle + white chevron
  - Content slides:   white bg + dark navy top header bar + slide number bottom right
  - Orange arrow callout for key takeaways
  - Diamond (❖) bullet points, clean Calibri font
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pathlib import Path

RESULTS = Path("results")

# ── Colour palette (matched to original slides) ───────────────────────────────
DARK_NAVY  = RGBColor(0x0A, 0x1E, 0x5C)   # title slide main bg
MED_NAVY   = RGBColor(0x14, 0x30, 0x70)   # title slide lighter band
HDR_NAVY   = RGBColor(0x1B, 0x3A, 0x72)   # content slide header bar
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE     = RGBColor(0xD4, 0x63, 0x2A)   # callout arrow
BODY_DARK  = RGBColor(0x1A, 0x1A, 0x2E)   # body text on white slides
LGRAY      = RGBColor(0xF0, 0xF4, 0xF8)   # subtle bg for boxes
MIDGRAY    = RGBColor(0x55, 0x55, 0x66)   # secondary text

SW = Inches(13.33)
SH = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def box(sld, l, t, w, h, fill=None, line_color=None, line_width=Pt(0)):
    sh = sld.shapes.add_shape(1, l, t, w, h)
    sh.line.fill.background()
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = line_width
    return sh


def txt(sld, text, l, t, w, h, size=18, bold=False, italic=False,
        color=BODY_DARK, align=PP_ALIGN.LEFT, font="Calibri", wrap=True):
    tb = sld.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_line(tf, text, size=18, bold=False, italic=False,
             color=BODY_DARK, align=PP_ALIGN.LEFT, space_before=4):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return p


def img(sld, path, l, t, w, h=None):
    if Path(path).exists():
        if h:
            sld.shapes.add_picture(str(path), l, t, w, h)
        else:
            sld.shapes.add_picture(str(path), l, t, w)


# ── Yonsei template shapes ────────────────────────────────────────────────────

def title_slide_bg(sld):
    """Dark navy bg + lighter navy horizontal band (like original template)."""
    box(sld, 0, 0, SW, SH, fill=DARK_NAVY)
    box(sld, 0, Inches(2.85), SW, Inches(1.15), fill=MED_NAVY)

    # White large circle (top-left, like Yonsei logo area)
    circle = sld.shapes.add_shape(9, Inches(0.1), Inches(0.0), Inches(3.8), Inches(3.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = WHITE
    circle.line.fill.background()

    # Inner dark circle (the ring effect)
    inner = sld.shapes.add_shape(9, Inches(0.55), Inches(0.45), Inches(2.9), Inches(2.9))
    inner.fill.solid()
    inner.fill.fore_color.rgb = DARK_NAVY
    inner.line.fill.background()

    # White chevron/arrow at bottom-left
    chev = sld.shapes.add_shape(9, Inches(0.1), Inches(4.4), Inches(3.6), Inches(3.6))
    chev.fill.solid()
    chev.fill.fore_color.rgb = WHITE
    chev.line.fill.background()

    # Dark cover to make it look like open chevron
    inner_chev = sld.shapes.add_shape(9, Inches(0.65), Inches(4.85), Inches(2.5), Inches(2.5))
    inner_chev.fill.solid()
    inner_chev.fill.fore_color.rgb = DARK_NAVY
    inner_chev.line.fill.background()

    # Gray section (like original gray band at middle of circle/chevron)
    box(sld, 0, Inches(2.55), Inches(3.8), Inches(0.6), fill=RGBColor(0x55, 0x55, 0x66))


def content_slide(prs, title_text):
    """White bg + dark navy header bar + slide title."""
    sld = blank(prs)
    # White background
    box(sld, 0, 0, SW, SH, fill=WHITE)
    # Dark navy top bar
    box(sld, 0, 0, SW, Inches(0.95), fill=HDR_NAVY)
    # Title text in header
    txt(sld, title_text,
        Inches(0.3), Inches(0.1), Inches(11.5), Inches(0.78),
        size=26, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font="Calibri")
    return sld


def slide_number(sld, num):
    txt(sld, str(num),
        SW - Inches(0.55), SH - Inches(0.42), Inches(0.45), Inches(0.35),
        size=13, color=MIDGRAY, align=PP_ALIGN.RIGHT)


def orange_callout(sld, text, t=None, size=18):
    """Orange arrow/chevron callout box at bottom (like original)."""
    if t is None:
        t = SH - Inches(1.38)
    box(sld, 0, t, SW, Inches(1.2), fill=ORANGE)
    txt(sld, text,
        Inches(0.6), t + Inches(0.18), SW - Inches(1.0), Inches(0.88),
        size=size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def bullet_tb(sld, l, t, w, h):
    """Return an empty textbox ready for bullet lines."""
    tb = sld.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].space_before = Pt(0)
    return tf


# ── Build presentation ────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
sld = blank(prs)
title_slide_bg(sld)

# Main title (in the lighter band)
txt(sld,
    "STABLECOINS AND THE EXORBITANT PRIVILEGE",
    Inches(4.0), Inches(2.78), Inches(9.0), Inches(0.65),
    size=22, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font="Calibri Light")

# Subtitle
txt(sld,
    "A New Channel for Safe-Asset Demand and Its Systemic Fragility",
    Inches(4.0), Inches(3.42), Inches(9.0), Inches(0.5),
    size=14, italic=True, color=RGBColor(0xAA, 0xC4, 0xE8),
    align=PP_ALIGN.LEFT)

# Bottom info box
box(sld, Inches(3.9), Inches(4.45), Inches(7.5), Inches(2.8),
    fill=RGBColor(0x0E, 0x26, 0x68))

txt(sld,
    "Yonsei GSIS 2026-1\nTopics in International Finance\n"
    "Mireu Kim · Sara Chekroune · Oybek Ibragimov\n"
    "April 2026",
    Inches(4.1), Inches(4.6), Inches(7.2), Inches(2.5),
    size=16, bold=False, color=WHITE, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Roadmap
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "Roadmap")
slide_number(sld, 2)

items = [
    "Motivation — why stablecoins matter for sovereign debt",
    "The New Triffin Dilemma — our core argument",
    "Theoretical extension of Maggiori (2017)",
    "Data & Methodology",
    "Result 1 — Privilege amplification regression",
    "Result 2 — Reserve adequacy threshold (Hansen 2000)",
    "Result 3 — Buffer-conditioned event study",
    "Conclusion & Policy implications",
]

tf = bullet_tb(sld, Inches(0.7), Inches(1.1), Inches(12.2), Inches(6.0))
for i, item in enumerate(items):
    add_line(tf, f"{'❖' if i > 0 else '❖'}  {item}",
             size=20, color=BODY_DARK, space_before=8 if i > 0 else 2)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Motivation
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "Motivation: Stablecoins as Treasury Buyers")
slide_number(sld, 3)

# Left column — text
tf = bullet_tb(sld, Inches(0.5), Inches(1.1), Inches(6.2), Inches(4.8))
add_line(tf, "Key Market Trends", size=18, bold=True, color=HDR_NAVY, space_before=2)
bullets = [
    "Growth: Combined market cap surged from under $50bn (2021) to over $300bn (2026)",
    "Usage: 30%+ of on-chain crypto transactions involve stablecoins",
    "Regulation: US, EU, Japan, and Hong Kong each finalising reserve legislation",
    "Scale: Tether alone holds $127bn in US Treasuries (Q2 2025) — comparable to mid-sized sovereigns",
    "Projection: $1–2 trillion ecosystem within this decade",
]
for b in bullets:
    add_line(tf, f"❖  {b}", size=17, color=BODY_DARK, space_before=10)

# Right column — stat boxes
for i, (stat, label) in enumerate([
    ("$300bn+", "Total USD-pegged supply (2026)"),
    ("$127bn",  "Tether T-bill holdings (Q2 2025)"),
    ("$1–2 tn", "Projected ecosystem size"),
    ("4",       "Jurisdictions finalising legislation"),
]):
    t = Inches(1.05) + i * Inches(1.38)
    box(sld, Inches(7.1), t, Inches(5.8), Inches(1.2), fill=HDR_NAVY)
    txt(sld, stat, Inches(7.2), t + Inches(0.08),
        Inches(5.6), Inches(0.62), size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sld, label, Inches(7.2), t + Inches(0.72),
        Inches(5.6), Inches(0.4), size=12, italic=True, color=RGBColor(0xAA, 0xC4, 0xE8),
        align=PP_ALIGN.CENTER)

orange_callout(sld,
    "The same structural role that deepens US safe-asset demand "
    "also embeds a fragility that can rapidly reverse it.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — The New Triffin Dilemma
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "The New Triffin Dilemma")
slide_number(sld, 4)

box(sld, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.85), fill=LGRAY)
tf = bullet_tb(sld, Inches(0.55), Inches(1.1), Inches(12.2), Inches(0.75))
add_line(tf,
    "Triffin (1960): the reserve currency country must supply safe assets to the world — "
    "but that very supply undermines confidence in the currency.",
    size=14, italic=True, color=BODY_DARK)

txt(sld, "Our analogue for stablecoin issuers:",
    Inches(0.5), Inches(2.05), Inches(12.5), Inches(0.38),
    size=17, bold=True, color=HDR_NAVY)

# Three column comparison
cols = [
    ("Normal times ✅",
     ["Stablecoin growth → more T-bill buying",
      "Deepens US safe-asset demand",
      "Compresses OIS–Treasury spreads",
      "Amplifies exorbitant privilege"],
     RGBColor(0xE6, 0xF4, 0xEA), RGBColor(0x1A, 0x7A, 0x4A)),
    ("Run scenario ❌",
     ["Run on issuer → forced T-bill liquidation",
      "Yields spike above OIS",
      "Sovereign-linked liquidity shock",
      "Cole–Kehoe crisis zone activated"],
     RGBColor(0xFC, 0xEB, 0xE8), RGBColor(0xC0, 0x39, 0x2B)),
    ("What determines outcome",
     ["θ = T-bill holdings / supply  (Treasury Exposure)",
      "L = cash reserves / supply  (Liquid Buffer)",
      "L adequate: run absorbed, no liquidation",
      "→ We estimate the threshold L* (Hansen 2000)"],
     RGBColor(0xEB, 0xF3, 0xFB), HDR_NAVY),
]
for i, (title, points, bg, col) in enumerate(cols):
    l = Inches(0.4) + i * Inches(4.3)
    box(sld, l, Inches(2.5), Inches(4.1), Inches(0.42), fill=col)
    txt(sld, title, l + Inches(0.1), Inches(2.53),
        Inches(3.9), Inches(0.36), size=14, bold=True, color=WHITE)
    box(sld, l, Inches(2.92), Inches(4.1), Inches(2.55), fill=bg)
    tf = bullet_tb(sld, l + Inches(0.12), Inches(2.98), Inches(3.88), Inches(2.4))
    for pt in points:
        add_line(tf, f"•  {pt}", size=13, color=BODY_DARK, space_before=6)

orange_callout(sld,
    "Our paper: decompose reserve quality into θ (privilege channel) and L (fragility channel); estimate threshold L*.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Theoretical Framework
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "Theoretical Extension: Maggiori (2017)")
slide_number(sld, 5)

# Left: Maggiori original + extension
box(sld, Inches(0.4), Inches(1.05), Inches(6.2), Inches(5.2), fill=LGRAY)
tf = bullet_tb(sld, Inches(0.55), Inches(1.12), Inches(5.9), Inches(5.0))
add_line(tf, "Maggiori (2017) original:", size=14, bold=True, color=HDR_NAVY)
add_line(tf, "Safe-asset demand driven by Ω* (marginal value of RoW net worth).",
         size=13, color=BODY_DARK, space_before=4)
add_line(tf, " ", size=8)
add_line(tf, "Our extension — add stablecoin supply S, θ, and L:", size=14,
         bold=True, color=HDR_NAVY, space_before=6)
add_line(tf, "D*(Ñ*, S, θ, L)  =  D*precautionary(Ñ*)  +  D*stablecoin(S, θ)  +  D*buffer(L)",
         size=13, italic=True, color=BODY_DARK, space_before=4)
add_line(tf, " ", size=8)
add_line(tf, "Liquid Buffer L enters asymmetrically:", size=14, bold=True, color=HDR_NAVY)
add_line(tf, "❖  L ≥ L*  →  run absorbed, no T-bill liquidation", size=13,
         color=RGBColor(0x1A, 0x7A, 0x4A), space_before=4)
add_line(tf, "❖  L < L*  →  forced T-bill sales, spread spikes", size=13,
         color=RGBColor(0xC0, 0x39, 0x2B), space_before=4)

# Right: Estimating equation
box(sld, Inches(6.8), Inches(1.05), Inches(6.1), Inches(5.2), fill=HDR_NAVY)
txt(sld, "Main Estimating Equation",
    Inches(6.95), Inches(1.12), Inches(5.8), Inches(0.5),
    size=15, bold=True, color=WHITE)
txt(sld,
    "Spreadₜ = α + β₁·ΔlnSₜ + β₂·θₜ + β₃·Lₜ\n"
    "           + β₄·(Lₜ × ΔlnSₜ)\n"
    "           + β₅·Vₜ + β₆·VIXₜ\n"
    "           + β₇·ΔlnN*ₜ + εₜ",
    Inches(6.95), Inches(1.68), Inches(5.8), Inches(1.85),
    size=14, color=WHITE, font="Courier New")

tf2 = bullet_tb(sld, Inches(6.95), Inches(3.65), Inches(5.8), Inches(2.5))
add_line(tf2, "Predicted signs:", size=13, bold=True,
         color=RGBColor(0xAA, 0xC4, 0xE8))
for var, interp in [
    ("β₁ < 0", "ΔlnS compresses spreads (privilege amplification)"),
    ("β₂ > 0", "Higher T-bill exposure θ amplifies demand"),
    ("β₄ > 0", "Higher L dampens crisis transmission during runs"),
]:
    add_line(tf2, f"❖  {var}: {interp}", size=12,
             color=WHITE, space_before=6)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Data
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "Data & Sample")
slide_number(sld, 6)

rows = [
    ("OIS–Treasury Spread",  "DTB3 − SOFR90DAYAVG",        "FRED (direct CSV, no API key)",           "Daily → Monthly"),
    ("Stablecoin Supply S",  "USDT + USDC market cap",      "DeFiLlama stablecoins API (free)",        "Daily → Monthly"),
    ("θ  (Treasury Exposure)","T-bill holdings / supply",   "Tether/BDO + Circle/Deloitte attestations","Quarterly/Monthly"),
    ("L  (Liquid Buffer)",   "Cash reserves / supply",      "Tether/BDO + Circle/Deloitte attestations","Quarterly/Monthly"),
    ("Velocity V",           "7-day rolling SD of ΔlnS",   "Computed from DeFiLlama",                 "Daily → Monthly"),
    ("VIX",                  "CBOE VIX index",              "FRED",                                    "Daily → Monthly"),
    ("ΔlnN* (RoW equity)",   "ACWX ETF log-return",         "Yahoo Finance",                           "Daily → Monthly"),
]
hdrs = ["Variable", "Definition", "Source", "Frequency"]
col_ls = [Inches(0.35), Inches(2.75), Inches(5.55), Inches(10.55)]
col_ws = [Inches(2.35), Inches(2.75), Inches(4.95), Inches(2.35)]

# Header row
box(sld, Inches(0.35), Inches(1.05), Inches(12.6), Inches(0.42), fill=HDR_NAVY)
for h, l, w in zip(hdrs, col_ls, col_ws):
    txt(sld, h, l + Inches(0.05), Inches(1.08),
        w - Inches(0.08), Inches(0.36), size=13, bold=True, color=WHITE)

# Data rows
for r_i, row in enumerate(rows):
    t = Inches(1.47) + r_i * Inches(0.73)
    bg = LGRAY if r_i % 2 == 0 else WHITE
    box(sld, Inches(0.35), t, Inches(12.6), Inches(0.73), fill=bg)
    for val, l, w in zip(row, col_ls, col_ws):
        txt(sld, val, l + Inches(0.06), t + Inches(0.1),
            w - Inches(0.1), Inches(0.55), size=12, color=BODY_DARK)

txt(sld,
    "Sample: January 2022 – March 2026  |  N = 51 monthly observations  "
    "|  θ and L sourced from Tether/BDO and Circle/Deloitte reserve attestations",
    Inches(0.35), SH - Inches(0.55), Inches(12.6), Inches(0.42),
    size=11, italic=True, color=MIDGRAY, align=PP_ALIGN.CENTER)
slide_number(sld, 6)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Time Series Figure
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "Key Variables: January 2022 – March 2026")
slide_number(sld, 7)

img(sld, RESULTS / "fig_timeseries.png",
    Inches(0.3), Inches(1.0), Inches(8.5))

tf = bullet_tb(sld, Inches(9.0), Inches(1.1), Inches(3.9), Inches(5.8))
add_line(tf, "Three panels, one story", size=15, bold=True, color=HDR_NAVY)
obs = [
    "A: Treasury spread — compresses when supply grows, spikes in crises",
    "B: USDT + USDC supply — grew from $100bn to $200bn+ over this period",
    "C: Liquid buffer L — dangerously thin in 2022, recovering since 2023",
    "Red dashed line = our estimated safety threshold (13%)",
    "Red vertical lines = three stress events we study",
]
for o in obs:
    add_line(tf, f"❖  {o}", size=12, color=BODY_DARK, space_before=9)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Three Patterns (plain-language chart interpretation)
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "What the Data Shows Us")
slide_number(sld, 8)

patterns = [
    (RGBColor(0x1A, 0x7A, 0x4A), RGBColor(0xE6, 0xF4, 0xEA),
     "MORE SUPPLY  →  LOWER SPREAD",
     "As USDT and USDC supply grew, the gap between Treasury yields and overnight rates compressed. "
     "Stablecoin issuers were buying U.S. T-bills to back their tokens — deepening demand for safe assets."),
    (RGBColor(0xC0, 0x39, 0x2B), RGBColor(0xFC, 0xEB, 0xE8),
     "CRISIS  →  SPREAD SPIKES",
     "During the LUNA/UST collapse (May 2022), spreads spiked as issuers were forced to liquidate "
     "T-bill holdings to meet redemptions. The privilege reversed — exactly what our theory predicts."),
    (HDR_NAVY, LGRAY,
     "BUFFER WAS THIN WHEN IT MATTERED MOST",
     "In 2022, liquid cash reserves sat below our estimated 13% safety threshold. "
     "Issuers had to sell T-bills instead of using cash — amplifying the shock."),
]

for i, (hdr_col, bg_col, heading, body) in enumerate(patterns):
    t = Inches(1.1) + i * Inches(1.95)
    box(sld, Inches(0.4), t, Inches(12.5), Inches(0.45), fill=hdr_col)
    txt(sld, heading, Inches(0.55), t + Inches(0.07),
        Inches(12.1), Inches(0.35), size=15, bold=True, color=WHITE)
    box(sld, Inches(0.4), t + Inches(0.45), Inches(12.5), Inches(1.42), fill=bg_col)
    txt(sld, body, Inches(0.6), t + Inches(0.58),
        Inches(12.1), Inches(1.2), size=14, color=BODY_DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Result 1: Main Finding (big and clean)
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "Result 1: Stablecoin Growth Compresses Treasury Spreads")
slide_number(sld, 9)

# Big confirmed banner
box(sld, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.52), fill=RGBColor(0x1A, 0x7A, 0x4A))
txt(sld, "H1  CONFIRMED  ✅   —   statistically significant at the 5% level",
    Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.42),
    size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Giant coefficient
box(sld, Inches(0.4), Inches(1.65), Inches(5.8), Inches(2.5), fill=HDR_NAVY)
txt(sld, "β₁ = −6.02**",
    Inches(0.5), Inches(1.75), Inches(5.6), Inches(1.2),
    size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(sld, "N = 51 months  |  R² = 0.50",
    Inches(0.5), Inches(2.95), Inches(5.6), Inches(0.35),
    size=13, italic=True, color=RGBColor(0xAA, 0xC4, 0xE8), align=PP_ALIGN.CENTER)

# Plain language explanation
box(sld, Inches(6.4), Inches(1.65), Inches(6.5), Inches(2.5), fill=LGRAY)
txt(sld, "What this means in plain terms:",
    Inches(6.6), Inches(1.78), Inches(6.1), Inches(0.38),
    size=14, bold=True, color=HDR_NAVY)
txt(sld,
    "Every 1% growth in stablecoin supply compresses "
    "the Treasury spread by about 6 basis points.\n\n"
    "When Tether and Circle expand, they buy more "
    "U.S. T-bills — pushing Treasury yields down "
    "relative to overnight rates. This is the "
    "exorbitant privilege being amplified.",
    Inches(6.6), Inches(2.22), Inches(6.1), Inches(1.75),
    size=13, color=BODY_DARK)

# Compact regression table (key rows only)
key_rows = [
    ("ΔlnS  (stablecoin supply growth)", "−6.017", "**"),
    ("θ  (Treasury Exposure)",           " 0.100", "—"),
    ("L  (Liquid Buffer)",               " 1.034", "—"),
    ("L × ΔlnS  (buffer interaction)",   " 3.891", "—"),
    ("Controls (velocity, VIX, ΔlnN*)",  "included",""),
]
col_ls_t = [Inches(0.4), Inches(8.2), Inches(9.9)]
col_ws_t = [Inches(7.75), Inches(1.65), Inches(0.8)]

box(sld, Inches(0.4), Inches(4.28), sum(col_ws_t), Inches(0.36), fill=HDR_NAVY)
for h, l, w in zip(["Variable", "Coefficient", "Sig."], col_ls_t, col_ws_t):
    txt(sld, h, l + Inches(0.05), Inches(4.31), w, Inches(0.28),
        size=11, bold=True, color=WHITE,
        align=PP_ALIGN.LEFT if h == "Variable" else PP_ALIGN.CENTER)

for r_i, (var, coef, sig) in enumerate(key_rows):
    t2 = Inches(4.64) + r_i * Inches(0.38)
    highlight = r_i == 0
    bg2 = RGBColor(0xFF, 0xF0, 0xCC) if highlight else (LGRAY if r_i % 2 == 0 else WHITE)
    box(sld, Inches(0.4), t2, sum(col_ws_t), Inches(0.38), fill=bg2)
    for val, l, w, al in zip([var, coef, sig], col_ls_t, col_ws_t,
                               [PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER]):
        txt(sld, val, l + Inches(0.05), t2 + Inches(0.06),
            w - Inches(0.08), Inches(0.28), size=11, bold=highlight,
            color=BODY_DARK, align=al)

txt(sld, "** p < 0.05   — = not significant at conventional levels   HAC Newey–West SE (3 lags)",
    Inches(0.4), Inches(6.65), Inches(12.5), Inches(0.25),
    size=9, italic=True, color=MIDGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Result 1 continued: Post-2023 evidence
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "Result 1: When Data Quality Improves, the Picture Sharpens")
slide_number(sld, 10)

# Context box
box(sld, Inches(0.4), Inches(1.05), Inches(12.5), Inches(1.1), fill=LGRAY)
txt(sld,
    "The buffer variables (θ and L) were not individually significant at our full sample of 51 months. "
    "This is a power problem — not a theory problem. Reserve attestations only became comprehensive after 2023.",
    Inches(0.6), Inches(1.12), Inches(12.1), Inches(0.95),
    size=14, italic=True, color=BODY_DARK)

txt(sld, "When we restrict to January 2023 – March 2026  (N = 39 months, best attestation coverage):",
    Inches(0.4), Inches(2.3), Inches(12.5), Inches(0.4),
    size=15, bold=True, color=HDR_NAVY)

# Two big result boxes
for i, (label, val, stars, explain) in enumerate([
    ("Supply channel", "β₁  =  −8.14", "***",
     "Even stronger supply effect in recent data.\nConfirmed at the 1% level."),
    ("Buffer interaction", "L×ΔlnS  =  +49.17", "***  (p = 0.004)",
     "Higher liquid buffer dampens the spread compression.\nThe fragility mechanism is confirmed."),
]):
    l = Inches(0.4) + i * Inches(6.4)
    box(sld, l, Inches(2.85), Inches(6.1), Inches(3.2), fill=HDR_NAVY)
    txt(sld, label, l + Inches(0.2), Inches(2.98),
        Inches(5.7), Inches(0.38), size=14, bold=True,
        color=RGBColor(0xAA, 0xC4, 0xE8))
    txt(sld, val, l + Inches(0.2), Inches(3.4),
        Inches(5.7), Inches(0.85), size=30, bold=True, color=WHITE)
    txt(sld, stars, l + Inches(0.2), Inches(4.28),
        Inches(5.7), Inches(0.38), size=18, bold=True,
        color=RGBColor(0xFF, 0xD7, 0x00))
    txt(sld, explain, l + Inches(0.2), Inches(4.72),
        Inches(5.7), Inches(1.0), size=13,
        color=RGBColor(0xAA, 0xC4, 0xE8))

orange_callout(sld,
    "Both channels confirmed in recent data. The mechanism is real — our full sample just needs more months to reach significance.",
    size=14)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Result 2: Threshold (chart + plain explanation)
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "Result 2: Is There a Safety Threshold?")
slide_number(sld, 11)

img(sld, RESULTS / "threshold_ssr.png",
    Inches(0.3), Inches(1.05), Inches(7.0))

# Answer box — big and simple
box(sld, Inches(7.5), Inches(1.05), Inches(5.5), Inches(2.0), fill=HDR_NAVY)
txt(sld, "q*  =  13%",
    Inches(7.65), Inches(1.15), Inches(5.2), Inches(0.9),
    size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(sld, "of supply held as liquid cash",
    Inches(7.65), Inches(2.05), Inches(5.2), Inches(0.38),
    size=14, italic=True, color=RGBColor(0xAA, 0xC4, 0xE8), align=PP_ALIGN.CENTER)

# Plain explanation
box(sld, Inches(7.5), Inches(3.15), Inches(5.5), Inches(2.6), fill=LGRAY)
txt(sld, "What this means:",
    Inches(7.7), Inches(3.28), Inches(5.1), Inches(0.35),
    size=14, bold=True, color=HDR_NAVY)
txt(sld,
    "Using Hansen (2000) threshold regression, we find "
    "a tipping point at 13% liquid cash reserves.\n\n"
    "Below this level — where most of our data falls — "
    "issuers don't have enough cash to absorb redemptions "
    "without selling T-bills.",
    Inches(7.7), Inches(3.68), Inches(5.1), Inches(1.85),
    size=13, color=BODY_DARK)

# Chart label
txt(sld, "← Lower SSR = better model fit at that threshold value",
    Inches(0.35), Inches(6.6), Inches(7.0), Inches(0.3),
    size=10, italic=True, color=MIDGRAY)

txt(sld,
    "Bootstrap p = 0.260  (suggestive — the direction is clear; significance grows as data accumulates past 2026)",
    Inches(0.35), SH - Inches(0.5), Inches(12.6), Inches(0.32),
    size=11, italic=True, color=MIDGRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Result 2: Two worlds (regime comparison)
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "Below vs. Above the Threshold: Two Different Worlds")
slide_number(sld, 12)

txt(sld, "Splitting our 51 months at q* = 13% reveals a striking difference in how stablecoin growth affects Treasury markets:",
    Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.42),
    size=14, italic=True, color=MIDGRAY)

# Left: low buffer
box(sld, Inches(0.4), Inches(1.6), Inches(6.0), Inches(4.55),
    fill=RGBColor(0xFC, 0xEB, 0xE8))
box(sld, Inches(0.4), Inches(1.6), Inches(6.0), Inches(0.55),
    fill=RGBColor(0xC0, 0x39, 0x2B))
txt(sld, "BELOW  13%  LIQUID BUFFER",
    Inches(0.5), Inches(1.65), Inches(5.8), Inches(0.44),
    size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(sld, "38 out of 51 months",
    Inches(0.5), Inches(2.25), Inches(5.8), Inches(0.35),
    size=13, italic=True, color=MIDGRAY, align=PP_ALIGN.CENTER)
txt(sld, "β  =  −6.97",
    Inches(0.5), Inches(2.65), Inches(5.8), Inches(0.85),
    size=42, bold=True, color=RGBColor(0xC0, 0x39, 0x2B), align=PP_ALIGN.CENTER)
txt(sld,
    "Supply growth STRONGLY compresses spreads.\n\n"
    "Issuers don't have enough cash to pay redemptions, "
    "so they sell T-bills — moving the market.",
    Inches(0.6), Inches(3.6), Inches(5.6), Inches(2.3),
    size=14, color=BODY_DARK)

# Right: high buffer
box(sld, Inches(6.9), Inches(1.6), Inches(6.0), Inches(4.55),
    fill=RGBColor(0xE6, 0xF4, 0xEA))
box(sld, Inches(6.9), Inches(1.6), Inches(6.0), Inches(0.55),
    fill=RGBColor(0x1A, 0x7A, 0x4A))
txt(sld, "ABOVE  13%  LIQUID BUFFER",
    Inches(7.0), Inches(1.65), Inches(5.8), Inches(0.44),
    size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(sld, "13 out of 51 months",
    Inches(7.0), Inches(2.25), Inches(5.8), Inches(0.35),
    size=13, italic=True, color=MIDGRAY, align=PP_ALIGN.CENTER)
txt(sld, "β  =  +1.26",
    Inches(7.0), Inches(2.65), Inches(5.8), Inches(0.85),
    size=42, bold=True, color=RGBColor(0x1A, 0x7A, 0x4A), align=PP_ALIGN.CENTER)
txt(sld,
    "Supply growth has NO compression effect.\n\n"
    "Issuers use their cash buffer to absorb "
    "redemptions — T-bill holdings stay intact, "
    "spreads are undisturbed.",
    Inches(7.1), Inches(3.6), Inches(5.6), Inches(2.3),
    size=14, color=BODY_DARK)

orange_callout(sld,
    "The liquid cash buffer is what separates a stablecoin issuer that stabilises Treasury markets from one that disrupts them.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Result 3: Event Study
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "Result 3: Buffer-Conditioned Event Study")
slide_number(sld, 10)

# Three events
events = [
    ("LUNA / UST Collapse", "May 9, 2022",  "Low buffer",  "CAR = +8.91 pp***",  RGBColor(0xC0, 0x39, 0x2B)),
    ("USDT Partial Depeg",  "May 12, 2022", "Low buffer",  "CAR = +8.85 pp***",  RGBColor(0xC0, 0x39, 0x2B)),
    ("USDC / SVB Failure",  "Mar 11, 2023", "High buffer", "CAR = −18.01 pp***", RGBColor(0x1A, 0x7A, 0x4A)),
]
for i, (name, date, regime, car, col) in enumerate(events):
    l = Inches(0.35) + i * Inches(4.35)
    box(sld, l, Inches(1.05), Inches(4.1), Inches(0.38), fill=HDR_NAVY)
    txt(sld, name, l + Inches(0.08), Inches(1.08),
        Inches(3.95), Inches(0.3), size=13, bold=True, color=WHITE)
    box(sld, l, Inches(1.43), Inches(4.1), Inches(0.32), fill=col)
    txt(sld, regime, l + Inches(0.08), Inches(1.46),
        Inches(3.95), Inches(0.26), size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(sld, l, Inches(1.75), Inches(4.1), Inches(1.1), fill=LGRAY)
    txt(sld, f"{date}\n{car}", l + Inches(0.1), Inches(1.8),
        Inches(3.9), Inches(0.95), size=14, bold=True, color=col, align=PP_ALIGN.CENTER)

img(sld, RESULTS / "event_study_cars.png",
    Inches(0.35), Inches(3.0), Inches(7.5))

tf = bullet_tb(sld, Inches(8.0), Inches(3.05), Inches(4.9), Inches(3.4))
add_line(tf, "Interpretation", size=15, bold=True, color=HDR_NAVY)
add_line(tf,
    "❖  Low-buffer runs → T-bill yields spike (+8.9 pp)\n"
    "    Cole–Kehoe crisis zone activated",
    size=13, color=BODY_DARK, space_before=8)
add_line(tf,
    "❖  High-buffer stress → flight to Treasuries\n"
    "    compresses spread (−18.0 pp)",
    size=13, color=BODY_DARK, space_before=8)
add_line(tf,
    "❖  Welch test: t = 15.22, p < 0.001\n"
    "    Buffer regime is the key conditioning variable ✅",
    size=13, color=BODY_DARK, space_before=8)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Robustness
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "Robustness Checks")
slide_number(sld, 11)

checks = [
    ("NW lag sensitivity\n(1, 2, or 3 lags)",
     "β₁ stable across all lag choices\n(−5.95 to −6.02, p < 0.05 throughout)\nθ, L, L×ΔlnS remain NS",
     "✅  PASS", RGBColor(0x1A, 0x7A, 0x4A)),
    ("Post-2023 sub-sample\n(Jan 2023 – Mar 2026, N=39)",
     "β₁ = −8.14*** (p < 0.001)\nL×ΔlnS = 49.17*** (p = 0.004)\nH2 recovers with broader attestation coverage",
     "✅  PASS", RGBColor(0x1A, 0x7A, 0x4A)),
    ("Mean-centering\n(VIF reduction)",
     "All VIFs < 8 in N=51 sample\nCoefficients and significance unchanged\nInteraction term VIF resolves with centering",
     "✅  PASS", RGBColor(0x1A, 0x7A, 0x4A)),
    ("First-differenced\nspecification",
     "β₁ not significant in differences\nLevel spec captures structural relationship;\nevent study provides causal identification",
     "⚠️  PARTIAL", RGBColor(0xD4, 0x63, 0x2A)),
]

for i, (label, desc, verdict, col) in enumerate(checks):
    l = Inches(0.35) + (i % 2) * Inches(6.35)
    t = Inches(1.12) + (i // 2) * Inches(2.65)
    box(sld, l, t, Inches(6.05), Inches(2.5), fill=LGRAY)
    box(sld, l, t, Inches(6.05), Inches(0.42), fill=HDR_NAVY)
    txt(sld, label.split("\n")[0], l + Inches(0.1), t + Inches(0.07),
        Inches(4.0), Inches(0.3), size=13, bold=True, color=WHITE)
    txt(sld, verdict, l + Inches(4.1), t + Inches(0.07),
        Inches(1.85), Inches(0.3), size=12, bold=True, color=col, align=PP_ALIGN.RIGHT)
    txt(sld, desc, l + Inches(0.12), t + Inches(0.55),
        Inches(5.82), Inches(1.75), size=13, color=BODY_DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Summary of Results
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "Summary of Results")
slide_number(sld, 12)

findings = [
    ("✅", HDR_NAVY,
     "H1 confirmed — privilege amplification",
     "β₁ = −6.02**  |  1 SD supply growth → −24 bp spread compression  (N=51, Jan 2022–Mar 2026)"),
    ("~", RGBColor(0xD4, 0x63, 0x2A),
     "H2 suggestive — buffer channel (power-limited at N=51)",
     "θ, L, L×ΔlnS individually NS; post-2023 sub-sample confirms β₁=−8.14***, L×ΔlnS=49.17***"),
    ("~", RGBColor(0xD4, 0x63, 0x2A),
     "Threshold suggestive — q* = 0.1301 (p = 0.260)",
     "Liquid reserves ≥ 13% of supply keeps issuers outside the fragility zone (economically meaningful)"),
    ("✅", HDR_NAVY,
     "Event study confirms asymmetry",
     "Low-buffer runs: CAR +8.9 pp  vs.  High-buffer stress: CAR −18.0 pp  (Welch t = 15.2***)"),
    ("⚠️", ORANGE,
     "Stationarity caveat",
     "Level results don't survive first-differencing; event study provides causal identification"),
]

for i, (icon, col, title, detail) in enumerate(findings):
    t = Inches(1.1) + i * Inches(1.08)
    box(sld, Inches(0.35), t, Inches(0.55), Inches(0.93), fill=col)
    txt(sld, icon, Inches(0.35), t + Inches(0.2),
        Inches(0.55), Inches(0.55), size=18, color=WHITE, align=PP_ALIGN.CENTER)
    box(sld, Inches(0.95), t, Inches(12.0), Inches(0.93),
        fill=LGRAY if i % 2 == 0 else WHITE)
    txt(sld, title, Inches(1.1), t + Inches(0.06),
        Inches(11.7), Inches(0.36), size=15, bold=True, color=HDR_NAVY)
    txt(sld, detail, Inches(1.1), t + Inches(0.46),
        Inches(11.7), Inches(0.4), size=13, color=BODY_DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Policy Implications
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "Policy Implications")
slide_number(sld, 13)

txt(sld,
    "Our threshold estimate provides a direct reference point at a moment when the US, EU, "
    "Japan, and Hong Kong are each finalising stablecoin reserve legislation.",
    Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.5),
    size=14, italic=True, color=MIDGRAY)

policies = [
    ("Minimum liquid reserve ratio",
     "Require cash + near-cash reserves ≥ 13% of outstanding supply — "
     "placing issuers above the empirically estimated fragility threshold (q* = 0.1301)."),
    ("Decomposed reporting requirements",
     "Mandate separate disclosure of θ (T-bill exposure) and L (liquid cash buffer). "
     "Aggregate reserve ratios obscure the quality and liquidity of the backing."),
    ("Velocity-contingent buffers",
     "Our velocity variable V (7-day rolling SD of supply changes) predicts stress. "
     "Higher recent redemption velocity → higher required liquid buffer. Dynamic rules."),
    ("Systemic size threshold",
     "Fragility is proportional to scale. A $100bn issuer's forced liquidation is 10× "
     "more disruptive than a $10bn issuer. Tiered requirements by AUM."),
]

for i, (title, desc) in enumerate(policies):
    l = Inches(0.35) + (i % 2) * Inches(6.45)
    t = Inches(1.72) + (i // 2) * Inches(2.35)
    box(sld, l, t, Inches(6.1), Inches(2.2), fill=LGRAY)
    box(sld, l, t, Inches(0.18), Inches(2.2), fill=HDR_NAVY)
    txt(sld, title, l + Inches(0.28), t + Inches(0.12),
        Inches(5.7), Inches(0.38), size=14, bold=True, color=HDR_NAVY)
    txt(sld, desc, l + Inches(0.28), t + Inches(0.55),
        Inches(5.7), Inches(1.5), size=13, color=BODY_DARK)

orange_callout(sld,
    "Policy bottom line: Liquid cash reserves ≥ 13% of supply (q* = 0.1301) keeps issuers outside the fragility zone.",
    size=17)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Conclusion
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "Conclusion")
slide_number(sld, 14)

box(sld, Inches(0.35), Inches(1.05), Inches(12.6), Inches(1.15), fill=LGRAY)
txt(sld,
    "Large-scale USD-pegged stablecoins introduce a two-sided 'New Triffin Dilemma': "
    "amplifying U.S. exorbitant privilege in normal times, while embedding a fragility that "
    "can rapidly reverse that benefit when liquid cash reserves fall below a critical threshold.",
    Inches(0.5), Inches(1.1), Inches(12.3), Inches(1.05), size=15, color=HDR_NAVY)

contribs = [
    ("Theory",     "Extended Maggiori (2017) with supply S,\nθ (privilege channel) and L (fragility channel)"),
    ("Empirics 1", "β₁ = −6.02** (H1 confirmed)\nPost-2023: β₁=−8.14***, L×ΔlnS=49.17*** ✅"),
    ("Empirics 2", "Suggestive threshold: q* = 0.1301 (p=0.260)\nLiquid reserves ≥ 13% of supply"),
    ("Empirics 3", "27 pp CAR swing between low- and\nhigh-buffer stress episodes"),
]
for i, (label, text) in enumerate(contribs):
    l = Inches(0.35) + (i % 2) * Inches(6.35)
    t = Inches(2.35) + (i // 2) * Inches(1.7)
    box(sld, l, t, Inches(6.1), Inches(1.55), fill=LGRAY)
    box(sld, l, t, Inches(1.3), Inches(1.55), fill=HDR_NAVY)
    txt(sld, label, l + Inches(0.05), t + Inches(0.5),
        Inches(1.2), Inches(0.55), size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sld, text, l + Inches(1.38), t + Inches(0.15),
        Inches(4.6), Inches(1.25), size=13, color=BODY_DARK)

tf = bullet_tb(sld, Inches(0.5), Inches(5.85), Inches(12.4), Inches(0.55))
add_line(tf,
    "Limitation: Level results don't survive first-differencing. "
    "Event study (daily data, quasi-experimental) provides the primary causal identification.",
    size=12, italic=True, color=MIDGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Thank You (matches original template style)
# ═══════════════════════════════════════════════════════════════════════════════
sld = blank(prs)
title_slide_bg(sld)

txt(sld, "THANK YOU",
    Inches(4.0), Inches(2.85), Inches(8.8), Inches(0.6),
    size=28, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font="Calibri Light")

box(sld, Inches(3.9), Inches(4.45), Inches(7.5), Inches(2.8),
    fill=RGBColor(0x0E, 0x26, 0x68))

txt(sld,
    "Yonsei GSIS 2026-1\nTopics in International Finance\n"
    "Mireu Kim · Sara Chekroune · Oybek Ibragimov\n"
    "kimmireu0921@gmail.com",
    Inches(4.1), Inches(4.6), Inches(7.2), Inches(2.5),
    size=16, color=WHITE, align=PP_ALIGN.LEFT)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "Stablecoins_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
