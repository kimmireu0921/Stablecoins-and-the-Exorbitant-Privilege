"""
add_notes.py — adds speaker notes to the existing pptx pulled from GitHub.
Run once: python add_notes.py
"""

from pptx import Presentation
from pathlib import Path

PATH = Path("presentations/0526_Stablecoin_Exorbitant_Privilege.pptx")

SCRIPTS = {
    1: (
        "Good afternoon everyone. Our paper is called Stablecoins and the Exorbitant Privilege. "
        "The central argument is that large USD-pegged stablecoin issuers — mainly Tether and Circle — "
        "have become significant buyers of U.S. Treasury bills. "
        "In normal times, that deepens demand for safe assets and amplifies America's borrowing advantage. "
        "But when a run hits and their cash reserves are too thin, they are forced to liquidate T-bills, "
        "which rapidly reverses that benefit. "
        "Today we want to address some methodological questions from last session "
        "and walk clearly through our three stress events."
    ),
    2: (
        "Let me quickly orient everyone. On the left is everything we have completed. "
        "We extended Maggiori's 2017 model to incorporate stablecoin supply, treasury exposure, and the liquid buffer. "
        "Our main regression finds that a one standard deviation increase in stablecoin supply "
        "compresses the OIS-Treasury spread by roughly 6 basis points — that is the privilege amplification result. "
        "We also estimated a safety threshold using Hansen's 2000 threshold regression: "
        "issuers need to hold at least 13 percent of supply in liquid cash to stay outside the fragility zone. "
        "And we ran a buffer-conditioned event study across three stress episodes, "
        "which produced a 27 percentage point swing between the low and high buffer regimes. "
        "This deck specifically addresses the clarification questions from last time — "
        "why we use the OIS spread, what the buffer condition actually means, and how to read the CAR graphs. "
        "Baptiste will walk us through this slide, then hand over to Sara."
    ),
    3: (
        "A question came up about why we use the OIS-Treasury spread instead of just Treasury yields. "
        "Here is the intuition. A raw T-bill yield moves for two completely different reasons at the same time — "
        "the Fed raising or cutting rates, and investors specifically wanting to hold safe Treasuries. "
        "We only care about the second one. "
        "The OIS rate tracks market expectations of the Fed Funds path — it is pure monetary policy. "
        "So when we subtract it from the T-bill yield, we cancel out the monetary policy component "
        "and are left with just the premium investors pay to hold Treasuries specifically. "
        "That premium is exactly what stablecoin issuers affect when they buy T-bills. "
        "This approach is well established in the literature. "
        "Krishnamurthy and Vissing-Jorgensen in 2012 use this same spread to measure safe-asset demand. "
        "Nagel in 2016 calls it the convenience yield of near-money assets. "
        "Greenwood, Hanson and Stein in 2015 show it captures demand that is insensitive to monetary policy. "
        "The practical point is this: if we had used raw yields, "
        "a Fed rate cut and a Tether buying surge would look identical in our data. "
        "The spread cleanly separates them."
    ),
    4: (
        "Now let me explain precisely what the buffer condition is, because it is the central variable in our paper. "
        "L is the liquid buffer — the share of an issuer's reserves held in cash or cash equivalents, "
        "divided by total stablecoin supply. "
        "Our Hansen threshold regression estimates a critical cutoff at 13 percent. "
        "Here is why it matters. When investors redeem stablecoins for dollars, the issuer has to pay them back. "
        "If they hold enough cash — above 13 percent — they pay from cash and T-bill holdings are untouched. "
        "The Treasury market feels nothing. "
        "But if the cash buffer is below 13 percent, the issuer cannot cover redemptions from cash alone. "
        "They have to sell Treasury bills quickly to raise the dollars. "
        "That sudden selling pushes T-bill yields up relative to OIS — the spread spikes. "
        "In May 2022, Tether had only around 5 to 8 percent in liquid cash — below the threshold. "
        "This is the low-buffer regime. Theory predicts that forced T-bill selling in this regime "
        "should push the OIS-Treasury spread upward. "
        "The buffer condition is what separates an issuer that stabilises Treasury markets from one that disrupts them."
    ),
    5: (
        "CAR stands for Cumulative Abnormal Return — but in our context it measures the cumulative abnormal spread, "
        "not a stock return. Let me explain what the y-axis means and an important correction we are making. "
        "We estimate how the OIS-Treasury spread normally behaves using six months of data before each crisis. "
        "The abnormal spread on any given day is the actual spread minus that prediction. "
        "CAR is the running total of those daily abnormal values. "
        "I want to flag a methodological correction we identified after last week's presentation. "
        "Our original model estimated the LEVEL of the spread each day. "
        "The problem is that the estimation window for the 2022 events coincided with the Fed hiking cycle — "
        "the spread rose from 5 basis points in January to 78 basis points by May 2022, "
        "entirely due to the Federal Reserve raising rates rapidly. "
        "Our model's baseline was the SIX-MONTH AVERAGE of 37 basis points, "
        "but the events happened when the spread was already at 71 basis points. "
        "So the model attributed the entire Fed-driven increase as abnormal — inflating the CARs. "
        "We have corrected to a first-difference model, which examines day-over-day CHANGES in the spread. "
        "This removes trend contamination, consistent with how standard finance event studies work. "
        "With the corrected model, the 2022 CARs are small and not statistically significant. "
        "The direction during LUNA and USDT was slightly negative, meaning the spread actually compressed slightly — "
        "this is because the general risk-off buying of T-bills during a crisis dominated any Tether selling pressure. "
        "The economic story is still directionally correct but less clean than the original numbers suggested. "
        "The USDT partial depeg on May 12th, 2022 — three days after LUNA — "
        "was a direct bank-run on Tether itself: over 7 billion dollars redeemed in a single day. "
        "But with the corrected model, the CAR was not significantly different from zero."
    ),
    6: (
        "Our third event is the SVB failure in March 2023. "
        "SVB collapsed on March 10th — the second largest U.S. bank failure in history. "
        "Circle disclosed that 3.3 billion dollars of USDC reserves — roughly 8 percent of total reserves — "
        "were held at SVB and temporarily inaccessible. "
        "USDC depegged to around 87 cents. "
        "Two forces explain the market response shown on this slide. "
        "Force A is government intervention. "
        "U.S. regulators guaranteed all SVB deposits within 72 hours of the collapse. "
        "Because of that guarantee, Circle never needed to sell a single Treasury bill. "
        "The crisis was resolved before any forced liquidation became necessary. "
        "Force B is market-wide flight to safety. "
        "SVB's collapse triggered investors across the financial system to buy U.S. Treasury bills. "
        "That buying compressed the OIS-Treasury spread. "
        "Now, the visual on this slide shows the original event study results using our old model. "
        "With our corrected first-difference model, the CAR for SVB is near zero and not significant, "
        "because the spread compression during SVB is fully explained by the VIX spike and equity selloff — "
        "those general risk factors already predict flight-to-safety buying. "
        "There is no residual stablecoin-specific effect above and beyond what the banking crisis itself caused. "
        "This is consistent with the high-buffer theory: Circle never had to sell T-bills, "
        "so there was no stablecoin-specific disruption to the Treasury market. "
        "The flight-to-safety was a banking crisis effect, not a stablecoin effect."
    ),
    7: (
        "Let me summarize where the event study stands and where the paper's contribution actually lies. "
        "We completed the placebo test: three non-crisis dates with the same buffer regimes. "
        "With our corrected first-difference model, the actual events and the placebos all produce "
        "CARs close to zero and statistically indistinguishable. "
        "This is an honest result. The event study does not find a clean, statistically significant "
        "stablecoin-specific effect on the OIS-Treasury spread around these three events. "
        "The reason is that general market conditions — VIX spikes, equity selloffs — "
        "already explain most of the spread movement during crises. "
        "What the event study DOES confirm is the directional mechanism: "
        "during 2022 low-buffer events, the spread did not spike above what risk factors predict, "
        "consistent with general risk-off buying of T-bills offsetting any Tether selling. "
        "During the high-buffer SVB event, the spread compression was entirely driven by the banking crisis, "
        "not by stablecoin-specific T-bill selling. "
        "The PRIMARY empirical contribution of our paper is the main regression: "
        "a one standard deviation increase in stablecoin supply compresses the OIS-Treasury spread "
        "by approximately 6 basis points in normal times. "
        "That is a clean, statistically significant finding with economic meaning. "
        "The buffer threshold at 13 percent is the second key finding. "
        "The event study provides qualitative narrative support for the mechanism — not primary evidence."
    ),
    8: (
        "Thank you for your time. "
        "Our paper makes two main contributions. "
        "First: stablecoin issuers have become structurally important buyers of U.S. Treasury bills. "
        "Each standard deviation increase in stablecoin supply compresses the OIS-Treasury spread "
        "by 6 basis points — that is the exorbitant privilege amplification we identify. "
        "Second: this relationship is conditional on the liquid buffer. "
        "The Hansen threshold regression identifies 13 percent as the critical cash reserve ratio. "
        "Above it, issuers absorb redemption shocks without touching T-bills. "
        "Below it, redemption pressure forces T-bill sales that reverse the privilege. "
        "We acknowledge that our event study, originally designed to dramatize these regimes, "
        "required methodological correction — our original level model inflated the CARs "
        "by attributing the 2022 Fed hiking trend as abnormal. "
        "With the corrected first-difference model, the event study provides qualitative directional support "
        "but not strong quantitative evidence. "
        "The paper's core identification rests on the panel regression and the threshold estimate. "
        "We are happy to take any questions."
    ),
}


def set_notes(slide, text: str):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


prs = Presentation(str(PATH))

if len(prs.slides) != len(SCRIPTS):
    print(f"WARNING: {len(prs.slides)} slides in file but {len(SCRIPTS)} scripts written. Check counts.")

for i, slide in enumerate(prs.slides, 1):
    if i in SCRIPTS:
        set_notes(slide, SCRIPTS[i])
        print(f"  Slide {i}: notes added ({len(SCRIPTS[i])} chars)")
    else:
        print(f"  Slide {i}: no script defined — skipped")

prs.save(str(PATH))
print(f"\nSaved: {PATH}")
