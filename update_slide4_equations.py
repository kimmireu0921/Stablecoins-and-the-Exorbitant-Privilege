"""
update_slide4_equations.py
- Removes floating equation image (Picture 10)
- Narrows bullet TextBox 9 to left column
- Adds 3 stacked equation blocks on the right column
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

TEAL   = RGBColor(0x1B, 0x99, 0x8B)
GOLD   = RGBColor(0xE8, 0xA1, 0x2C)
NAVY   = RGBColor(0x13, 0x29, 0x4B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LTGRAY = RGBColor(0xF2, 0xF4, 0xF8)
DKGRAY = RGBColor(0x33, 0x42, 0x5C)
MGRAY  = RGBColor(0x5B, 0x6B, 0x82)
RED    = RGBColor(0xC1, 0x12, 0x1F)
SANS   = "Calibri"
ANS    = "http://schemas.openxmlformats.org/drawingml/2006/main"


def get_shape(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None


def clear_tf(tf):
    txBody = tf._txBody
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)
    txBody.append(parse_xml(f'<a:p xmlns:a="{ANS}"/>'))


def fp(tf, text, sz, bold=False, color=NAVY, align=PP_ALIGN.LEFT, sb=0, sa=0):
    p = tf.paragraphs[0]
    p.alignment = align
    if sb: p.space_before = Pt(sb)
    if sa: p.space_after  = Pt(sa)
    r = p.add_run()
    r.text = text; r.font.name = SANS; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.color.rgb = color


def ap(tf, text, sz, bold=False, color=NAVY, align=PP_ALIGN.LEFT, sb=0, sa=0):
    p = tf.add_paragraph()
    p.alignment = align
    if sb: p.space_before = Pt(sb)
    if sa: p.space_after  = Pt(sa)
    r = p.add_run()
    r.text = text; r.font.name = SANS; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.color.rgb = color


def add_rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def add_tb(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


# ══════════════════════════════════════════════════════════════════════════════
prs = Presentation("presentations/0616_Stablecoin_Exorbitant_Privilege.pptx")
s4  = prs.slides[3]

# ── 1. Remove floating equation image ─────────────────────────────────────────
pic = get_shape(s4, "Picture 10")
if pic:
    s4.shapes._spTree.remove(pic._element)
    print("Removed Picture 10")

# ── 2. Narrow bullet column (TextBox 9) and rewrite at 9pt ───────────────────
BULLET_W = 5500000
tb9 = get_shape(s4, "TextBox 9")
tb9.width = BULLET_W
tf = tb9.text_frame
tf.word_wrap = True
clear_tf(tf)

fp(tf, "β₁ — Privilege Amplification", 11, bold=True, color=NAVY, sa=3)
ap(tf, "→  Measures the baseline direct effect of stablecoin supply growth on yield spreads, independent of the buffer level", 9, color=DKGRAY, sb=2)
ap(tf, "→  Before: β₁ = −7.57 (p = 0.004) — statistically significant", 9, color=DKGRAY)
ap(tf, "→  After re-run: β₁ = +2.74 — not significant; sign flipped", 9, bold=True, color=GOLD)
ap(tf, "→  Verdict: amplification claim does not hold", 9, color=DKGRAY, sa=2)

ap(tf, "β₄ — Buffer Interaction", 11, bold=True, color=NAVY, sb=12, sa=3)
ap(tf, "→  Tests whether the reserve buffer level amplifies or dampens the effect of stablecoin supply growth on spreads", 9, color=DKGRAY, sb=2)
ap(tf, "→  Levels regression: β₄ = −35.89 (p = 0.032) — appeared significant", 9, color=DKGRAY)
ap(tf, "→  First-differenced: β₄ = −107 (p = 0.46) — collapses to insignificance", 9, bold=True, color=GOLD)
ap(tf, "→  Result vanishes under differencing — consistent with spurious regression", 9, color=DKGRAY)

# ── 3. Equation column layout ─────────────────────────────────────────────────
# Card:    L=548640, right edge=11640185
# Bullets: L=822960, right edge=822960+5500000=6322960
EQ_L     = 6422960          # equation column left (100k gap after bullets)
EQ_W     = 11640185 - EQ_L - 100000   # ≈ 5117225
ACCENT   = 55000            # thin left accent bar width
TXT_L    = EQ_L + ACCENT + 90000
TXT_W    = EQ_W - ACCENT - 90000 - 60000

# Block tops (card content area: T≈1912000 to T≈5236000)
B1_T = 1960000
B2_T = 3060000
B3_T = 4200000
BH   = 970000   # block height

# ── Block 1: Original equation ────────────────────────────────────────────────
add_rect(s4, EQ_L, B1_T, ACCENT, BH, NAVY)
tf1 = add_tb(s4, TXT_L, B1_T + 50000, TXT_W, BH - 60000)
fp(tf1, "① Original  —  before professor's feedback", 8.5, bold=True, color=NAVY, sa=5)
ap(tf1, "Spread = α + β₁·ΔlnS + β₂·θ + β₃·L + β₄·(L×ΔlnS) + controls", 8.5, color=DKGRAY, sa=5)
ap(tf1, "β₁ = −7.57   (p = 0.004) ***", 9, bold=True, color=RED)

# ── Block 2: Corrected equation ───────────────────────────────────────────────
add_rect(s4, EQ_L, B2_T, ACCENT, BH + 100000, TEAL)
tf2 = add_tb(s4, TXT_L, B2_T + 50000, TXT_W, BH + 40000)
fp(tf2, "② Corrected  —  θ dropped (θ+L≈1, collinear), L kept", 8.5, bold=True, color=TEAL, sa=5)
ap(tf2, "Spread = α + β₁·ΔlnS + β₃·L + β₄·(L×ΔlnS) + controls", 8.5, color=DKGRAY, sa=5)
ap(tf2, "β₁ = +2.74   (p = 0.228)   not significant", 9, bold=True, color=GOLD, sa=3)
ap(tf2, "β₄ = −35.89  (p = 0.032)   significant in levels", 9, bold=True, color=GOLD)

# ── Block 3: First-differenced equation ──────────────────────────────────────
add_rect(s4, EQ_L, B3_T, ACCENT, BH, GOLD)
tf3 = add_tb(s4, TXT_L, B3_T + 50000, TXT_W, BH - 60000)
fp(tf3, "③ First-differenced  —  L replaced by ΔL (removes drift)", 8.5, bold=True, color=GOLD, sa=5)
ap(tf3, "ΔSpread = α + β₁·Δ²lnS + β₃·ΔL + β₄·(ΔL×Δ²lnS) + controls", 8.5, color=DKGRAY, sa=5)
ap(tf3, "β₄ = −107   (p = 0.46)   collapses — spurious confirmed", 9, bold=True, color=RED)

prs.save("presentations/0616_Stablecoin_Exorbitant_Privilege.pptx")
print("Done — slide 4 equations updated.")
