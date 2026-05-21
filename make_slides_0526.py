"""
make_slides_0526.py — response to professor's comments (May 26 presentation).

Addresses:
  1. Why OIS-Treasury spread instead of raw yields (literature-backed rationale)
  2. What is the buffer condition — precise definition
  3. What is CAR — explain the y-axis in plain terms
  4. Each event explained (not just graphs): LUNA/UST, USDT depeg, USDC/SVB
  5. Why the SVB CAR is negative, not flat

Same Yonsei GSIS theme as make_slides.py.
Saves to: presentations/0526_Stablecoin_Exorbitant_Privilege.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

RESULTS       = Path("results")
PRESENTATIONS = Path("presentations")
PRESENTATIONS.mkdir(exist_ok=True)

# ── Colour palette (identical to make_slides.py) ─────────────────────────────
DARK_NAVY  = RGBColor(0x0A, 0x1E, 0x5C)
MED_NAVY   = RGBColor(0x14, 0x30, 0x70)
HDR_NAVY   = RGBColor(0x1B, 0x3A, 0x72)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE     = RGBColor(0xD4, 0x63, 0x2A)
BODY_DARK  = RGBColor(0x1A, 0x1A, 0x2E)
LGRAY      = RGBColor(0xF0, 0xF4, 0xF8)
MIDGRAY    = RGBColor(0x55, 0x55, 0x66)
GREEN      = RGBColor(0x1A, 0x7A, 0x4A)
RED        = RGBColor(0xC0, 0x39, 0x2B)
LIGHT_RED  = RGBColor(0xFC, 0xEB, 0xE8)
LIGHT_GRN  = RGBColor(0xE6, 0xF4, 0xEA)
GOLD       = RGBColor(0xFF, 0xD7, 0x00)

SW = Inches(13.33)
SH = Inches(7.5)


# ── Helpers (identical to make_slides.py) ────────────────────────────────────

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def box(sld, l, t, w, h, fill=None, line_color=None, line_width=Pt(0)):
    sh = sld.shapes.add_shape(1, l, t, w, h)
    sh.line.fill.background()
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = line_width
    return sh


def txt(sld, text, l, t, w, h, size=18, bold=False, italic=False,
        color=BODY_DARK, align=PP_ALIGN.LEFT, font="Calibri", wrap=True):
    tb = sld.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_line(tf, text, size=18, bold=False, italic=False,
             color=BODY_DARK, align=PP_ALIGN.LEFT, space_before=4):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return p


def img(sld, path, l, t, w, h=None):
    if Path(path).exists():
        if h:
            sld.shapes.add_picture(str(path), l, t, w, h)
        else:
            sld.shapes.add_picture(str(path), l, t, w)


def title_slide_bg(sld):
    box(sld, 0, 0, SW, SH, fill=DARK_NAVY)
    box(sld, 0, Inches(2.85), SW, Inches(1.15), fill=MED_NAVY)
    circle = sld.shapes.add_shape(9, Inches(0.1), Inches(0.0), Inches(3.8), Inches(3.8))
    circle.fill.solid(); circle.fill.fore_color.rgb = WHITE; circle.line.fill.background()
    inner = sld.shapes.add_shape(9, Inches(0.55), Inches(0.45), Inches(2.9), Inches(2.9))
    inner.fill.solid(); inner.fill.fore_color.rgb = DARK_NAVY; inner.line.fill.background()
    chev = sld.shapes.add_shape(9, Inches(0.1), Inches(4.4), Inches(3.6), Inches(3.6))
    chev.fill.solid(); chev.fill.fore_color.rgb = WHITE; chev.line.fill.background()
    inner_chev = sld.shapes.add_shape(9, Inches(0.65), Inches(4.85), Inches(2.5), Inches(2.5))
    inner_chev.fill.solid(); inner_chev.fill.fore_color.rgb = DARK_NAVY; inner_chev.line.fill.background()
    box(sld, 0, Inches(2.55), Inches(3.8), Inches(0.6), fill=MIDGRAY)


def content_slide(prs, title_text):
    sld = blank(prs)
    box(sld, 0, 0, SW, SH, fill=WHITE)
    box(sld, 0, 0, SW, Inches(0.95), fill=HDR_NAVY)
    txt(sld, title_text,
        Inches(0.3), Inches(0.1), Inches(11.5), Inches(0.78),
        size=26, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font="Calibri")
    return sld


def slide_number(sld, num):
    txt(sld, str(num),
        SW - Inches(0.55), SH - Inches(0.42), Inches(0.45), Inches(0.35),
        size=13, color=MIDGRAY, align=PP_ALIGN.RIGHT)


def orange_callout(sld, text, t=None, size=17):
    if t is None:
        t = SH - Inches(1.38)
    box(sld, 0, t, SW, Inches(1.2), fill=ORANGE)
    txt(sld, text,
        Inches(0.6), t + Inches(0.18), SW - Inches(1.0), Inches(0.88),
        size=size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def bullet_tb(sld, l, t, w, h):
    tb = sld.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].space_before = Pt(0)
    return tf


def notes(sld, script: str):
    """Add presenter notes to a slide."""
    notes_slide = sld.notes_slide
    notes_slide.notes_text_frame.text = script


# ── Build presentation ────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
sld = blank(prs)
title_slide_bg(sld)

txt(sld, "STABLECOINS AND THE EXORBITANT PRIVILEGE",
    Inches(4.0), Inches(2.78), Inches(9.0), Inches(0.65),
    size=22, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font="Calibri Light")

txt(sld, "Clarifying the Methodology: Spread Measure, Buffer Condition, and Event Study",
    Inches(4.0), Inches(3.42), Inches(9.0), Inches(0.5),
    size=13, italic=True, color=RGBColor(0xAA, 0xC4, 0xE8), align=PP_ALIGN.LEFT)

box(sld, Inches(3.9), Inches(4.45), Inches(7.5), Inches(2.8),
    fill=RGBColor(0x0E, 0x26, 0x68))

txt(sld,
    "Yonsei GSIS 2026-1  |  Topics in International Finance\n"
    "Mireu Kim · Sara Chekroune · Oybek Ibragimov\n"
    "May 26, 2026",
    Inches(4.1), Inches(4.6), Inches(7.2), Inches(2.5),
    size=16, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

notes(sld,
    "Good afternoon, everyone. Our paper is called Stablecoins and the Exorbitant Privilege. "
    "The central argument is that large stablecoin issuers like Tether and Circle have become "
    "significant buyers of U.S. Treasury bills — which amplifies America's safe-asset privilege "
    "in normal times, but can sharply reverse that benefit during a run if their cash reserves "
    "are too thin. Today we want to address some of the methodological questions that came up "
    "last time and walk you through our three key findings clearly.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Roadmap
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "Roadmap")
slide_number(sld, 2)

# Left column — What we have done
box(sld, Inches(0.35), Inches(1.05), Inches(6.1), Inches(0.42), fill=HDR_NAVY)
txt(sld, "What we have done",
    Inches(0.5), Inches(1.08), Inches(5.9), Inches(0.34),
    size=15, bold=True, color=WHITE)

done_items = [
    ("Theory",       "Extended Maggiori (2017) with stablecoin supply S, treasury exposure θ, and liquid buffer L"),
    ("Regression",   "β₁ = −6.02** — stablecoin growth compresses OIS–Treasury spreads (privilege amplification)"),
    ("Threshold",    "Hansen (2000) threshold at q* = 13% liquid cash — separates safe from fragile regime"),
    ("Event Study",  "3 stress episodes (LUNA, USDT depeg, SVB) — 27 pp CAR swing between low/high buffer regimes"),
    ("Robustness",   "Mean-centering (VIF), Engle-Granger cointegration, first-differenced spec, post-2023 subsample"),
    ("This deck",    "Clarified: why OIS spread, what the buffer condition means, how to read CAR, what each event was"),
]

for i, (tag, desc) in enumerate(done_items):
    t = Inches(1.55) + i * Inches(0.83)
    box(sld, Inches(0.35), t, Inches(1.2), Inches(0.76), fill=HDR_NAVY)
    txt(sld, tag, Inches(0.38), t + Inches(0.2),
        Inches(1.14), Inches(0.38), size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(sld, Inches(1.55), t, Inches(4.9), Inches(0.76), fill=LGRAY if i % 2 == 0 else WHITE)
    txt(sld, f"✓  {desc}", Inches(1.68), t + Inches(0.18),
        Inches(4.7), Inches(0.5), size=12, color=BODY_DARK)

# Right column — What we plan to do next week
box(sld, Inches(6.8), Inches(1.05), Inches(6.15), Inches(0.42), fill=ORANGE)
txt(sld, "Planned for next week",
    Inches(6.95), Inches(1.08), Inches(5.95), Inches(0.34),
    size=15, bold=True, color=WHITE)

box(sld, Inches(6.8), Inches(1.55), Inches(6.15), Inches(5.7), fill=RGBColor(0xFF, 0xF5, 0xEE))

tf = bullet_tb(sld, Inches(7.0), Inches(1.65), Inches(5.8), Inches(5.5))
add_line(tf, "Placebo Test", size=16, bold=True, color=ORANGE, space_before=4)
add_line(tf,
    "Run the same event study methodology on 3 ordinary days "
    "with no known stablecoin stress event:",
    size=13, color=BODY_DARK, space_before=8)

placebo_rows = [
    ("Jun 15, 2021", "Low buffer",  "Quiet summer — no redemption pressure"),
    ("Oct 12, 2021", "Low buffer",  "Pre-LUNA, no crypto-specific stress"),
    ("Jul 15, 2025", "High buffer", "Post-SVB recovery, stable spread period"),
]
for date, regime, note in placebo_rows:
    add_line(tf, f"❖  {date}  ({regime})  —  {note}", size=12, color=BODY_DARK, space_before=7)

add_line(tf, " ", size=6, space_before=2)
add_line(tf, "What we expect to show:", size=14, bold=True, color=HDR_NAVY, space_before=4)
add_line(tf,
    "CARs on placebo days should be near 0 pp — confirming that "
    "the 9–18 pp swings in actual events are caused by the "
    "crisis itself, not by random market drift.",
    size=13, color=BODY_DARK, space_before=6)
add_line(tf, " ", size=6, space_before=2)
add_line(tf,
    "Preliminary result: placebo mean |CAR| = 0.55 pp vs. "
    "actual mean |CAR| = 11.92 pp  →  21.7× difference",
    size=12, italic=True, color=ORANGE, bold=True, space_before=4)

notes(sld,
    "Let me quickly walk you through what we have done and what is coming. "
    "On the left, you can see everything we have completed: the theoretical framework, "
    "the main regression, the threshold test, the event study across three crises, "
    "and several robustness checks. This deck specifically addresses the clarification "
    "questions from last session — why we use the OIS spread, what the buffer condition "
    "is, and how to read the event study graphs. "
    "On the right is what we are planning for next week. We are running a placebo test — "
    "we take three ordinary, uneventful days with the same buffer conditions as our real "
    "events, run the exact same methodology, and show that nothing unusual happens. "
    "Our preliminary result already shows the placebo CARs are less than one basis point, "
    "compared to nine to eighteen basis points in the actual events. That is a 22 times "
    "difference, which strongly supports that our results are real and not just noise.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Why OIS-Treasury Spread, Not Raw Yields?
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "Why OIS–Treasury Spread, Not Raw Yields?")
slide_number(sld, 3)

# Top context strip
box(sld, Inches(0.35), Inches(1.05), Inches(12.6), Inches(0.5), fill=LGRAY)
txt(sld,
    "A T-bill yield moves for two reasons simultaneously. We need to separate them.",
    Inches(0.55), Inches(1.12), Inches(12.1), Inches(0.38),
    size=14, italic=True, color=HDR_NAVY)

# Three-column breakdown
cols = [
    ("Raw T-bill Yield\n(what we did NOT use)",
     ["Driven by Fed rate expectations",
      "Driven by inflation expectations",
      "Driven by safe-asset demand",
      "Cannot distinguish these three forces"],
     RED, LIGHT_RED),
    ("OIS Rate\n(the subtracted benchmark)",
     ["Reflects expected Fed Funds path",
      "Pure monetary policy expectations",
      "No T-bill specific demand embedded",
      "Moves with macro, not stablecoin flows"],
     HDR_NAVY, LGRAY),
    ("OIS–Treasury Spread\n(our dependent variable)",
     ["= T-bill yield  −  OIS rate",
      "Cancels out monetary policy",
      "Isolates convenience / safety premium",
      "Only moves when T-bill demand shifts"],
     GREEN, LIGHT_GRN),
]

for i, (title, points, hcol, bcol) in enumerate(cols):
    l = Inches(0.35) + i * Inches(4.3)
    box(sld, l, Inches(1.65), Inches(4.1), Inches(0.5), fill=hcol)
    title_line1 = title.split("\n")[0]
    title_line2 = title.split("\n")[1] if "\n" in title else ""
    txt(sld, title_line1, l + Inches(0.08), Inches(1.68),
        Inches(3.95), Inches(0.26), size=13, bold=True, color=WHITE)
    if title_line2:
        txt(sld, title_line2, l + Inches(0.08), Inches(1.93),
            Inches(3.95), Inches(0.22), size=11, italic=True, color=WHITE)
    box(sld, l, Inches(2.15), Inches(4.1), Inches(2.3), fill=bcol)
    tf = bullet_tb(sld, l + Inches(0.12), Inches(2.22), Inches(3.88), Inches(2.15))
    for j, pt in enumerate(points):
        add_line(tf, f"{'❖' if i == 2 else '•'}  {pt}", size=13,
                 color=BODY_DARK, space_before=7 if j > 0 else 2)

# Literature row
box(sld, Inches(0.35), Inches(4.55), Inches(12.6), Inches(0.38), fill=HDR_NAVY)
txt(sld, "Literature precedent for using this spread as the convenience yield:",
    Inches(0.55), Inches(4.58), Inches(12.1), Inches(0.3),
    size=13, bold=True, color=WHITE)

lit = [
    ("Krishnamurthy &\nVissing-Jorgensen (2012)",
     "JF — use T-bill spread to measure how much investors pay for Treasuries' safety and liquidity"),
    ("Nagel (2016)",
     "RFS — defines OIS–T-bill spread as the 'convenience yield' of near-money assets"),
    ("Greenwood, Hanson\n& Stein (2015)",
     "QJE — show T-bill spread captures special demand that is insensitive to monetary policy"),
]
for i, (auth, desc) in enumerate(lit):
    l = Inches(0.35) + i * Inches(4.3)
    box(sld, l, Inches(4.93), Inches(4.1), Inches(1.55), fill=LGRAY)
    txt(sld, auth, l + Inches(0.12), Inches(5.0),
        Inches(3.88), Inches(0.45), size=12, bold=True, color=HDR_NAVY)
    txt(sld, desc, l + Inches(0.12), Inches(5.48),
        Inches(3.88), Inches(0.88), size=11, color=BODY_DARK)

orange_callout(sld,
    "When stablecoin issuers buy T-bills, T-bill yields fall but OIS does not move. "
    "Only the spread captures this demand shock. Raw yields would conflate it with Fed policy.")

notes(sld,
    "A question came up about why we use the OIS-Treasury spread instead of just Treasury yields. "
    "Here is the intuition. A raw T-bill yield moves for two completely different reasons at the same time: "
    "one is the Fed raising or cutting interest rates, and the other is investors specifically wanting "
    "to hold safe U.S. Treasury bills. We only care about the second one — the safe-asset demand. "
    "The OIS rate, which stands for Overnight Index Swap, tracks market expectations of the Fed Funds rate. "
    "So when we subtract it from the T-bill yield, we cancel out the monetary policy part and are left "
    "with just the premium that investors pay to hold Treasuries specifically. "
    "That premium is exactly what stablecoin issuers affect when they buy T-bills. "
    "This approach is standard in the literature — Krishnamurthy and Vissing-Jorgensen in 2012, "
    "and Nagel in 2016, both use this same spread as their measure of safe-asset convenience yield. "
    "The practical implication: if we had used raw yields, a Fed rate cut and a Tether buying surge "
    "would look identical in our data. The spread separates them cleanly.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — What is the Buffer Condition? (renumbered SLIDE 4)
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "What is the Buffer Condition?")
slide_number(sld, 4)

# Definition box
box(sld, Inches(0.35), Inches(1.05), Inches(12.6), Inches(1.5), fill=HDR_NAVY)
txt(sld, "Definition:",
    Inches(0.55), Inches(1.1), Inches(12.1), Inches(0.35),
    size=14, bold=True, color=RGBColor(0xAA, 0xC4, 0xE8))
txt(sld,
    "L  =  Liquid Buffer  =  Cash & Cash-Equivalent Reserves  ÷  Total Stablecoin Supply",
    Inches(0.55), Inches(1.45), Inches(12.1), Inches(0.42),
    size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Courier New")
txt(sld,
    "Estimated threshold from Hansen (2000) threshold regression:   L*  =  13%  of supply",
    Inches(0.55), Inches(1.9), Inches(12.1), Inches(0.38),
    size=13, italic=True, color=RGBColor(0xAA, 0xC4, 0xE8), align=PP_ALIGN.CENTER)

# Two regime boxes
for i, (regime, lval, what_happens, why_it_matters, col, bcol) in enumerate([
    ("LOW BUFFER  (L < 13%)",
     "e.g. Tether in May 2022: ~5–8% cash",
     [
         "Issuer holds mostly T-bills, very little cash",
         "When investors redeem stablecoins for dollars...",
         "Issuer MUST sell T-bills to raise the cash",
         "This selling pushes T-bill yields UP",
         "OIS-Treasury spread SPIKES (CAR is positive)",
     ],
     "The issuer becomes a forced seller into an already-stressed market.",
     RED, LIGHT_RED),
    ("HIGH BUFFER  (L ≥ 13%)",
     "e.g. Circle in March 2023: ~18–22% cash",
     [
         "Issuer holds substantial cash reserves",
         "When investors redeem stablecoins for dollars...",
         "Issuer pays out from cash — T-bills untouched",
         "No forced selling, no T-bill yield pressure",
         "Spread does NOT spike from issuer behavior",
     ],
     "The cash buffer absorbs the redemption shock. T-bill market is insulated.",
     GREEN, LIGHT_GRN),
]):
    l = Inches(0.35) + i * Inches(6.45)
    box(sld, l, Inches(2.65), Inches(6.1), Inches(0.45), fill=col)
    txt(sld, regime, l + Inches(0.1), Inches(2.68),
        Inches(5.9), Inches(0.38), size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sld, lval, l + Inches(0.1), Inches(3.13),
        Inches(5.9), Inches(0.28), size=11, italic=True, color=MIDGRAY, align=PP_ALIGN.CENTER)
    box(sld, l, Inches(3.41), Inches(6.1), Inches(2.55), fill=bcol)
    tf = bullet_tb(sld, l + Inches(0.15), Inches(3.48), Inches(5.82), Inches(2.4))
    for j, pt in enumerate(what_happens):
        add_line(tf, f"{'→' if j > 0 else '❖'}  {pt}", size=13,
                 color=BODY_DARK, space_before=6 if j > 0 else 2)
    box(sld, l, Inches(5.96), Inches(6.1), Inches(0.42), fill=col)
    txt(sld, why_it_matters, l + Inches(0.12), Inches(5.99),
        Inches(5.88), Inches(0.35), size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

orange_callout(sld,
    "The buffer condition is what separates a stablecoin issuer that stabilises Treasury markets from one that disrupts them.")

notes(sld,
    "The buffer condition is the central variable in our paper, so let me explain it precisely. "
    "L is the liquid buffer — it is the share of an issuer's reserves held in cash or cash equivalents, "
    "divided by the total stablecoin supply outstanding. "
    "Our Hansen threshold regression estimates a critical cutoff of 13 percent. "
    "Here is why it matters. When investors want to redeem their stablecoins for dollars, "
    "the issuer has to pay them back. If the issuer holds enough cash — above 13 percent — "
    "they simply pay from that cash and their Treasury bill holdings are untouched. "
    "The T-bill market does not feel anything. "
    "But if the cash buffer is below 13 percent, the issuer cannot cover redemptions from cash alone. "
    "They have to sell Treasury bills quickly to raise the dollars. "
    "That sudden selling pushes T-bill yields up relative to OIS — the spread spikes. "
    "So the buffer condition is essentially asking: does this issuer have enough cash to absorb "
    "a redemption shock without being forced to liquidate their T-bill portfolio? "
    "In May 2022, the answer was no. In March 2023, the answer was yes.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — What is CAR? (Explaining the Y-axis)
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "What is CAR? Reading the Y-Axis")
slide_number(sld, 5)

# Step-by-step construction
box(sld, Inches(0.35), Inches(1.05), Inches(12.6), Inches(0.38), fill=HDR_NAVY)
txt(sld, "CAR = Cumulative Abnormal Return on the OIS–Treasury Spread",
    Inches(0.55), Inches(1.08), Inches(12.1), Inches(0.3),
    size=15, bold=True, color=WHITE)

steps = [
    (HDR_NAVY, "Step 1 — Build a Normal Model",
     "Using data from 120 to 6 trading days BEFORE the event, we run a regression:\n"
     "Spread = f(VIX, ΔlnN*)   →   this captures how the spread normally responds to market conditions."),
    (MIDGRAY,  "Step 2 — Compute the Abnormal Spread",
     "For each day in the event window (5 days before to 20 days after the event):\n"
     "Abnormal Spread = Actual Spread  −  Predicted Spread\n"
     "A positive value means the spread was HIGHER than normal market conditions would predict."),
    (GREEN,    "Step 3 — Cumulate into CAR",
     "CAR on day τ = sum of all abnormal spreads from day −5 to day τ\n"
     "Y-axis unit: percentage points (pp) of spread deviation accumulated since 5 days before the event."),
]

for i, (col, title, body) in enumerate(steps):
    t = Inches(1.55) + i * Inches(1.3)
    box(sld, Inches(0.35), t, Inches(1.1), Inches(1.2), fill=col)
    txt(sld, str(i + 1), Inches(0.35), t + Inches(0.3),
        Inches(1.1), Inches(0.6), size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(sld, Inches(1.5), t, Inches(11.45), Inches(1.2), fill=LGRAY if i % 2 == 0 else WHITE)
    txt(sld, title, Inches(1.65), t + Inches(0.08),
        Inches(11.1), Inches(0.3), size=14, bold=True, color=HDR_NAVY)
    txt(sld, body, Inches(1.65), t + Inches(0.42),
        Inches(11.1), Inches(0.68), size=12, color=BODY_DARK)

# Reading guide
box(sld, Inches(0.35), Inches(5.55), Inches(12.6), Inches(0.38), fill=ORANGE)
txt(sld, "How to read the graph:",
    Inches(0.55), Inches(5.58), Inches(12.1), Inches(0.3),
    size=13, bold=True, color=WHITE)

guides = [
    ("CAR = +8.9 pp (LUNA event)",
     "By day +20, the spread had accumulated 8.9 pp MORE than expected — sustained T-bill selling pressure"),
    ("CAR = −18.0 pp (SVB event)",
     "By day +20, the spread had accumulated 18 pp LESS than expected — flight-to-safety buying compressed spreads"),
    ("CAR = 0 (flat line)",
     "Would mean the event had no effect — spread behaved exactly as the normal model predicted"),
]
for i, (val, explain) in enumerate(guides):
    l = Inches(0.35) + i * Inches(4.3)
    box(sld, l, Inches(5.93), Inches(4.1), Inches(1.35), fill=LGRAY)
    txt(sld, val, l + Inches(0.1), Inches(6.0),
        Inches(3.9), Inches(0.3), size=12, bold=True, color=HDR_NAVY)
    txt(sld, explain, l + Inches(0.1), Inches(6.33),
        Inches(3.9), Inches(0.88), size=11, color=BODY_DARK)

notes(sld,
    "CAR stands for Cumulative Abnormal Return — but in our context it is the cumulative "
    "abnormal spread, not a stock return. Let me explain what that means in plain terms. "
    "Before each crisis event, we estimate how the OIS-Treasury spread normally behaves "
    "using data from about six months prior. We build a simple model that says: given today's "
    "VIX level and global equity returns, what would we expect the spread to be? "
    "The abnormal spread on any given day is just the actual spread minus that prediction. "
    "If nothing unusual is happening, this should be close to zero. "
    "CAR is simply the running total of those daily abnormal spreads from five days before "
    "the event through twenty days after. "
    "So the y-axis is in percentage points of spread deviation accumulated over time. "
    "A CAR of plus 8.9 means that by day twenty, the spread had built up 8.9 percentage points "
    "more than we would normally expect — that is a large and persistent deviation. "
    "A CAR of minus 18 means the spread was 18 percentage points below normal — "
    "which happened because investors were rushing into T-bills for safety, compressing yields. "
    "A flat line at zero would mean the event had no measurable effect on Treasury markets.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Event 1: LUNA/UST Collapse
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "Event 1: LUNA / UST Collapse — May 9, 2022")
slide_number(sld, 6)

# Event banner
box(sld, Inches(0.35), Inches(1.05), Inches(12.6), Inches(0.42), fill=RED)
txt(sld, "Buffer Condition: LOW   |   L ≈ 5–8% cash   |   CAR = +8.91 pp***   |   t-stat > 30",
    Inches(0.55), Inches(1.08), Inches(12.1), Inches(0.34),
    size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# What happened (left)
box(sld, Inches(0.35), Inches(1.55), Inches(7.1), Inches(4.1), fill=LGRAY)
txt(sld, "What happened", Inches(0.5), Inches(1.62),
    Inches(6.9), Inches(0.35), size=15, bold=True, color=HDR_NAVY)

tf = bullet_tb(sld, Inches(0.5), Inches(2.05), Inches(6.9), Inches(3.5))
events_text = [
    "Terra's UST was an algorithmic stablecoin — its $1 peg was maintained by minting/burning LUNA, not by holding real dollars",
    "On May 7–9, a large coordinated sell-off drained UST's liquidity pool on Anchor Protocol",
    "UST lost its peg → panic triggered mass redemptions → LUNA was minted in huge quantities → LUNA price collapsed",
    "Within 72 hours: ~$40 billion in combined USDT+LUNA market cap was wiped out",
    "Contagion: even though USDT and USDC were NOT algorithmic, markets feared all stablecoins might fail",
    "Massive redemption pressure hit Tether (USDT) — and Tether had very little cash to pay back",
]
for j, line in enumerate(events_text):
    add_line(tf, f"❖  {line}", size=12, color=BODY_DARK, space_before=7 if j > 0 else 2)

# Why this affected the Treasury spread (right)
box(sld, Inches(7.6), Inches(1.55), Inches(5.1), Inches(4.1), fill=HDR_NAVY)
txt(sld, "Why the spread spiked", Inches(7.75), Inches(1.62),
    Inches(4.9), Inches(0.35), size=15, bold=True, color=WHITE)

tf2 = bullet_tb(sld, Inches(7.75), Inches(2.05), Inches(4.9), Inches(3.5))
mechanism = [
    "Tether's buffer was LOW (~5–8% cash)",
    "Redemption demand: investors wanted dollars, not T-bills",
    "Tether had to SELL T-bills rapidly to raise cash",
    "Sudden T-bill supply → T-bill yields rose above OIS",
    "Spread spiked — exactly our theory predicts",
    "CAR accumulated to +8.91 pp over 20 trading days",
]
for j, line in enumerate(mechanism):
    col = GOLD if j >= 3 else RGBColor(0xAA, 0xC4, 0xE8)
    add_line(tf2, f"{'→' if j >= 3 else '❖'}  {line}", size=12, color=col,
             space_before=7 if j > 0 else 2)

orange_callout(sld,
    "This event confirms the forced-selling channel: Low buffer + redemption pressure = T-bill liquidation = spread spike.")

notes(sld,
    "Our first event is the LUNA and UST collapse in May 2022. "
    "To understand why this affected Tether and Circle, you first need to understand what Terra was. "
    "UST was an algorithmic stablecoin — it did not hold real dollars or Treasury bills. "
    "Instead, its one-dollar peg was maintained by a mathematical relationship with its sister token, LUNA. "
    "When a large coordinated sell-off drained UST's liquidity on May 7th, "
    "the peg broke, a death spiral began, and roughly 40 billion dollars in market value was destroyed within days. "
    "Now, Tether and Circle had nothing to do with Terra. Their reserves are real assets. "
    "But markets panicked — investors began asking: what if all stablecoins are fragile? "
    "This caused a wave of redemptions across the industry, hitting Tether especially hard. "
    "Here is where our buffer condition becomes decisive. "
    "At the time, Tether held only about 5 to 8 percent of its reserves in liquid cash. "
    "The rest was in Treasury bills. So when redemptions came in, Tether had to sell T-bills to pay out. "
    "That forced selling pushed Treasury yields above OIS, and the spread spiked. "
    "Our event study shows a CAR of plus 8.9 percentage points over the following 20 days — "
    "a large, sustained, and statistically highly significant increase in the spread.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Event 2: USDT Partial Depeg
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "Event 2: USDT Partial Depeg — May 12, 2022")
slide_number(sld, 7)

box(sld, Inches(0.35), Inches(1.05), Inches(12.6), Inches(0.42), fill=RED)
txt(sld, "Buffer Condition: LOW   |   L ≈ 5–8% cash   |   CAR = +8.85 pp***   |   t-stat > 30",
    Inches(0.55), Inches(1.08), Inches(12.1), Inches(0.34),
    size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# What happened
box(sld, Inches(0.35), Inches(1.55), Inches(7.1), Inches(4.1), fill=LGRAY)
txt(sld, "What happened", Inches(0.5), Inches(1.62),
    Inches(6.9), Inches(0.35), size=15, bold=True, color=HDR_NAVY)

tf = bullet_tb(sld, Inches(0.5), Inches(2.05), Inches(6.9), Inches(3.5))
events2 = [
    "Three days after the LUNA collapse, markets turned their attention to USDT itself",
    "On May 12, Tether (USDT) briefly fell to $0.95 — a 5% depeg — on major exchanges",
    "This was the direct test: 'Does Tether actually have the reserves it claims?'",
    "Over 7 billion USDT was redeemed in a single day — the largest single-day redemption in Tether history",
    "Unlike LUNA/UST, this was NOT algorithmic — it was a bank-run-style confidence crisis on Tether",
    "Tether successfully maintained the peg and returned to $1.00 within days",
]
for j, line in enumerate(events2):
    add_line(tf, f"❖  {line}", size=12, color=BODY_DARK, space_before=7 if j > 0 else 2)

# Why nearly identical CAR to Event 1
box(sld, Inches(7.6), Inches(1.55), Inches(5.1), Inches(4.1), fill=HDR_NAVY)
txt(sld, "Why CAR ≈ +8.85 pp (same as Event 1)", Inches(7.75), Inches(1.62),
    Inches(4.9), Inches(0.35), size=13, bold=True, color=WHITE)

tf2 = bullet_tb(sld, Inches(7.75), Inches(2.05), Inches(4.9), Inches(3.5))
mechanism2 = [
    "Same mechanism, different trigger: direct run on USDT",
    "Same low buffer (~5–8% cash) — same constraint",
    "Same forced T-bill selling to meet $7bn in redemptions",
    "Near-identical CAR (+8.85 vs +8.91) is not coincidence:",
    "→ It reflects the SAME mechanical constraint",
    "→ The buffer level determined the market impact, not the specific trigger",
]
for j, line in enumerate(mechanism2):
    col = GOLD if j >= 3 else RGBColor(0xAA, 0xC4, 0xE8)
    add_line(tf2, f"{'❖' if j < 3 else ''}  {line}", size=12, color=col,
             space_before=7 if j > 0 else 2)

orange_callout(sld,
    "Two separate events, same buffer level, same CAR magnitude. The buffer is the key variable — not the nature of the shock.")

notes(sld,
    "Three days after the LUNA collapse, markets turned their attention directly to Tether itself. "
    "On May 12th, USDT briefly fell to 95 cents — a 5 percent depeg — on major exchanges. "
    "This was not contagion from an unrelated token. This was a direct bank-run-style attack on Tether's credibility. "
    "Investors were essentially asking: does Tether actually have the reserves it claims? "
    "In a single day, over 7 billion dollars of USDT was redeemed — the largest single-day redemption in Tether's history. "
    "Tether managed to maintain the peg and return to one dollar within days, which actually validated their reserves. "
    "But during those days of uncertainty, the forced selling mechanism activated again: "
    "same low buffer of around 5 to 8 percent, same constraint, same outcome. "
    "The CAR came in at plus 8.85 percentage points — essentially identical to the LUNA event three days earlier. "
    "We want to highlight that this near-identical magnitude is not a coincidence. "
    "It tells us that what drove the market impact was not the specific trigger — LUNA contagion versus direct Tether run — "
    "but rather the underlying constraint: Tether did not have enough cash to absorb redemptions "
    "without selling T-bills. The buffer level determined the outcome, not the nature of the shock.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Event 3: USDC / SVB Failure  +  Why Is the CAR Negative?
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "Event 3: USDC / SVB Failure — March 11, 2023")
slide_number(sld, 8)

box(sld, Inches(0.35), Inches(1.05), Inches(12.6), Inches(0.42), fill=GREEN)
txt(sld, "Buffer Condition: HIGH   |   L ≈ 18–22% cash   |   CAR = −18.01 pp***   |   t-stat > 30",
    Inches(0.55), Inches(1.08), Inches(12.1), Inches(0.34),
    size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# What happened
box(sld, Inches(0.35), Inches(1.55), Inches(5.9), Inches(2.5), fill=LGRAY)
txt(sld, "What happened", Inches(0.5), Inches(1.62),
    Inches(5.7), Inches(0.3), size=14, bold=True, color=HDR_NAVY)
tf = bullet_tb(sld, Inches(0.5), Inches(1.98), Inches(5.7), Inches(2.0))
events3 = [
    "SVB failed on March 10, 2023 — the second-largest US bank failure in history",
    "Circle disclosed: $3.3 billion of USDC reserves were held at SVB — temporarily inaccessible",
    "USDC depegged to ~$0.87 on March 11 — a larger nominal depeg than USDT in 2022",
    "But by March 13, after the US government guaranteed SVB deposits, USDC returned to $1.00",
]
for j, line in enumerate(events3):
    add_line(tf, f"❖  {line}", size=12, color=BODY_DARK, space_before=6 if j > 0 else 2)

# The key question: why negative, not flat?
box(sld, Inches(0.35), Inches(4.15), Inches(5.9), Inches(1.45), fill=HDR_NAVY)
txt(sld, "Professor's question: if Circle had enough reserves, why isn't the CAR flat?",
    Inches(0.5), Inches(4.22), Inches(5.7), Inches(0.55),
    size=13, bold=True, color=GOLD)
txt(sld,
    "Because the buffer condition tells us what Circle does. "
    "It does NOT tell us what OTHER investors do.",
    Inches(0.5), Inches(4.8), Inches(5.7), Inches(0.65),
    size=12, italic=True, color=WHITE)

# Two forces explanation (right side - two panels)
box(sld, Inches(6.5), Inches(1.55), Inches(6.5), Inches(1.88), fill=LIGHT_GRN)
box(sld, Inches(6.5), Inches(1.55), Inches(6.5), Inches(0.38), fill=GREEN)
txt(sld, "Force A: Circle's behavior (HIGH BUFFER)",
    Inches(6.65), Inches(1.58), Inches(6.2), Inches(0.3),
    size=13, bold=True, color=WHITE)
tf2 = bullet_tb(sld, Inches(6.65), Inches(2.0), Inches(6.2), Inches(1.3))
add_line(tf2, "Circle used its large cash buffer to meet redemptions", size=12, color=BODY_DARK, space_before=2)
add_line(tf2, "Did NOT need to sell T-bills", size=12, color=BODY_DARK, space_before=5)
add_line(tf2, "No forced selling pressure from the issuer side", size=12, color=BODY_DARK, space_before=5)
add_line(tf2, "→ This is what a flat CAR would look like IF nothing else happened", size=12,
         color=GREEN, bold=True, space_before=5)

box(sld, Inches(6.5), Inches(3.52), Inches(6.5), Inches(2.05), fill=LIGHT_RED)
box(sld, Inches(6.5), Inches(3.52), Inches(6.5), Inches(0.38), fill=RED)
txt(sld, "Force B: Market-wide flight to safety",
    Inches(6.65), Inches(3.55), Inches(6.2), Inches(0.3),
    size=13, bold=True, color=WHITE)
tf3 = bullet_tb(sld, Inches(6.65), Inches(3.97), Inches(6.2), Inches(1.5))
add_line(tf3, "SVB's failure caused widespread panic — investors fled to safety", size=12, color=BODY_DARK, space_before=2)
add_line(tf3, "Massive buying of U.S. T-bills (the safest short-term asset)", size=12, color=BODY_DARK, space_before=5)
add_line(tf3, "This BUYING pushed T-bill yields DOWN relative to OIS", size=12, color=BODY_DARK, space_before=5)
add_line(tf3, "→ Spread compressed far BELOW what the normal model predicted", size=12,
         color=RED, bold=True, space_before=5)

# Combined effect box
box(sld, Inches(6.5), Inches(5.65), Inches(6.5), Inches(0.7), fill=HDR_NAVY)
txt(sld,
    "Net effect: Force A (no selling) + Force B (active buying) = CAR of −18 pp\n"
    "The graph is NEGATIVE, not flat, because of pure market flight-to-safety — not issuer behavior.",
    Inches(6.65), Inches(5.7), Inches(6.2), Inches(0.6),
    size=11, bold=True, color=WHITE)

orange_callout(sld,
    "High buffer → no forced selling → the flight-to-safety buying dominates → CAR goes deeply negative. "
    "This is the OPPOSITE of a crisis, not the absence of one.")

notes(sld,
    "Our third event is the Silicon Valley Bank failure in March 2023, and this one needs careful explanation "
    "because the graph looks counterintuitive. "
    "SVB collapsed on March 10th — the second largest U.S. bank failure in history. "
    "Circle, which issues USDC, disclosed that 3.3 billion dollars of its reserves were held at SVB "
    "and temporarily inaccessible. USDC depegged to around 87 cents — actually a larger nominal depeg "
    "than Tether in 2022. "
    "So why does the CAR go to minus 18 instead of spiking upward like the 2022 events? "
    "The answer is the buffer condition. By early 2023, Circle had rebuilt its cash reserves to around "
    "18 to 22 percent — well above our 13 percent threshold. "
    "When redemptions came in, Circle paid from cash. It did not need to sell a single Treasury bill. "
    "So there was no forced selling pressure on the T-bill market from Circle's side. "
    "But here is what actually happened: SVB's collapse triggered a massive market-wide flight to safety. "
    "Investors across the entire financial system rushed to buy U.S. Treasury bills — "
    "the safest short-term asset available. That buying pressure pushed T-bill yields sharply below OIS, "
    "compressing the spread far below what our normal model would predict. "
    "That is why the CAR is minus 18 and not flat. "
    "The high buffer prevented the forced-selling story. But it could not prevent the flight-to-safety buying story. "
    "What we observe in the SVB event is the pure flight-to-safety effect, "
    "completely uncontaminated by any issuer-side selling pressure.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — All Three Events: The Natural Experiment
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "The Natural Experiment: Buffer Regime Determines the Outcome")
slide_number(sld, 9)

txt(sld,
    "Three major stablecoin crises. Similar scale of shock. Opposite outcomes. The only systematic difference: the liquid buffer level.",
    Inches(0.35), Inches(1.05), Inches(12.6), Inches(0.38),
    size=13, italic=True, color=MIDGRAY)

# Three event cards
events_data = [
    ("LUNA / UST\nCollapse", "May 9, 2022", "LOW", "~5–8%",
     "+8.91 pp", "T-bill forced selling\nSpread spiked sharply", RED, LIGHT_RED),
    ("USDT\nPartial Depeg", "May 12, 2022", "LOW", "~5–8%",
     "+8.85 pp", "Same constraint, same result\nNear-identical magnitude", RED, LIGHT_RED),
    ("USDC / SVB\nFailure", "Mar 11, 2023", "HIGH", "~18–22%",
     "−18.01 pp", "No forced selling\nFlight-to-safety buying", GREEN, LIGHT_GRN),
]
for i, (name, date, regime, pct, car, interp, col, bcol) in enumerate(events_data):
    l = Inches(0.35) + i * Inches(4.3)
    box(sld, l, Inches(1.55), Inches(4.1), Inches(0.45), fill=col)
    txt(sld, name.replace("\n", " "), l + Inches(0.08), Inches(1.58),
        Inches(3.95), Inches(0.35), size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(sld, l, Inches(2.0), Inches(4.1), Inches(3.5), fill=bcol)
    txt(sld, date, l + Inches(0.1), Inches(2.08),
        Inches(3.9), Inches(0.28), size=12, italic=True, color=MIDGRAY, align=PP_ALIGN.CENTER)
    txt(sld, f"Buffer: {regime}  ({pct})", l + Inches(0.1), Inches(2.42),
        Inches(3.9), Inches(0.3), size=12, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(sld, f"CAR = {car}", l + Inches(0.1), Inches(2.78),
        Inches(3.9), Inches(0.65), size=26, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(sld, interp, l + Inches(0.1), Inches(3.5),
        Inches(3.9), Inches(0.72), size=12, color=BODY_DARK, align=PP_ALIGN.CENTER)

# CAR chart
img(sld, RESULTS / "event_study_cars.png",
    Inches(0.35), Inches(3.9) - Inches(0.3), Inches(8.2))

# Right side summary
box(sld, Inches(8.7), Inches(3.6), Inches(4.3), Inches(2.25), fill=HDR_NAVY)
txt(sld, "What the 27 pp swing means:", Inches(8.85), Inches(3.67),
    Inches(4.0), Inches(0.35), size=13, bold=True, color=WHITE)
tf = bullet_tb(sld, Inches(8.85), Inches(4.08), Inches(4.0), Inches(1.65))
add_line(tf, "Low buffer avg CAR:  +8.9 pp", size=13, color=GOLD, bold=True, space_before=2)
add_line(tf, "High buffer CAR:  −18.0 pp", size=13, color=GOLD, bold=True, space_before=6)
add_line(tf, "Difference:  26.9 pp swing", size=15, color=WHITE, bold=True, space_before=8)
add_line(tf, "Welch t = 15.22***\np < 0.001", size=12, color=RGBColor(0xAA, 0xC4, 0xE8), space_before=6)

orange_callout(sld,
    "Same type of shock, same asset class. The buffer condition alone explains a 27 pp swing in Treasury market impact. "
    "This is the core finding of the paper.")

notes(sld,
    "This slide brings all three events together and shows why we call this a natural experiment. "
    "Look at what is the same across all three events: stablecoin issuers, U.S. Treasury bills, "
    "a sudden confidence shock, and mass redemptions. "
    "The only systematic difference is the liquid buffer level. "
    "In 2022, the buffer was low — around 5 to 8 percent — and both events produced a CAR of about plus 9 percentage points. "
    "In 2023, the buffer was high — around 18 to 22 percent — and the event produced a CAR of minus 18 percentage points. "
    "That is a swing of 27 percentage points driven entirely by the buffer condition. "
    "The Welch t-statistic comparing the two groups is 15.22 with a p-value below 0.001, "
    "meaning this difference is not remotely explained by chance. "
    "This is the core empirical finding of our paper: the liquid cash buffer is what determines "
    "whether a stablecoin issuer amplifies or absorbs a shock to the Treasury market. "
    "And next week, the placebo test will confirm that these large CARs only appear on real crisis days — "
    "not on ordinary days with the same buffer conditions.")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Thank You
# ═══════════════════════════════════════════════════════════════════════════════
sld = blank(prs)
title_slide_bg(sld)

txt(sld, "THANK YOU",
    Inches(4.0), Inches(2.85), Inches(8.8), Inches(0.6),
    size=28, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font="Calibri Light")

box(sld, Inches(3.9), Inches(4.45), Inches(7.5), Inches(2.8),
    fill=RGBColor(0x0E, 0x26, 0x68))

txt(sld,
    "Yonsei GSIS 2026-1  |  Topics in International Finance\n"
    "Mireu Kim · Sara Chekroune · Oybek Ibragimov\n"
    "kimmireu0921@gmail.com",
    Inches(4.1), Inches(4.6), Inches(7.2), Inches(2.5),
    size=16, color=WHITE, align=PP_ALIGN.LEFT)

notes(sld,
    "Thank you for listening. To summarize in two sentences: "
    "stablecoin issuers have become structurally important buyers of U.S. Treasury bills, "
    "which deepens America's safe-asset privilege in normal times. "
    "But when their liquid cash buffers fall below roughly 13 percent of supply, "
    "a run forces them to liquidate T-bills — turning a privilege amplifier into a systemic disruptor. "
    "We are happy to take any questions.")


# ── Save ──────────────────────────────────────────────────────────────────────
out = PRESENTATIONS / "0526_Stablecoin_Exorbitant_Privilege.pptx"
prs.save(str(out))
print(f"Saved: {out}  ({len(prs.slides)} slides)")
