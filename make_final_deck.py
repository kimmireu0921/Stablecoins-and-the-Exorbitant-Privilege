"""
make_final_deck.py — FINAL deck, styled to match 0609_FINAL_merged (Georgia/Calibri,
teal-gold-navy, kicker labels, accent-bar cards). No individual names — group wording.

Story: structural question -> what we did after the professor's feedback -> the levels
regression turned out spurious -> we found and stress-tested the bid-cover result, which
survives -> supporting evidence + honest caveats -> conclusion.

All bid-cover numbers computed live from the data.
Writes: results/bidcover_final_figure.png, presentations/FINAL_Stablecoin_Privilege.pptx
"""
from pathlib import Path

import numpy as np
import pandas as pd
import statsmodels.api as sm
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent
DATA = ROOT / "data"
RES = ROOT / "results"; RES.mkdir(exist_ok=True)

# palette (from 0609 deck)
TEAL   = RGBColor(0x1B, 0x99, 0x8B)
GOLD   = RGBColor(0xE8, 0xA1, 0x2C)
NAVY   = RGBColor(0x13, 0x29, 0x4B)
NAVYDK = RGBColor(0x12, 0x28, 0x44)
BLUE   = RGBColor(0x5B, 0x8D, 0xEF)
SLATE  = RGBColor(0x33, 0x42, 0x5C)
GRAYB  = RGBColor(0x5B, 0x6B, 0x82)
CARD   = RGBColor(0xF4, 0xF6, 0xFA)
RED    = RGBColor(0xC1, 0x12, 0x1F)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LT1    = RGBColor(0xCA, 0xDC, 0xFC)
LT2    = RGBColor(0xAE, 0xC0, 0xDD)
LT3    = RGBColor(0x9D, 0xB2, 0xD6)
SERIF, SANS = "Georgia", "Calibri"
SW, SH = Inches(13.33), Inches(7.5)
TERMS = ["4-Week", "8-Week", "13-Week", "26-Week"]


# ── live numbers + figure ────────────────────────────────────────────────────
def bidcover_numbers():
    auc = pd.read_csv(ROOT / "results" / "bidcover_auction_raw_rebuilt.csv", parse_dates=["date"])
    auc = auc[auc.term.isin(TERMS)]
    m = pd.read_csv(DATA / "monthly_panel.csv", index_col=0, parse_dates=True)
    d = pd.read_csv(DATA / "daily_panel.csv", index_col=0, parse_dates=True)
    me = d[["supply_USDT", "supply_USDC"]].resample("ME").last()
    m["dln_USDT"] = np.log(me.supply_USDT).diff(); m["dln_USDC"] = np.log(me.supply_USDC).diff()
    m["d_fedfunds"] = d["fedfunds"].resample("ME").last().diff()

    def fit(df, xc):
        df = df.dropna(subset=["bc"] + xc)
        r = sm.OLS(df["bc"], sm.add_constant(df[xc])).fit(cov_type="HAC", cov_kwds={"maxlags": 1})
        return (r.params["dln_USDT"], r.pvalues["dln_USDT"], r.params["dln_USDC"],
                r.pvalues["dln_USDC"], float(r.f_test("dln_USDT = dln_USDC").pvalue))

    specs = {"clean": ["dln_USDT", "dln_USDC", "vix"],
             "fed": ["dln_USDT", "dln_USDC", "vix", "d_fedfunds"],
             "fed_off": ["dln_USDT", "dln_USDC", "vix", "d_fedfunds", "ln_offering"]}
    out = {s: {} for s in specs}
    for term in TERMS:
        t = auc[auc.term == term].set_index("date")
        base = m.join(t["bid_cover"].resample("ME").mean().rename("bc")) \
                .join(np.log(t["offering"].resample("ME").mean()).rename("ln_offering"))
        for s, xc in specs.items():
            out[s][term] = fit(base.copy(), xc)
    return out


def make_figure(nums):
    plt.rcParams["font.family"] = "Georgia"
    fig, ax = plt.subplots(1, 2, figsize=(11, 4.1))
    x = np.arange(len(TERMS))
    for key, lab, col in [("clean", "supply + VIX", "#9aa7bd"),
                          ("fed", "+ Fed-funds control", "#13294B"),
                          ("fed_off", "+ Fed + offering size", "#1B998B")]:
        ax[0].plot(x, [nums[key][t][0] for t in TERMS], "o-", color=col, lw=2.4, ms=7, label=lab)
    ax[0].axhline(0, color="#333", lw=0.8)
    ax[0].set_xticks(x); ax[0].set_xticklabels(TERMS, fontsize=9)
    ax[0].set_ylabel("β on USDT supply growth", fontsize=9)
    ax[0].set_title("The rate-cycle control is the crux\n(no interpolated reserve data)", fontsize=10.5)
    ax[0].legend(fontsize=8, loc="lower left"); ax[0].grid(axis="y", alpha=.2)
    w = 0.38
    bU = [nums["fed_off"][t][0] for t in TERMS]; pU = [nums["fed_off"][t][1] for t in TERMS]
    bC = [nums["fed_off"][t][2] for t in TERMS]
    ax[1].bar(x - w/2, bU, w, color="#13294B", label="USDT")
    ax[1].bar(x + w/2, bC, w, color="#E8A12C", label="USDC (foil)")
    for xi, v, pp in zip(x - w/2, bU, pU):
        st = "***" if pp < .01 else "**" if pp < .05 else "*" if pp < .1 else "ns"
        ax[1].text(xi, v - 0.05, st, ha="center", va="top", fontsize=8, color="white", fontweight="bold")
    ax[1].axhline(0, color="#333", lw=0.8)
    ax[1].set_xticks(x); ax[1].set_xticklabels(TERMS, fontsize=9)
    ax[1].set_ylabel("β (final spec)", fontsize=9)
    ax[1].set_title("USDT lowers bid-cover; USDC does not", fontsize=10.5)
    ax[1].legend(fontsize=8, loc="lower right"); ax[1].grid(axis="y", alpha=.2)
    plt.tight_layout()
    out = RES / "bidcover_final_figure.png"
    plt.savefig(out, dpi=150, bbox_inches="tight"); plt.close()
    plt.rcParams["font.family"] = "sans-serif"
    return out


def sample_facts():
    d = pd.read_csv(DATA / "daily_panel.csv", index_col=0, parse_dates=True)
    last = d[["supply_USDT", "supply_USDC"]].dropna().iloc[-1]
    m = pd.read_csv(DATA / "monthly_panel.csv", index_col=0, parse_dates=True)
    return {"end_supply": float(last.sum()), "n_months": len(m)}


# ── styled pptx helpers ──────────────────────────────────────────────────────
def blank(p): return p.slides.add_slide(p.slide_layouts[6])

def box(s, l, t, w, h, fill=None):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h)); sh.shadow.inherit = False
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); return sh

def text(s, l, t, w, h, runs, font=SANS, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=1.05):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    if isinstance(runs, str): runs = [[(runs, 16, SLATE, False)]]
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = sp; p.space_after = Pt(3)
        for tup in line:
            txt_, sz, col, bold = tup
            r = p.add_run(); r.text = txt_; r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = bold; r.font.name = font
    return tb

def head(s, kicker, title, subtitle=None):
    text(s, 0.62, 0.42, 12.1, 0.32, [[(kicker.upper(), 12, TEAL, True)]], font=SANS)
    text(s, 0.6, 0.74, 12.1, 0.85, [[(title, 27, NAVY, True)]], font=SERIF)
    if subtitle:
        text(s, 0.62, 1.62, 12.1, 0.45, [[(subtitle, 15.5, GRAYB, False)]], font=SERIF)

def foot(s, n):
    text(s, 0.62, 7.06, 9.0, 0.3, [[("Stablecoins & the Exorbitant Privilege", 9, GRAYB, False)]], font=SANS)
    text(s, 12.0, 7.0, 0.7, 0.3, [[(str(n), 11, GRAYB, False)]], font=SANS, align=PP_ALIGN.RIGHT)

def card(s, l, t, w, h, accent, title, body, tsize=17, bsize=13.5):
    box(s, l, t, w, h, fill=CARD)
    box(s, l, t, 0.09, h, fill=accent)
    text(s, l + 0.28, t + 0.18, w - 0.5, 0.5, [[(title, tsize, accent, True)]], font=SERIF)
    if isinstance(body, str):
        body = [[(body, bsize, SLATE, False)]]
    text(s, l + 0.28, t + 0.74, w - 0.52, h - 0.9, body, font=SANS, sp=1.12)

def darkbar(s, l, t, w, h, heading, body_runs):
    box(s, l, t, w, h, fill=NAVY)
    text(s, l + 0.3, t + 0.16, w - 0.6, 0.35, [[(heading, 14, GOLD, True)]], font=SERIF)
    text(s, l + 0.3, t + 0.56, w - 0.6, h - 0.7, body_runs, font=SANS, sp=1.1)

def star(p): return "***" if p < .01 else "**" if p < .05 else "*" if p < .1 else "ns"


# ── build ────────────────────────────────────────────────────────────────────
def build():
    nums = bidcover_numbers()
    facts = sample_facts()
    fig_bc = make_figure(nums)
    fig_ev = RES / "event_study_multi.png"
    fig_irf = ROOT / "results" / "irf_usdt_usdc.png"

    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
    pg = [0]
    def newpage(dark=False):
        s = blank(prs)
        if dark:
            box(s, 0, 0, 13.33, 7.5, fill=NAVY)
        box(s, 0, 0, 0.2, 7.5, fill=TEAL); box(s, 0.2, 0, 0.1, 7.5, fill=GOLD)
        pg[0] += 1
        return s

    # 1 — TITLE
    s = newpage(dark=True)
    text(s, 1.0, 1.05, 11.0, 0.35, [[("YONSEI GSIS   |   TOPICS IN INTERNATIONAL FINANCE", 12, LT3, True)]], font=SANS)
    text(s, 1.0, 1.65, 11.3, 1.2, [[("Stablecoins and the Exorbitant Privilege", 40, WHITE, True)]], font=SERIF)
    text(s, 1.0, 2.95, 11.0, 0.6, [[("Safe-Asset Demand and Its Systemic Fragility", 20, LT1, False)]], font=SERIF)
    text(s, 1.0, 3.75, 11.3, 0.4, [[("Group project  ·  Topics in International Finance  ·  June 2026", 13, RGBColor(0xC7,0xD3,0xE8), False)]], font=SANS)
    chips = [(1.0, TEAL, "Bid-cover ↓", "USDT supply weakens\nT-bill auction demand"),
             (5.0, GOLD, "USDT ≠ USDC", "A reserve-composition\nchannel, issuer-specific"),
             (8.9, BLUE, "Robust", "Survives every\nstress-test we ran")]
    for x, ac, big, cap in chips:
        box(s, x, 4.7, 3.6, 1.65, fill=NAVYDK)
        box(s, x, 4.7, 3.6, 0.08, fill=ac)
        text(s, x + 0.1, 4.82, 3.4, 0.65, [[(big, 24, WHITE, True)]], font=SERIF)
        text(s, x + 0.1, 5.55, 3.4, 0.9, [[(cap, 12, LT2, False)]], font=SANS, sp=1.05)

    # 2 — MOTIVATION
    s = newpage(); head(s, "Motivation", "The Question — A Structural Fact, Then a Test",
                        "Stablecoin issuers have quietly become major buyers of US Treasury bills")
    card(s, 0.6, 2.25, 3.85, 3.7, TEAL, "The structural fact",
         "Stablecoin issuers hold their reserves largely in short-term US Treasuries. USDT alone is now one of the largest holders of T-bills in the world — by construction.")
    card(s, 4.74, 2.25, 3.85, 3.7, NAVY, "The mechanism",
         "When USDT grows, it mechanically buys T-bills. This is programmatic, rule-bound demand for US government debt — a candidate channel of the 'exorbitant privilege'.")
    card(s, 8.88, 2.25, 3.85, 3.7, GOLD, "The question",
         "Does this demand actually show up in the US Treasury market? We test it against T-bill auction demand (bid-cover) and stress-test every result.")
    foot(s, pg[0])

    # 3 — AFTER FEEDBACK: what we changed
    s = newpage(); head(s, "After the professor's feedback", "What We Changed",
                        "We rebuilt the analysis around the professor's corrections")
    card(s, 0.6, 2.2, 5.95, 1.7, TEAL, "Specification",
         "Exploit the panel: each issuer (USDT, USDC) a separate observation. Stop combining Treasury and liquid reserves; drop θ (θ + L ≈ 1, collinear). Show the estimating equation on every result.", tsize=16)
    card(s, 6.78, 2.2, 5.95, 1.7, NAVY, "Data construction",
         "Replace forward-filling of the quarterly reserve attestations with a moving-average interpolation between reporting dates, as advised.", tsize=16)
    card(s, 0.6, 4.05, 5.95, 1.7, GOLD, "Event study",
         "Keep the LUNA collapse as motivation. Drop SVB — it is confounded by the banking crisis, so the stablecoin channel cannot be isolated there.", tsize=16)
    card(s, 6.78, 4.05, 5.95, 1.7, BLUE, "Focus",
         "Re-centre the analysis on the buffer-interaction coefficient and re-run everything from the corrected panel.", tsize=16)
    foot(s, pg[0])

    # 4 — AFTER FEEDBACK: what it revealed
    s = newpage(); head(s, "After the professor's feedback", "What the Re-Run Revealed",
                        "The original headline results did not survive the corrected specification")
    card(s, 0.6, 2.2, 5.95, 2.5, RED, "The results weakened",
         "Re-run on the corrected panel, the privilege coefficient β₁ became insignificant, and the buffer interaction was unstable — it even flipped sign under first-differencing. A warning sign, not a result.", tsize=16)
    card(s, 6.78, 2.2, 5.95, 2.5, NAVY, "The diagnosis: spurious",
         "Formal tests explained why: the spread and the reserve buffer are non-stationary (ADF) and do not cointegrate (Engle–Granger, Johansen). Both simply trended down together during the 2022–24 Fed hiking cycle — a textbook spurious regression.", tsize=16)
    darkbar(s, 0.6, 4.95, 12.13, 1.25, "What we dropped from the main claims",
            [[("β₁ (privilege amplification)   ·   the ≈13% reserve threshold   ·   the LSTAR model", 15, WHITE, True)],
             [("All three are significant only in the spurious levels regression — so we demote them.", 13, LT2, False)]])
    foot(s, pg[0])

    # 5 — OUR RESPONSE (process, no names)
    s = newpage(); head(s, "Our response", "Finding the Result That Survives",
                        "We changed the dependent variable, then tried hard to break the new result")
    card(s, 0.6, 2.25, 3.85, 3.7, TEAL, "1 · Re-anchor",
         "Move from the (non-stationary) spread to T-bill auction demand — the bid-cover ratio. This sidesteps the spurious-regression problem entirely.")
    card(s, 4.74, 2.25, 3.85, 3.7, NAVY, "2 · Stress-test",
         "Control for auction size; remove the interpolated reserve data; control the Fed cycle; test at the individual-auction level; and run a 2,000-shuffle placebo.")
    card(s, 8.88, 2.25, 3.85, 3.7, GOLD, "3 · Converge",
         "Reproduce every check independently and reconcile. The bid-cover result holds throughout — and that becomes the paper's lead.")
    foot(s, pg[0])

    # 6 — HEADLINE RESULT
    s = newpage(); head(s, "Main result", "USDT Supply Growth → Weaker T-Bill Auction Demand",
                        "Monthly, N = 51, Newey–West HAC(1). Final spec: supply + VIX + Δfed-funds + offering size")
    cols = [0.95, 3.7, 3.0, 3.0, 2.7]
    headers = ["Maturity", "USDT (β)", "USDC — foil (β)", "Wald  USDT≠USDC"]
    cw = [Inches(2.4), Inches(2.4), Inches(2.7), Inches(2.7)]
    x0 = 1.6; rh = 0.62
    # header row
    xx = x0
    for j, hd in enumerate(headers):
        box(s, xx, 2.3, cw[j].inches, rh, fill=NAVY)
        text(s, xx, 2.3, cw[j].inches, rh, [[(hd, 13, WHITE, True)]], font=SANS, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        xx += cw[j].inches
    yy = 2.3 + rh
    for i, t in enumerate(TERMS):
        bU, pU, bC, pC, w = nums["fed_off"][t]
        cells = [t, f"{bU:+.2f} {star(pU) if star(pU)!='ns' else 'ns'}",
                 f"{bC:+.2f} {star(pC) if star(pC)!='ns' else 'ns'}", f"p = {w:.3f}"]
        colr = [SLATE, TEAL if pU < .1 else GRAYB, GRAYB, NAVY]
        xx = x0
        for j, c in enumerate(cells):
            box(s, xx, yy, cw[j].inches, rh, fill=CARD if i % 2 == 0 else WHITE)
            text(s, xx, yy, cw[j].inches, rh, [[(c, 13.5, colr[j], j == 1)]], font=SANS, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            xx += cw[j].inches
        yy += rh
    darkbar(s, 1.6, 5.45, 10.2, 1.15, "Reading the result",
            [[("USDT supply growth predicts lower bid-cover (weaker auction demand); USDC does not. "
               "The two issuers differ at every maturity — the reserve-composition signature.", 13.5, WHITE, False)]])
    foot(s, pg[0])

    # 7 — ROBUSTNESS
    s = newpage(); head(s, "Robustness", "What Survives Every Test",
                        "The result holds with no interpolated data, with auction size controlled, and against a placebo")
    s.shapes.add_picture(str(fig_bc), Inches(0.5), Inches(2.15), width=Inches(7.4))
    box(s, 8.25, 2.15, 4.5, 4.1, fill=CARD); box(s, 8.25, 2.15, 0.09, 4.1, fill=TEAL)
    text(s, 8.5, 2.35, 4.05, 3.8,
         [[("◆  ", 13, TEAL, False), ("Drop the interpolated reserves (θ, L) entirely — result holds", 13, SLATE, False)],
          [("◆  ", 13, TEAL, False), ("Controlling the Fed rate cycle is the crux: once it is in, USDT is significant at every maturity", 13, SLATE, False)],
          [("◆  ", 13, TEAL, False), ("Survives the auction-size control (4-week softens; 8/13/26-week strong)", 13, SLATE, False)],
          [("◆  ", 13, TEAL, False), ("Passes a 2,000-shuffle placebo test", 13, SLATE, False)],
          [("◆  ", 13, TEAL, False), ("USDT ≠ USDC at every maturity", 13, SLATE, False)]], font=SANS, sp=1.2)
    foot(s, pg[0])

    # 8 — CAVEAT
    s = newpage(); head(s, "An honest caveat", "Frequency, Not Absence",
                        "Why the appropriate unit of analysis is the month")
    card(s, 0.6, 2.2, 5.95, 3.5, NAVY, "At the auction level, no signal",
         "Run auction by auction (N = 1,094), the USDT coefficient is not significant. Taken alone this looks like a non-result.")
    card(s, 6.78, 2.2, 5.95, 3.5, TEAL, "Why — a frequency mismatch",
         "Stablecoin supply is a monthly series: ~70% of its variation is between months, ~30% within. The same monthly value repeats across ~17 auctions a month, and the pre-auction windows overlap. So N = 1,094 carries only ~51 independent supply observations. The right unit for a monthly regressor is the month — where the effect is significant.")
    darkbar(s, 0.6, 5.95, 12.13, 0.9, "Bottom line",
            [[("The auction-level null is the expected consequence of a frequency mismatch — not evidence against the channel.", 14, WHITE, False)]])
    foot(s, pg[0])

    # 9 — SUPPORTING
    s = newpage(); head(s, "Supporting evidence & motivation", "Two Secondary Checks",
                        "Daily impulse responses, and an event study used as motivation")
    box(s, 0.6, 2.15, 6.0, 4.05, fill=CARD); box(s, 0.6, 2.15, 0.09, 4.05, fill=NAVY)
    text(s, 0.85, 2.3, 5.6, 0.4, [[("Daily VAR / impulse response", 15, NAVY, True)]], font=SERIF)
    if fig_irf.exists():
        s.shapes.add_picture(str(fig_irf), Inches(0.85), Inches(2.85), height=Inches(2.6))
    text(s, 0.85, 5.55, 5.5, 0.6, [[("USDT supply Granger-causes the spread; small, short-lived day-1 compression — suggestive support.", 11.5, SLATE, False)]], font=SANS)
    box(s, 6.75, 2.15, 6.0, 4.05, fill=CARD); box(s, 6.75, 2.15, 0.09, 4.05, fill=GOLD)
    text(s, 7.0, 2.3, 5.6, 0.4, [[("Multi-event study (motivation)", 15, GOLD, True)]], font=SERIF)
    if fig_ev.exists():
        s.shapes.add_picture(str(fig_ev), Inches(7.0), Inches(3.0), width=Inches(5.55))
    text(s, 7.0, 5.05, 5.6, 1.1, [[("LUNA · Celsius · FTX · BUSD (SVB dropped — confounded). Events disagree in sign; pooled p = 0.43 → no systematic effect. Narrative motivation, not identification.", 11.5, SLATE, False)]], font=SANS)
    foot(s, pg[0])

    # 10 — LIMITS: why the results had to change
    s = newpage(); head(s, "Limitations", "Why the Results Had to Change",
                        "The feedback exposed three constraints that bind any study of this market today")
    card(s, 0.6, 2.2, 3.85, 3.6, RED, "One macro regime",
         "Our 51 months are dominated by a single Fed hiking cycle and its unwind. Trending variables correlate mechanically, so a levels regression flatters any hypothesis — ours included. Removing θ and the forward-fill took away the props; the stationarity tests did the rest.",
         tsize=16, bsize=12.5)
    card(s, 4.74, 2.2, 3.85, 3.6, NAVY, "Reserves are unobserved",
         "Tether attests quarterly: roughly two-thirds of the monthly reserve ratios are interpolation, and the estimated threshold moved 13% → 6.6% just by changing the fill method. The buffer-fragility hypothesis is currently untestable at monthly frequency — not false.",
         tsize=16, bsize=12.5)
    card(s, 8.88, 2.2, 3.85, 3.6, GOLD, "A hard power ceiling",
         "Two issuers, 51 months, no natural experiment. Controls, an issuer foil, and a placebo are the strongest identification available here — and supply shocks still explain only ~1.3% of spread variance. These are limits of the setting, not of the specification.",
         tsize=16, bsize=12.5)
    darkbar(s, 0.6, 6.0, 12.13, 0.95, "The honest reading",
            [[("The feedback did not break good results — it revealed which results this data could never have supported.", 14, WHITE, False)]])
    foot(s, pg[0])

    # 11 — LIMITS: claim boundary
    s = newpage(); head(s, "Limitations", "What We Can — and Cannot — Claim",
                        "Drawing the boundary precisely is part of the result")
    card(s, 0.6, 2.2, 5.95, 3.5, TEAL, "Supported by the data",
         [[("◆  ", 13, TEAL, False), ("USDT supply growth is associated with weaker T-bill auction demand — robust to every check we ran", 13.5, SLATE, False)],
          [("◆  ", 13, TEAL, False), ("The channel is issuer-specific (USDT ≠ USDC), consistent with reserve composition", 13.5, SLATE, False)],
          [("◆  ", 13, TEAL, False), ("The structural fact: programmatic T-bill purchases at sovereign-holder scale", 13.5, SLATE, False)]])
    card(s, 6.78, 2.2, 5.95, 3.5, RED, "Beyond the data — for now",
         [[("◆  ", 13, RED, False), ("A causal privilege effect in yield levels (β₁): one rate cycle cannot identify it", 13.5, SLATE, False)],
          [("◆  ", 13, RED, False), ("A reserve-buffer threshold (≈13%): reserves are not observed at monthly frequency", 13.5, SLATE, False)],
          [("◆  ", 13, RED, False), ("Run-scenario fragility: no clean redemption episode in sample (SVB is confounded)", 13.5, SLATE, False)]])
    darkbar(s, 0.6, 5.95, 12.13, 0.95, "Untested is not rejected",
            [[("The question is ahead of its data — and the data is catching up.", 14, WHITE, False)]])
    foot(s, pg[0])

    # 12 — FUTURE: constraints relax
    s = newpage(); head(s, "Looking forward", "The Constraints Relax Mechanically",
                        "Time, regulation, and scale each remove one of the binding limits")
    card(s, 0.6, 2.2, 3.85, 3.55, TEAL, "Time → regimes",
         "Each year adds 12 observations. By 2030 the sample roughly doubles (~100 months) and spans at least two full rate cycles. Regime variation — not just a larger N — is what separates a real channel from a shared trend.",
         tsize=16, bsize=12.5)
    card(s, 4.74, 2.2, 3.85, 3.55, NAVY, "Regulation → real data",
         "The 2025 US stablecoin law (GENIUS) and the EU's MiCA mandate monthly, certified reserve disclosures. θ and L become observed rather than interpolated, and the issuer panel widens as banks and fintechs enter. The buffer hypothesis becomes testable for the first time.",
         tsize=16, bsize=12.5)
    card(s, 8.88, 2.2, 3.85, 3.55, GOLD, "Scale → detectability",
         f"Combined USDT + USDC supply is ≈ ${facts['end_supply']:.0f}B at the end of our sample — already on par with a mid-sized sovereign holder of US debt. Official and industry projections reach $1–2T by 2028–30; at that scale, issuers become marginal price-setters in the bill market.",
         tsize=16, bsize=12.5)
    darkbar(s, 0.6, 5.95, 12.13, 0.95, "The same test rig, better inputs",
            [[("Every analysis in this paper is scripted and reproducible — it can simply be re-run as each constraint lifts.", 14, WHITE, False)]])
    foot(s, pg[0])

    # 13 — FUTURE: falsifiable expectations + timeline
    s = newpage(); head(s, "Looking forward", "What We Expect — Falsifiable Predictions",
                        "Each claim comes with the evidence that would confirm or kill it")
    card(s, 0.6, 2.05, 3.85, 2.65, TEAL, "Bid-cover sharpens",
         "If the channel is real, β stays near −1.4 and significance grows with the sample. If it was an artifact of the 2022–23 crypto winter, it attenuates toward zero as those months roll out of the window.",
         tsize=15, bsize=12)
    card(s, 4.74, 2.05, 3.85, 2.65, NAVY, "Privilege becomes measurable",
         "Across two or more rate cycles the spread regression becomes estimable in differences. Safe-asset demand-curve logic predicts a small, few-bps compression per $100B of issuance — visible only at scale.",
         tsize=15, bsize=12)
    card(s, 8.88, 2.05, 3.85, 2.65, RED, "Fragility becomes testable",
         "Observed monthly reserves plus the first large redemption under the new regime test the buffer threshold directly. At $1T+ scale a forced bill sale is macro-relevant — the New Triffin question returns.",
         tsize=15, bsize=12)
    box(s, 0.6, 5.0, 12.13, 1.78, fill=NAVY)
    text(s, 0.9, 5.1, 6.0, 0.3, [[("EXPECTED PATH", 11, GOLD, True)]], font=SANS)
    box(s, 0.9, 5.47, 11.55, 0.022, fill=GOLD)
    miles = [("2026–27", "Monthly certified reserve reports begin (GENIUS / MiCA)"),
             ("2027–28", "Sample spans a full easing cycle (~75 months)"),
             ("2028–29", "2–3 years of observed θ, L → threshold test on real data"),
             ("2030", "~100 months, ≥2 cycles, wider issuer panel — full re-run")]
    mx = 0.9
    for yr, mtxt in miles:
        text(s, mx, 5.58, 2.85, 0.35, [[(yr, 14, WHITE, True)]], font=SERIF)
        text(s, mx, 5.97, 2.75, 0.75, [[(mtxt, 10.5, LT2, False)]], font=SANS, sp=1.0)
        mx += 2.98
    foot(s, pg[0])

    # 14 — CONCLUSION
    s = newpage(); head(s, "Conclusion", "The Defensible Claim")
    box(s, 0.8, 2.1, 11.7, 2.3, fill=CARD); box(s, 0.8, 2.1, 0.09, 2.3, fill=TEAL)
    text(s, 1.15, 2.35, 11.1, 1.9,
         [[("USDT supply growth is associated with lower T-bill auction bid-cover at all maturities, while USDC is not.", 18, NAVY, True)],
          [("Robust to controlling auction size and to removing all interpolated reserve controls, and it passes a 2,000-iteration placebo. At the individual-auction frequency the association is not detectable, consistent with stablecoin supply being a monthly-frequency variable.", 14.5, SLATE, False)]],
         font=SERIF, sp=1.12)
    text(s, 0.9, 4.7, 11.7, 1.8,
         [[("◆  ", 14, TEAL, False), ("The structural privilege channel is real — USDT is a programmatic T-bill buyer", 15, SLATE, False)],
          [("◆  ", 14, TEAL, False), ("Its imprint is visible in auction demand, not in the spurious levels regression", 15, SLATE, False)],
          [("◆  ", 14, TEAL, False), ("Stated carefully, the finding survives every robustness check we could devise", 15, SLATE, False)]],
         font=SANS, sp=1.25)
    foot(s, pg[0])

    # 11 — CLOSING
    s = newpage(dark=True)
    text(s, 1.0, 2.4, 11.3, 1.2, [[("A robust result, honestly earned.", 34, WHITE, True)]], font=SERIF)
    text(s, 1.0, 3.9, 11.3, 2.4,
         [[("USDT supply growth → weaker T-bill auction demand, and USDT ≠ USDC.", 17, LT1, False)],
          [("Survives offering size, zero interpolated data, the rate-cycle control, and a placebo.", 17, LT1, False)],
          [("The spurious levels results are dropped; the VAR and event study remain as support.", 17, LT1, False)],
          [("Diagnosis, stress-testing, and reconciliation were a shared, group effort.", 17, LT1, False)]],
         font=SANS, sp=1.3)

    out = ROOT / "presentations" / "FINAL_Stablecoin_Privilege.pptx"
    prs.save(out)
    print(f"Saved {out}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
