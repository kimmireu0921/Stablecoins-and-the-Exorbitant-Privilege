"""
update_slide4.py
Updates slide 4 (index 3, 4A: "What the Re-Run Revealed"):
  - Replaces dense β₁/β₄ paragraphs with clean bullet-point layout
  - Expands results card to fit the new content
  - Shifts the "What we dropped" dark bar down to match
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

TEAL   = RGBColor(0x1B, 0x99, 0x8B)
GOLD   = RGBColor(0xE8, 0xA1, 0x2C)
NAVY   = RGBColor(0x13, 0x29, 0x4B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DKGRAY = RGBColor(0x33, 0x42, 0x5C)
MGRAY  = RGBColor(0x5B, 0x6B, 0x82)
SANS   = "Calibri"

ANS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def get_shape(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None


def clear_tf(tf):
    txBody = tf._txBody
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)
    txBody.append(parse_xml(f'<a:p xmlns:a="{ANS}"/>'))


def fp(tf, text, font, sz, bold=False, color=NAVY,
       align=PP_ALIGN.LEFT, sb=0, sa=0):
    p = tf.paragraphs[0]
    p.alignment = align
    if sb: p.space_before = Pt(sb)
    if sa: p.space_after  = Pt(sa)
    r = p.add_run()
    r.text = text; r.font.name = font; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.color.rgb = color


def ap(tf, text, font, sz, bold=False, color=NAVY,
       align=PP_ALIGN.LEFT, sb=0, sa=0):
    p = tf.add_paragraph()
    p.alignment = align
    if sb: p.space_before = Pt(sb)
    if sa: p.space_after  = Pt(sa)
    r = p.add_run()
    r.text = text; r.font.name = font; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.color.rgb = color


# ══════════════════════════════════════════════════════════════════════════════
prs = Presentation("presentations/0616_Stablecoin_Exorbitant_Privilege.pptx")
s4  = prs.slides[3]

CARD_TOP     = 2011680
OLD_CARD_H   = 2286000
NEW_CARD_H   = 2950000   # expanded to fit bullet content
OLD_BAR_TOP  = 4526280   # Rectangle 14 original top
GAP          = OLD_BAR_TOP - (CARD_TOP + OLD_CARD_H)   # preserve the gap
NEW_BAR_TOP  = CARD_TOP + NEW_CARD_H + GAP

# ── 1. Expand results card background & accent bar ────────────────────────────
get_shape(s4, "Rectangle 6").height = NEW_CARD_H
get_shape(s4, "Rectangle 7").height = NEW_CARD_H

# ── 2. Expand and rewrite card body (TextBox 9) ───────────────────────────────
tb9 = get_shape(s4, "TextBox 9")
tb9.height = NEW_CARD_H - 700000   # leaves room below title
tf = tb9.text_frame
tf.word_wrap = True
clear_tf(tf)

# β₁ block
fp(tf, "β₁ — Privilege Amplification", SANS, 12, bold=True, color=NAVY, sa=4)
ap(tf, "→  Tests whether stablecoin supply amplifies spread sensitivity to reserve buffers",
   SANS, 10.5, color=DKGRAY, sb=2)
ap(tf, "→  Before: β₁ = −6.02 (p < 0.01) — statistically significant",
   SANS, 10.5, color=DKGRAY)
ap(tf, "→  After re-run: β₁ = +2.74 — not significant; sign flipped",
   SANS, 10.5, bold=True, color=GOLD)
ap(tf, "→  Verdict: amplification claim does not hold",
   SANS, 10.5, color=DKGRAY, sa=2)

# β₄ block
ap(tf, "β₄ — Buffer Interaction", SANS, 12, bold=True, color=NAVY, sb=14, sa=4)
ap(tf, "→  Tests direct effect of reserve buffers on bid-cover spreads",
   SANS, 10.5, color=DKGRAY, sb=2)
ap(tf, "→  Levels regression: β₄ = −35.89 (p = 0.032) — appeared significant",
   SANS, 10.5, color=DKGRAY)
ap(tf, "→  First-differenced: β₄ = +8.86 (p = 0.31) — sign reversed, insignificant",
   SANS, 10.5, bold=True, color=GOLD)
ap(tf, "→  Sign reversal under differencing = textbook spurious regression signal",
   SANS, 10.5, color=DKGRAY)

# ── 3. Shift the "What we dropped" dark bar down ─────────────────────────────
delta = NEW_BAR_TOP - OLD_BAR_TOP

for name in ["Rectangle 14", "TextBox 15", "TextBox 16"]:
    sh = get_shape(s4, name)
    if sh:
        sh.top = sh.top + delta

prs.save("presentations/0616_Stablecoin_Exorbitant_Privilege.pptx")
print(f"Done — slide 4 updated. Dark bar shifted by {delta} EMU.")
