"""
update_s4title_s6.py
  - Slide 4: rename card title "The results weakened" → "Neither Result Survived"
  - Slide 6: replace dense paragraph bodies with approved bullet-point content
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

TEAL   = RGBColor(0x1B, 0x99, 0x8B)
GOLD   = RGBColor(0xE8, 0xA1, 0x2C)
NAVY   = RGBColor(0x13, 0x29, 0x4B)
NAVYDK = RGBColor(0x12, 0x28, 0x44)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DKGRAY = RGBColor(0x33, 0x42, 0x5C)
MGRAY  = RGBColor(0x5B, 0x6B, 0x82)
SANS   = "Calibri"
SERIF  = "Georgia"

ANS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def get_shape(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None


def clear_tf(tf):
    txBody = tf._txBody
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)
    txBody.append(parse_xml(f'<a:p xmlns:a="{ANS}"/>'))


def fp(tf, text, font, sz, bold=False, color=NAVY,
       align=PP_ALIGN.LEFT, sb=0, sa=0):
    p = tf.paragraphs[0]
    p.alignment = align
    if sb: p.space_before = Pt(sb)
    if sa: p.space_after  = Pt(sa)
    r = p.add_run()
    r.text = text; r.font.name = font; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.color.rgb = color


def ap(tf, text, font, sz, bold=False, color=NAVY,
       align=PP_ALIGN.LEFT, sb=0, sa=0):
    p = tf.add_paragraph()
    p.alignment = align
    if sb: p.space_before = Pt(sb)
    if sa: p.space_after  = Pt(sa)
    r = p.add_run()
    r.text = text; r.font.name = font; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.color.rgb = color


# ══════════════════════════════════════════════════════════════════════════════
prs = Presentation("presentations/0616_Stablecoin_Exorbitant_Privilege.pptx")

# ── SLIDE 4: rename card title ────────────────────────────────────────────────
s4   = prs.slides[3]
tb8  = get_shape(s4, "TextBox 8")
tf8  = tb8.text_frame
# Preserve the existing run's font by just replacing .text
for r in tf8.paragraphs[0].runs:
    r.text = "Neither Result Survived"
    break  # only the first run needed

# ── SLIDE 6: bullet rewrite ───────────────────────────────────────────────────
s6 = prs.slides[5]

CARDS = [
    dict(
        tb="TextBox 9",
        rows=[
            ("Problem",  TEAL,
             "→  The spread drifts over time (non-stationary) — slides 4–5 showed this makes regression spurious"),
            ("Fix",      TEAL,
             "→  Switch the outcome variable to the bid-cover ratio"),
            ("What is bid-cover", TEAL,
             "→  Bids received per dollar of T-bills offered at auction — a direct, stationary measure of demand strength"),
            ("Why it works", TEAL,
             "→  It fluctuates around a stable mean → no drift, no spurious regression risk"),
        ]
    ),
    dict(
        tb="TextBox 13",
        rows=[
            ("Auction size",       TEAL,
             "→  Large supply mechanically lowers bid-cover; we add offering size as a control to isolate USDT's effect"),
            ("Interpolated data",  TEAL,
             "→  Reserve ratios were quarterly, filled in monthly; we re-ran with all interpolated values removed"),
            ("Fed cycle",          TEAL,
             "→  USDT and T-bill demand both move with rates; we added Δfedfunds as a control"),
            ("Placebo",            GOLD,
             "→  Shuffled USDT supply 2,000 times at random — the real result ranks above 95%+ of random draws"),
        ]
    ),
    dict(
        tb="TextBox 17",
        rows=[
            ("Independent replication", TEAL,
             "→  Each team member reproduced the full analysis independently"),
            ("Reconciliation",          TEAL,
             "→  Results aligned across all specifications and all members"),
            ("Verdict",                 GOLD,
             "→  The bid-cover result survived every check → it becomes the paper's lead finding"),
        ]
    ),
]

for card in CARDS:
    sh = get_shape(s6, card["tb"])
    sh.text_frame.word_wrap = True
    tf = sh.text_frame
    clear_tf(tf)

    first = True
    for label, label_color, body in card["rows"]:
        if first:
            fp(tf, label, SANS, 10.5, bold=True, color=label_color, sa=3)
            first = False
        else:
            ap(tf, label, SANS, 10.5, bold=True, color=label_color, sb=11, sa=3)
        ap(tf, body, SANS, 10, color=DKGRAY)

prs.save("presentations/0616_Stablecoin_Exorbitant_Privilege.pptx")
print("Done — slide 4 title updated, slide 6 bullets written.")
