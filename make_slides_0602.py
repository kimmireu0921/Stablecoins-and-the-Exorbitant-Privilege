"""
make_slides_0602.py — June 2026 update deck with graphs embedded.

Slides:
  1.  Title
  2.  Roadmap
  3.  Event Study Update — corrected CARs insignificant (1 slide)
  4.  Threshold Robustness I — TRIM Sensitivity
  5.  Threshold Robustness II — Smooth Transition Regression (LSTAR)
  6.  What the Paper Proves (β₁ + q* + robustness summary)
  7.  Thank You

Saves to: presentations/0602_Stablecoin_Exorbitant_Privilege.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

RESULTS       = Path("results")
PRESENTATIONS = Path("presentations")
PRESENTATIONS.mkdir(exist_ok=True)

# ── Colours ───────────────────────────────────────────────────────────────────
DARK_NAVY  = RGBColor(0x0A, 0x1E, 0x5C)
MED_NAVY   = RGBColor(0x14, 0x30, 0x70)
HDR_NAVY   = RGBColor(0x1B, 0x3A, 0x72)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE     = RGBColor(0xD4, 0x63, 0x2A)
BODY_DARK  = RGBColor(0x1A, 0x1A, 0x2E)
LGRAY      = RGBColor(0xF0, 0xF4, 0xF8)
MIDGRAY    = RGBColor(0x55, 0x55, 0x66)
GREEN      = RGBColor(0x1A, 0x7A, 0x4A)
RED        = RGBColor(0xC0, 0x39, 0x2B)
LIGHT_RED  = RGBColor(0xFC, 0xEB, 0xE8)
LIGHT_GRN  = RGBColor(0xE6, 0xF4, 0xEA)
GOLD       = RGBColor(0xFF, 0xD7, 0x00)
AMBER      = RGBColor(0xE6, 0x7E, 0x22)
TEAL       = RGBColor(0x00, 0x8B, 0x8B)
LIGHT_TEAL = RGBColor(0xE0, 0xF5, 0xF5)
PURPLE     = RGBColor(0x5B, 0x2C, 0x8D)
LIGHT_PURP = RGBColor(0xF0, 0xE8, 0xF8)

SW = Inches(13.33)
SH = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def box(sld, l, t, w, h, fill=None, line_color=None, line_width=Pt(0)):
    sh = sld.shapes.add_shape(1, l, t, w, h)
    sh.line.fill.background()
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = line_width
    return sh


def txt(sld, text, l, t, w, h, size=18, bold=False, italic=False,
        color=BODY_DARK, align=PP_ALIGN.LEFT, font="Calibri", wrap=True):
    tb = sld.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_line(tf, text, size=18, bold=False, italic=False,
             color=BODY_DARK, align=PP_ALIGN.LEFT, space_before=4):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return p


def img(sld, path, l, t, w, h=None):
    p = Path(path)
    if p.exists():
        if h:
            sld.shapes.add_picture(str(p), l, t, w, h)
        else:
            sld.shapes.add_picture(str(p), l, t, w)
    else:
        box(sld, l, t, w, h or Inches(2), fill=LGRAY, line_color=MIDGRAY, line_width=Pt(1))
        txt(sld, f"[image not found:\n{p.name}]", l + Inches(0.1), t + Inches(0.1),
            w - Inches(0.2), (h or Inches(2)) - Inches(0.2),
            size=11, italic=True, color=MIDGRAY)


def title_slide_bg(sld):
    box(sld, 0, 0, SW, SH, fill=DARK_NAVY)
    box(sld, 0, Inches(2.85), SW, Inches(1.15), fill=MED_NAVY)
    for shape_args in [
        (9, Inches(0.1),  Inches(0.0),  Inches(3.8), Inches(3.8), WHITE,     DARK_NAVY),
        (9, Inches(0.55), Inches(0.45), Inches(2.9), Inches(2.9), DARK_NAVY, None),
        (9, Inches(0.1),  Inches(4.4),  Inches(3.6), Inches(3.6), WHITE,     DARK_NAVY),
        (9, Inches(0.65), Inches(4.85), Inches(2.5), Inches(2.5), DARK_NAVY, None),
    ]:
        shape_type, sl, st, sw, sh_, fill_col, _ = shape_args
        sh = sld.shapes.add_shape(shape_type, sl, st, sw, sh_)
        sh.fill.solid(); sh.fill.fore_color.rgb = fill_col; sh.line.fill.background()
    box(sld, 0, Inches(2.55), Inches(3.8), Inches(0.6), fill=MIDGRAY)


def content_slide(prs, title_text):
    sld = blank(prs)
    box(sld, 0, 0, SW, SH, fill=WHITE)
    box(sld, 0, 0, SW, Inches(0.95), fill=HDR_NAVY)
    txt(sld, title_text,
        Inches(0.3), Inches(0.1), Inches(11.5), Inches(0.78),
        size=26, color=WHITE, align=PP_ALIGN.LEFT, font="Calibri")
    return sld


def slide_number(sld, num):
    txt(sld, str(num),
        SW - Inches(0.55), SH - Inches(0.42), Inches(0.45), Inches(0.35),
        size=13, color=MIDGRAY, align=PP_ALIGN.RIGHT)


def orange_callout(sld, text, t=None, size=15):
    if t is None:
        t = SH - Inches(1.22)
    box(sld, 0, t, SW, Inches(1.05), fill=ORANGE)
    txt(sld, text,
        Inches(0.6), t + Inches(0.12), SW - Inches(1.0), Inches(0.82),
        size=size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def bullet_tb(sld, l, t, w, h):
    tb = sld.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].space_before = Pt(0)
    return tf


def notes(sld, script: str):
    sld.notes_slide.notes_text_frame.text = script


# ── Build ─────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
sld = blank(prs)
title_slide_bg(sld)
txt(sld, "STABLECOINS AND THE EXORBITANT PRIVILEGE",
    Inches(4.0), Inches(2.65), Inches(9.0), Inches(0.65),
    size=22, color=WHITE, align=PP_ALIGN.LEFT, font="Calibri Light")
txt(sld, "Event Study Update  ·  Threshold Robustness  ·  LSTAR Validation",
    Inches(4.0), Inches(3.3), Inches(9.0), Inches(0.5),
    size=13, italic=True, color=RGBColor(0xAA, 0xC4, 0xE8), align=PP_ALIGN.LEFT)
box(sld, Inches(3.9), Inches(4.45), Inches(7.5), Inches(2.8), fill=RGBColor(0x0E, 0x26, 0x68))
txt(sld,
    "Yonsei GSIS 2026-1  |  Topics in International Finance\n"
    "Mireu Kim · Sara Chekroune · Oybek Ibragimov\n"
    "June 2026",
    Inches(4.1), Inches(4.6), Inches(7.2), Inches(2.5),
    size=16, color=WHITE, align=PP_ALIGN.LEFT)
notes(sld,
    "Good afternoon. Since last week's presentation we have three updates. "
    "First: we corrected and retested the event study — the CARs are insignificant. "
    "Second: we ran two robustness checks on the 13% buffer threshold — "
    "TRIM sensitivity and a Smooth Transition Regression. "
    "Both independently confirm q* is near 13%. "
    "The main finding of the paper — β₁ = −6.02 bps and q* = 13% — stands.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Roadmap
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "What We Are Covering Today")
slide_number(sld, 2)

items = [
    (RED,     "Event Study Update",
     "We retested the event study with the corrected model. All three CARs are insignificant "
     "(−15/−4/−2 bps, all n.s.). We explain why and propose how the event study still contributes."),
    (TEAL,    "Threshold Robustness I — TRIM Sensitivity",
     "q* = 13% is stable across all TRIM parameter values (15%, 20%, 25%). "
     "Bootstrap 90% CI: [2.6%, 14.5%]. Not a boundary artifact."),
    (PURPLE,  "Threshold Robustness II — Smooth Transition (LSTAR)",
     "A completely different model that allows a smooth, gradual regime shift instead of a sharp cut. "
     "Transition midpoint c* = 13.1% — virtually identical to Hansen's q* = 13.0%."),
    (GREEN,   "What the Paper Proves",
     "Three independent methods agree: the tipping point is near 13%. "
     "β₁ = −6.02 bps (regression) + q* = 13% (Hansen) + LSTAR c* = 13.1% = the paper."),
]
for i, (col, title, body) in enumerate(items):
    t = Inches(1.15) + i * Inches(1.38)
    box(sld, Inches(0.35), t, Inches(0.65), Inches(1.25), fill=col)
    txt(sld, str(i + 1), Inches(0.35), t + Inches(0.32),
        Inches(0.65), Inches(0.5), size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(sld, Inches(1.05), t, Inches(11.9), Inches(1.25), fill=LGRAY if i % 2 == 0 else WHITE)
    txt(sld, title, Inches(1.2), t + Inches(0.1),
        Inches(11.6), Inches(0.35), size=15, bold=True, color=col)
    txt(sld, body,  Inches(1.2), t + Inches(0.48),
        Inches(11.6), Inches(0.65), size=13, color=BODY_DARK)

notes(sld,
    "Four items today. "
    "One: event study update — we corrected the model since last week, all CARs come back insignificant, "
    "and we will explain why in one slide and what role the event study still plays. "
    "Two and three: two robustness checks on the 13% buffer threshold. "
    "TRIM sensitivity confirms the threshold is not a boundary artifact, "
    "and LSTAR — a smooth transition regression — independently locates the transition at 13.1%. "
    "Four: what the paper proves, now with three lines of convergent evidence for q* = 13%.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Event Study Update (1 slide: correction + insignificance + why + alternative role)
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "Event Study Update — Corrected CARs Are Insignificant")
slide_number(sld, 3)

# Update banner
box(sld, Inches(0.35), Inches(1.05), Inches(12.6), Inches(0.34), fill=AMBER)
txt(sld,
    "Update since last week: we retested with the corrected first-difference model. "
    "Original CARs (+891/+885/−1,801 bps) were inflated ~120× by a Fed hiking trend in the estimation window.",
    Inches(0.5), Inches(1.08), Inches(12.1), Inches(0.27),
    size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Left: car_comparison.png (before/after visual)
img(sld, RESULTS / "car_comparison.png",
    Inches(0.35), Inches(1.47), Inches(7.5), Inches(3.8))
txt(sld, "Top row: original (level model)  ·  Bottom row: corrected (first-diff model)",
    Inches(0.35), Inches(5.3), Inches(7.5), Inches(0.25),
    size=10, italic=True, color=MIDGRAY)

# Right: explanation + numbers + alternative role
RX = Inches(8.05)
RW = Inches(5.0)

# Why it was wrong
box(sld, RX, Inches(1.47), RW, Inches(0.3), fill=RED)
txt(sld, "Why the original was wrong",
    RX + Inches(0.08), Inches(1.5), RW - Inches(0.12), Inches(0.24),
    size=11, bold=True, color=WHITE)
box(sld, RX, Inches(1.77), RW, Inches(0.82), fill=LIGHT_RED)
tf = bullet_tb(sld, RX + Inches(0.1), Inches(1.8), RW - Inches(0.15), Inches(0.76))
add_line(tf, "Estimation window (Jan–May 2022) coincided with Fed hiking.", size=11, color=BODY_DARK, space_before=2)
add_line(tf, "Spread rose 5→78 bps — entirely Fed policy, not stablecoins.", size=11, color=RED, bold=True, space_before=4)
add_line(tf, "Daily 'abnormal' = 34 bps × 26 days = 884 bps fake inflation.", size=11, color=BODY_DARK, space_before=4)

# Corrected numbers
box(sld, RX, Inches(2.67), RW, Inches(0.3), fill=MIDGRAY)
txt(sld, "Corrected results (first-difference model)",
    RX + Inches(0.08), Inches(2.70), RW - Inches(0.12), Inches(0.24),
    size=11, bold=True, color=WHITE)
box(sld, RX, Inches(2.97), RW, Inches(0.82), fill=LGRAY)
tf2 = bullet_tb(sld, RX + Inches(0.1), Inches(3.0), RW - Inches(0.15), Inches(0.76))
add_line(tf2, "LUNA/UST  (May 2022):   −15.3 bps   n.s.", size=12, bold=True, color=MIDGRAY, space_before=2)
add_line(tf2, "USDT depeg (May 2022):   −4.3 bps   n.s.", size=12, bold=True, color=MIDGRAY, space_before=4)
add_line(tf2, "USDC/SVB  (Mar 2023):    −2.4 bps   n.s.", size=12, bold=True, color=MIDGRAY, space_before=4)

# Placebo check
box(sld, RX, Inches(3.87), RW, Inches(0.3), fill=HDR_NAVY)
txt(sld, "Placebo check",
    RX + Inches(0.08), Inches(3.90), RW - Inches(0.12), Inches(0.24),
    size=11, bold=True, color=WHITE)
box(sld, RX, Inches(4.17), RW, Inches(0.55), fill=RGBColor(0xE8, 0xEE, 0xF8))
tf3 = bullet_tb(sld, RX + Inches(0.1), Inches(4.20), RW - Inches(0.15), Inches(0.49))
add_line(tf3, "Placebo mean |CAR| = 4.8 bps   ·   Actual mean |CAR| = 7.3 bps", size=11, color=BODY_DARK, space_before=2)
add_line(tf3, "Ratio = 1.5×  →  actuals indistinguishable from placebos.", size=11, bold=True, color=HDR_NAVY, space_before=4)

# Alternative role (green)
box(sld, RX, Inches(4.80), RW, Inches(0.3), fill=GREEN)
txt(sld, "Alternative role — still useful as qualitative context",
    RX + Inches(0.08), Inches(4.83), RW - Inches(0.12), Inches(0.24),
    size=11, bold=True, color=WHITE)
box(sld, RX, Inches(5.10), RW, Inches(0.75), fill=LIGHT_GRN)
tf4 = bullet_tb(sld, RX + Inches(0.1), Inches(5.13), RW - Inches(0.15), Inches(0.69))
add_line(tf4, "Directional pattern is consistent with β₁ = −6.02 bps:", size=11, color=BODY_DARK, space_before=2)
add_line(tf4, "  low-buffer events → spread falls; high-buffer → smaller response.", size=11, color=BODY_DARK, space_before=3)
add_line(tf4, "Illustrates the mechanism — quantitative proof is in the regression.", size=11, bold=True, color=GREEN, space_before=4)

orange_callout(sld,
    "The event study is insignificant and cannot be the main evidence. "
    "The paper's proof is β₁ = −6.02 bps + q* = 13%. The event study illustrates the mechanism.",
    size=13)

notes(sld,
    "Since last week, we corrected the event study. Here is the full story on one slide. "
    "The graph on the left shows the before and after. "
    "The top row is the original model — dramatic spikes of plus 891, plus 885, and minus 1,801 basis points. "
    "The bottom row is the corrected model — all three events produce CARs near zero. "
    "Why was the original wrong? "
    "For the 2022 events, our estimation window was January through May 2022 — "
    "exactly when the Fed was hiking rates at the fastest pace in 40 years. "
    "The spread rose from 5 to 78 basis points during that window. "
    "Every single day, the model computed about 34 basis points of fake 'abnormal' spread. "
    "Over 26 event days, that accumulated to 884 basis points — "
    "none of which was caused by stablecoins. "
    "The corrected first-difference model removes this trend by modelling daily changes, not levels. "
    "With the correction: LUNA minus 15.3, USDT minus 4.3, SVB minus 2.4 — all insignificant. "
    "The placebo test confirms this: placebos and actuals have essentially the same mean absolute CAR. "
    "A ratio of 1.5 times is not meaningful. "
    "However, the event study still contributes qualitatively. "
    "The directional pattern — low-buffer events show a larger spread response — "
    "is consistent with what the regression tells us. "
    "We present it as illustrative context, not as proof.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Threshold Robustness I: TRIM Sensitivity
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "Threshold Robustness I — TRIM Sensitivity")
slide_number(sld, 4)

# What TRIM is
box(sld, Inches(0.35), Inches(1.05), Inches(12.6), Inches(0.38), fill=TEAL)
txt(sld,
    "TRIM = the fraction of the threshold variable excluded from the grid search edges "
    "(Hansen recommends 15%). "
    "We tested 15%, 20%, and 25% to check that q* = 13% is not a boundary artifact.",
    Inches(0.5), Inches(1.09), Inches(12.1), Inches(0.3),
    size=12, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

# Left: trim sensitivity plot
img(sld, RESULTS / "threshold_trim_sensitivity.png",
    Inches(0.35), Inches(1.52), Inches(7.4), Inches(3.8))

# Right: results table + bootstrap CI
RX = Inches(7.95)
RW = Inches(5.1)

box(sld, RX, Inches(1.52), RW, Inches(0.32), fill=TEAL)
txt(sld, "TRIM stability results",
    RX + Inches(0.1), Inches(1.55), RW - Inches(0.15), Inches(0.25),
    size=13, bold=True, color=WHITE)

# Mini table
table_data = [
    ("TRIM", "Candidates", "q*",        "Stable?"),
    ("15%",  "31",         "0.1301",    "✓"),
    ("20%",  "28",         "0.1301",    "✓"),
    ("25%",  "24",         "0.1301",    "✓"),
]
col_ls_t = [RX + Inches(0.1), RX + Inches(1.4), RX + Inches(2.85), RX + Inches(4.2)]
col_ws_t = [Inches(1.3), Inches(1.45), Inches(1.35), Inches(0.75)]
for row_i, row_data in enumerate(table_data):
    row_t = Inches(1.84) + row_i * Inches(0.42)
    is_hdr = (row_i == 0)
    bg = MED_NAVY if is_hdr else (LGRAY if row_i % 2 == 1 else WHITE)
    for val, cl, cw in zip(row_data, col_ls_t, col_ws_t):
        box(sld, cl, row_t, cw, Inches(0.38), fill=bg)
        txt(sld, val, cl + Inches(0.05), row_t + Inches(0.1),
            cw - Inches(0.08), Inches(0.25),
            size=12, bold=is_hdr, color=WHITE if is_hdr else BODY_DARK,
            align=PP_ALIGN.CENTER)

# Bootstrap CI
box(sld, RX, Inches(3.57), RW, Inches(0.32), fill=TEAL)
txt(sld, "Bootstrap 90% CI for q*  (B = 1,000 replications)",
    RX + Inches(0.1), Inches(3.60), RW - Inches(0.15), Inches(0.25),
    size=12, bold=True, color=WHITE)
box(sld, RX, Inches(3.89), RW, Inches(0.55), fill=LIGHT_TEAL)
tf = bullet_tb(sld, RX + Inches(0.1), Inches(3.92), RW - Inches(0.15), Inches(0.49))
add_line(tf, "CI = [2.6%,  14.5%]      Point estimate: 13.0%", size=14, bold=True, color=TEAL, space_before=2)
add_line(tf, "Method: percentile bootstrap — resample rows, re-run grid search each time.", size=10, color=BODY_DARK, space_before=5)

# Multiple threshold result
box(sld, RX, Inches(4.52), RW, Inches(0.32), fill=HDR_NAVY)
txt(sld, "Multiple threshold test (H₀: 1 threshold vs H₁: 2 thresholds)",
    RX + Inches(0.1), Inches(4.55), RW - Inches(0.15), Inches(0.25),
    size=12, bold=True, color=WHITE)
box(sld, RX, Inches(4.84), RW, Inches(0.65), fill=LGRAY)
tf2 = bullet_tb(sld, RX + Inches(0.1), Inches(4.87), RW - Inches(0.15), Inches(0.59))
add_line(tf2, "Bootstrap p = 0.143  →  Cannot reject H₀.", size=13, bold=True, color=HDR_NAVY, space_before=2)
add_line(tf2, "Single threshold at 13% is sufficient. No evidence of a second break.", size=11, color=BODY_DARK, space_before=5)

orange_callout(sld,
    "q* = 13.0% is stable across all TRIM values. "
    "Bootstrap CI [2.6%, 14.5%] brackets it. One threshold is sufficient.",
    size=14)

notes(sld,
    "This is robustness check number one for the Hansen threshold. "
    "The TRIM parameter controls what fraction of observations near the edges of the buffer distribution "
    "are excluded from the grid search. This ensures each regime has enough observations to estimate. "
    "Hansen recommends 15 percent. We tested 15, 20, and 25 percent. "
    "In all three cases, the grid search finds exactly the same threshold: q* = 0.1301. "
    "The table on the right shows this. "
    "Below the table: we also estimated the bootstrap 90 percent confidence interval for q*. "
    "We resampled the data 1,000 times with replacement, re-ran the full grid search each time, "
    "and took the 5th and 95th percentiles of the resulting distribution. "
    "The CI is 2.6% to 14.5%, and the point estimate of 13.0% sits comfortably inside it. "
    "Finally, we tested whether the data supports two thresholds instead of one. "
    "The bootstrap p-value is 0.143 — we cannot reject the single-threshold model. "
    "One break at 13% is the right specification.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Threshold Robustness II: LSTAR
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "Threshold Robustness II — Smooth Transition Regression (LSTAR)")
slide_number(sld, 5)

# What LSTAR is
box(sld, Inches(0.35), Inches(1.05), Inches(12.6), Inches(0.38), fill=PURPLE)
txt(sld,
    "Hansen (2000) assumes a sharp binary switch at q*. "
    "LSTAR asks: what if the regime shift is gradual? "
    "It estimates a smooth logistic transition — γ controls sharpness, c controls the midpoint.",
    Inches(0.5), Inches(1.09), Inches(12.1), Inches(0.3),
    size=12, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

# Left: LSTAR transition plot
img(sld, RESULTS / "star_transition.png",
    Inches(0.35), Inches(1.52), Inches(7.5), Inches(3.9))
txt(sld,
    "Left panel: G(L) = logistic weight function — how much the low-buffer regime is active at each L value.  "
    "Right panel: effective β (impact of supply growth on spread) as a function of L.",
    Inches(0.35), Inches(5.45), Inches(7.5), Inches(0.35),
    size=10, italic=True, color=MIDGRAY)

# Right: results + comparison
RX = Inches(8.05)
RW = Inches(5.0)

box(sld, RX, Inches(1.52), RW, Inches(0.32), fill=PURPLE)
txt(sld, "LSTAR estimation results",
    RX + Inches(0.1), Inches(1.55), RW - Inches(0.12), Inches(0.25),
    size=13, bold=True, color=WHITE)
box(sld, RX, Inches(1.84), RW, Inches(1.15), fill=LIGHT_PURP)
tf = bullet_tb(sld, RX + Inches(0.1), Inches(1.87), RW - Inches(0.15), Inches(1.09))
add_line(tf, "Sharpness:  γ* = 2,768  (near-sharp transition)", size=12, bold=True, color=PURPLE, space_before=2)
add_line(tf, "Midpoint:   c* = 0.1314  (13.1%)", size=12, bold=True, color=PURPLE, space_before=5)
add_line(tf, "Bootstrap 90% CI for c*:  [2.6%,  14.5%]", size=12, color=BODY_DARK, space_before=5)
add_line(tf, "β at G=1 (low buffer):   −7.85 bps  |  β at G=0 (high):  −1.84 bps", size=11, color=BODY_DARK, space_before=6)

# Comparison vs Hansen
box(sld, RX, Inches(3.07), RW, Inches(0.32), fill=GREEN)
txt(sld, "Comparison with Hansen (2000)",
    RX + Inches(0.1), Inches(3.10), RW - Inches(0.12), Inches(0.25),
    size=13, bold=True, color=WHITE)
box(sld, RX, Inches(3.39), RW, Inches(1.05), fill=LIGHT_GRN)
tf2 = bullet_tb(sld, RX + Inches(0.1), Inches(3.42), RW - Inches(0.15), Inches(0.99))
add_line(tf2, "Hansen q*  = 0.1301  (sharp switch)", size=12, bold=True, color=HDR_NAVY, space_before=2)
add_line(tf2, "LSTAR c*   = 0.1314  (smooth midpoint)", size=12, bold=True, color=PURPLE, space_before=5)
add_line(tf2, "Difference = 0.0013  |  Hansen q* inside LSTAR CI: YES ✓", size=12, bold=True, color=GREEN, space_before=5)

# What γ* tells us
box(sld, RX, Inches(4.52), RW, Inches(0.32), fill=HDR_NAVY)
txt(sld, "What γ* = 2,768 means",
    RX + Inches(0.1), Inches(4.55), RW - Inches(0.12), Inches(0.25),
    size=12, bold=True, color=WHITE)
box(sld, RX, Inches(4.84), RW, Inches(0.75), fill=LGRAY)
tf3 = bullet_tb(sld, RX + Inches(0.1), Inches(4.87), RW - Inches(0.15), Inches(0.69))
add_line(tf3, "LSTAR tried to fit a smooth transition — and found that the data supports a near-discrete jump.", size=11, color=BODY_DARK, space_before=2)
add_line(tf3, "The model independently rediscovered the sharp switch. This validates Hansen's assumption.", size=11, bold=True, color=HDR_NAVY, space_before=5)

orange_callout(sld,
    "CONVERGENT VALIDITY: Sharp threshold (Hansen, 13.0%) and smooth transition (LSTAR, 13.1%) "
    "independently identify the same tipping point. γ* confirms the switch is near-discrete.",
    size=13)

notes(sld,
    "This is robustness check number two — a completely different model. "
    "Hansen 2000 assumes the regime switches sharply at q*. "
    "A critic could say: what if the transition is actually gradual? "
    "LSTAR, or Logistic Smooth Transition Regression, tests exactly that. "
    "Instead of a binary indicator, it uses a logistic weight function G that takes values between 0 and 1. "
    "When the buffer is far below the midpoint c, G is close to 1 — the low-buffer regime is fully active. "
    "When the buffer is far above c, G is close to 0 — the high-buffer regime. "
    "The parameter gamma controls how sharp the transition is. "
    "A small gamma means the transition is very gradual. "
    "A large gamma means the transition is essentially a sharp step function. "
    "We estimated the model using a grid search followed by Nelder-Mead optimisation. "
    "The results: the midpoint c-star is 0.1314, or 13.1 percent. "
    "The Hansen threshold was 13.0 percent. The difference is 0.0013 — negligible. "
    "And the sharpness parameter gamma-star is 2,768. "
    "This is an extremely large value — it means the transition is almost exactly a step function. "
    "LSTAR tried to find a smooth transition and concluded the data looks like a sharp switch. "
    "The left panel in the graph shows the G function — it is almost a vertical step at 13 percent. "
    "The right panel shows the effective beta as a function of the buffer. "
    "At low buffer, beta is about minus 7.9 bps. At high buffer, it is minus 1.8 bps. "
    "This is convergent validity: two completely different model architectures find the same threshold.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — What the Paper Proves
# ═══════════════════════════════════════════════════════════════════════════════
sld = content_slide(prs, "What Our Paper Proves — Three Convergent Lines of Evidence")
slide_number(sld, 6)

# Top strip
box(sld, Inches(0.35), Inches(1.05), Inches(12.6), Inches(0.5), fill=HDR_NAVY)
txt(sld,
    "The event study is insignificant and has been removed from the main evidence. "
    "The paper now rests on three independent results that all point to the same conclusion.",
    Inches(0.5), Inches(1.1), Inches(12.1), Inches(0.38),
    size=12, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

# Three evidence blocks
BLKS = [
    (GREEN,  LIGHT_GRN,
     "Finding 1 — Regression (β₁ = −6.02 bps)",
     [
         "Monthly panel OLS with Newey-West HAC standard errors (1 lag).",
         "Each σ increase in stablecoin supply → spread falls 6.02 bps.",
         "Stablecoins amplify the exorbitant privilege in normal times.",
         "Statistically significant: p < 0.01.",
     ]),
    (TEAL,   LIGHT_TEAL,
     "Finding 2 — Hansen Threshold (q* = 13%)",
     [
         "Grid search finds q* = 0.1301: below this, the privilege reverses.",
         "TRIM stable: q* = 0.1301 at TRIM = 15%, 20%, 25%.",
         "Bootstrap 90% CI: [2.6%, 14.5%] — well-defined, not a boundary artifact.",
         "One threshold is sufficient (two-threshold p = 0.143).",
     ]),
    (PURPLE, LIGHT_PURP,
     "Robustness — LSTAR (c* = 13.1%)",
     [
         "Smooth transition model relaxes sharp-switch assumption.",
         "Transition midpoint c* = 0.1314 ≈ Hansen's q* = 0.1301.",
         "Sharpness γ* = 2,768 → switch is near-discrete: validates Hansen.",
         "Hansen q* lies inside LSTAR 90% CI. Convergent validity confirmed.",
     ]),
]
for i, (hcol, bcol, title, bullets) in enumerate(BLKS):
    l = Inches(0.35) + i * Inches(4.35)
    w = Inches(4.2)
    box(sld, l, Inches(1.65), w, Inches(0.35), fill=hcol)
    txt(sld, title, l + Inches(0.1), Inches(1.68), w - Inches(0.15), Inches(0.28),
        size=12, bold=True, color=WHITE)
    box(sld, l, Inches(2.0), w, Inches(3.35), fill=bcol)
    tf = bullet_tb(sld, l + Inches(0.12), Inches(2.06), w - Inches(0.2), Inches(3.2))
    for j, btext in enumerate(bullets):
        add_line(tf, f"• {btext}", size=12, color=BODY_DARK, space_before=4 if j > 0 else 2)

# Threshold graph below Finding 2
img(sld, RESULTS / "threshold_ssr.png",
    Inches(4.7), Inches(5.42), Inches(4.2), Inches(1.45))
txt(sld, "SSR minimised at q* = 0.130 (red dashed line) — the regime switch point.",
    Inches(4.7), Inches(6.9), Inches(4.2), Inches(0.25),
    size=9, italic=True, color=MIDGRAY, align=PP_ALIGN.CENTER)

orange_callout(sld,
    "Same conclusion. Stronger evidence. "
    "β₁ = −6.02 bps  +  q* = 13.0% (Hansen)  +  c* = 13.1% (LSTAR)  =  the paper.",
    size=14)

notes(sld,
    "This is the summary slide. What does the paper actually prove? Three things. "
    "First, the regression. Each standard deviation of stablecoin supply growth compresses "
    "the OIS-Treasury spread by 6.02 basis points. "
    "Stablecoins are structural buyers of T-bills, deepening the US exorbitant privilege. "
    "This result is statistically significant, estimated on a monthly panel with robust standard errors. "
    "Second, the Hansen threshold. The optimal threshold in the data is q* = 13%. "
    "Below that buffer, the convenience yield compression reverses — stablecoins become sellers. "
    "We have confirmed this is stable across all TRIM values and that a single threshold is sufficient. "
    "Third, the LSTAR robustness check. "
    "A smooth transition model that makes no assumption about switch sharpness "
    "finds its transition midpoint at c* = 13.1% — essentially identical to Hansen. "
    "And the estimated sharpness gamma of 2,768 independently tells us the switch is near-discrete, "
    "validating Hansen's binary-regime assumption. "
    "The event study, once corrected, finds insignificant CARs and is presented as qualitative context only. "
    "The paper's conclusion is unchanged: stablecoins amplify the exorbitant privilege, "
    "up to the 13% buffer tipping point.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Thank You
# ═══════════════════════════════════════════════════════════════════════════════
sld = blank(prs)
title_slide_bg(sld)
txt(sld, "THANK YOU",
    Inches(4.0), Inches(2.85), Inches(8.8), Inches(0.6),
    size=28, color=WHITE, align=PP_ALIGN.LEFT, font="Calibri Light")
txt(sld,
    "Event study corrected (insignificant)  ·  "
    "TRIM confirms q* = 13%  ·  "
    "LSTAR validates the threshold",
    Inches(4.0), Inches(3.5), Inches(9.0), Inches(0.4),
    size=13, italic=True, color=RGBColor(0xAA, 0xC4, 0xE8), align=PP_ALIGN.LEFT)
box(sld, Inches(3.9), Inches(4.45), Inches(7.5), Inches(2.8), fill=RGBColor(0x0E, 0x26, 0x68))
txt(sld,
    "Yonsei GSIS 2026-1  |  Topics in International Finance\n"
    "Mireu Kim · Sara Chekroune · Oybek Ibragimov\n"
    "kimmireu0921@gmail.com",
    Inches(4.1), Inches(4.6), Inches(7.2), Inches(2.5),
    size=16, color=WHITE, align=PP_ALIGN.LEFT)
notes(sld,
    "Thank you. Three sentences to close. "
    "First: stablecoin issuers are now structurally important buyers of US Treasury bills. "
    "Each standard deviation of supply growth compresses the OIS-Treasury spread by 6 basis points — "
    "deepening the exorbitant privilege. "
    "Second: this amplification reverses when the issuer's liquid cash buffer falls below 13%. "
    "Third: two independent threshold models — Hansen's sharp switch and our LSTAR smooth transition — "
    "both locate this tipping point at 13%. "
    "We are happy to take questions.")


# ── Save ──────────────────────────────────────────────────────────────────────
out = PRESENTATIONS / "0602_Stablecoin_Exorbitant_Privilege.pptx"
prs.save(str(out))
print(f"Saved: {out}  ({len(prs.slides)} slides)")
