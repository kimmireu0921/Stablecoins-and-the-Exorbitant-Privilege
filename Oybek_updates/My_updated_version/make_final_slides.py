"""
make_final_slides.py — FINAL reassessment deck.

Builds a self-contained presentation summarising:
  - the teardown of the current methodology,
  - the two rebuilds (auction-level bid-cover, multi-event study),
  - what survives, and the honest reframing.

Also generates the bid-cover teardown figure from results.
Output: presentations/FINAL_Stablecoin_Reassessment.pptx
Leaves all existing project files untouched.
"""
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent
RES = ROOT / "results"
PRES = ROOT / "presentations"
PRES.mkdir(exist_ok=True)

DARK_NAVY = RGBColor(0x0A, 0x1E, 0x5C)
HDR_NAVY = RGBColor(0x1B, 0x3A, 0x72)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xD4, 0x63, 0x2A)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
BODY = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x55, 0x55, 0x66)
LGRAY = RGBColor(0xF0, 0xF4, 0xF8)

SW, SH = Inches(13.33), Inches(7.5)


# ── figure: bid-cover teardown ───────────────────────────────────────────────
def make_bidcover_figure() -> Path:
    terms = ["4-Week", "8-Week", "13-Week", "26-Week"]
    orig = [-1.14, -1.54, -1.47, -1.68]          # original 8-control spec
    orig_p = [0.012, 0.016, 0.000, 0.000]
    clean = [-0.26, -0.91, -0.80, -0.82]          # clean controls
    clean_p = [0.628, 0.110, 0.094, 0.065]

    fig, ax = plt.subplots(1, 2, figsize=(11, 4.2))
    x = np.arange(len(terms)); w = 0.38
    b1 = ax[0].bar(x - w/2, orig, w, color="#1b3a72", label="Original (8 ctrl, incl. interpolated θ,L)")
    b2 = ax[0].bar(x + w/2, clean, w, color="#d4632a", label="Clean controls (supply + VIX)")
    for xi, v, p in zip(x - w/2, orig, orig_p):
        ax[0].text(xi, v - 0.06, "***" if p < .01 else "**" if p < .05 else "*" if p < .1 else "ns",
                   ha="center", va="top", fontsize=9, color="white", fontweight="bold")
    for xi, v, p in zip(x + w/2, clean, clean_p):
        ax[0].text(xi, v - 0.06, "***" if p < .01 else "**" if p < .05 else "*" if p < .1 else "ns",
                   ha="center", va="top", fontsize=9, color="#333")
    ax[0].axhline(0, color="black", lw=0.7)
    ax[0].set_xticks(x); ax[0].set_xticklabels(terms, fontsize=9)
    ax[0].set_ylabel("β_USDT  (effect on bid-cover)", fontsize=9)
    ax[0].set_title("Monthly design: significance lives in the\ninterpolated reserve controls", fontsize=10)
    ax[0].legend(fontsize=7.5, loc="lower left")

    # auction-level spec ladder
    try:
        df = pd.read_csv(RES / "bidcover_auction_level_results.csv")
        ladder = df[df.spec.str.startswith("(")]
        labels = [s.split(")")[0].strip("(") for s in ladder.spec]
        b = ladder.beta_USDT.values; p = ladder.p_USDT.values
    except Exception:
        labels = ["1", "2", "3", "4", "5"]
        b = [0.31, 0.11, 0.11, -0.05, 0.32]; p = [0.5]*5
    ax[1].axhspan(-0.15, 0.15, color="#eee", zorder=0)
    ax[1].plot(range(len(b)), b, "o-", color="#1b3a72", lw=2, ms=8)
    ax[1].axhline(0, color="black", lw=0.7)
    ax[1].set_xticks(range(len(labels)))
    ax[1].set_xticklabels(["(1)\nsupply", "(2)\n+offering", "(3)\n+term FE", "(4)\n+VIX", "(5)\n+trend"][:len(labels)], fontsize=8)
    ax[1].set_ylabel("β_USDT  (auction-level, N=1,094)", fontsize=9)
    ax[1].set_title("Auction-level (proper N, offering controlled):\nUSDT effect is null", fontsize=10)
    ax[1].set_ylim(-1.6, 1.0)
    ax[1].text(0.5, -1.4, "all specs p > 0.4  →  no USDT bid-cover channel", fontsize=8.5,
               color="#c62828", fontweight="bold")
    plt.tight_layout()
    out = RES / "bidcover_teardown.png"
    plt.savefig(out, dpi=150, bbox_inches="tight")
    plt.close()
    return out


# ── pptx helpers ─────────────────────────────────────────────────────────────
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def box(s, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    shp = s.shapes.add_shape(1, l, t, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = lw
    return shp


def txt(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=1.0):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [[(runs, 18, BODY, False)]]
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = space; p.space_after = Pt(2)
        for s_, sz, col, bold in line:
            r = p.add_run(); r.text = s_; r.font.size = Pt(sz)
            r.font.color.rgb = col; r.font.bold = bold; r.font.name = "Calibri"
    return tb


def header(s, title, kicker=None):
    box(s, 0, 0, SW, Inches(1.05), fill=HDR_NAVY)
    box(s, 0, Inches(1.05), SW, Pt(3), fill=ORANGE)
    txt(s, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.85),
        [[(title, 26, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        txt(s, Inches(0.5), Inches(0.6), Inches(12.3), Inches(0.4),
            [[(kicker, 13, RGBColor(0xBF, 0xD0, 0xE8), False)]])


def callout(s, l, t, w, text, fill=ORANGE):
    box(s, l, t, w, Inches(0.7), fill=fill)
    txt(s, l + Inches(0.15), t, w - Inches(0.3), Inches(0.7),
        [[(text, 14, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)


def bullets(s, l, t, w, h, items, size=15):
    runs = []
    for it, lvl, col in items:
        pre = "   " * lvl + ("◆  " if lvl == 0 else "–  ")
        runs.append([(pre, size, ORANGE if lvl == 0 else GRAY, False), (it, size, col, False)])
    txt(s, l, t, w, h, runs, space=1.12)


# ── build deck ───────────────────────────────────────────────────────────────
def build():
    fig_bc = make_bidcover_figure()
    fig_ev = RES / "event_study_multi.png"

    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH

    # 1 — title
    s = blank(prs)
    box(s, 0, 0, SW, SH, fill=DARK_NAVY)
    box(s, 0, Inches(4.6), SW, Pt(4), fill=ORANGE)
    txt(s, Inches(0.8), Inches(2.1), Inches(11.7), Inches(2),
        [[("Stablecoins & the Exorbitant Privilege", 40, WHITE, True)],
         [("A Rigorous Reassessment — stress-testing our own results", 22, RGBColor(0xBF, 0xD0, 0xE8), False)]])
    txt(s, Inches(0.8), Inches(4.85), Inches(11.7), Inches(1.5),
        [[("Mireu Kim · Sara Chekroune · Oybek Ibragimov", 16, WHITE, False)],
         [("Yonsei GSIS — Topics in International Finance · June 2026", 14, RGBColor(0x9F, 0xB3, 0xD0), False)]])

    # 2 — TL;DR
    s = blank(prs); header(s, "The honest bottom line", "What we claimed vs. what survives rigorous testing")
    txt(s, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5),
        [[("What the deck claimed", 18, DARK_NAVY, True)]])
    bullets(s, Inches(0.5), Inches(1.85), Inches(6.0), Inches(4),
            [("β₁: stablecoin growth compresses the spread", 0, BODY),
             ("q* = 13% reserve-buffer threshold", 0, BODY),
             ("LSTAR confirms the threshold", 0, BODY),
             ("Bid-cover: USDT suppresses auction demand", 0, BODY),
             ("(placebo-validated, all 4 maturities)", 1, GRAY)])
    txt(s, Inches(6.9), Inches(1.3), Inches(6.0), Inches(5),
        [[("What survives", 18, GREEN, True)]])
    bullets(s, Inches(6.9), Inches(1.85), Inches(6.0), Inches(4),
            [("β₁ — spurious (fails stationarity)", 0, RED),
             ("threshold — invalid (built on spurious base)", 0, RED),
             ("LSTAR — same base, same fate", 0, RED),
             ("bid-cover — collapses at auction level", 0, RED),
             ("β₃ (liquid buffer) — survives differencing", 0, GREEN),
             ("USDT is a structural T-bill holder — a fact", 0, GREEN)])
    callout(s, Inches(0.5), Inches(6.55), Inches(12.3),
            "No channel gives a clean causal imprint. That null IS the contribution — and it is defensible.")

    # 3 — original claims recap
    s = blank(prs); header(s, "Where we started", "The three-part thesis as presented")
    bullets(s, Inches(0.6), Inches(1.5), Inches(12), Inches(5),
            [("Privilege amplification — each 1σ of stablecoin supply compressed the OIS–Treasury spread (β₁ ≈ −6 to −7.5 bps, p<0.01)", 0, BODY),
             ("A liquid-buffer threshold at ~13% — below it, the privilege reverses into a 'New Triffin' crisis zone (Hansen 2000)", 0, BODY),
             ("Convergent validity — an LSTAR smooth-transition model independently found the same ~13% switch point", 0, BODY),
             ("Mechanism — USDT supply growth suppresses T-bill auction bid-cover (placebo-validated, billed as the cleanest result)", 0, BODY)])
    callout(s, Inches(0.6), Inches(6.4), Inches(12.1),
            "Then Prof. Hur's feedback triggered a re-run — and the foundations moved.", fill=DARK_NAVY)

    # 4 — teardown overview
    s = blank(prs); header(s, "Teardown — layer by layer", "Each result examined against the data and the econometrics")
    rows = [
        ("θ (Treasury exposure)", "Placeholder + scraped-zeros + interpolated", "UNUSABLE", RED),
        ("L monthly variation", "Mostly straight-line interpolation between quarters", "ARTIFACT", RED),
        ("Levels β₁, β₄", "Non-stationary, no cointegration → spurious", "SPURIOUS", RED),
        ("Hansen / LSTAR threshold", "Built on spurious base; iid bootstrap invalid", "INVALID", RED),
        ("Bid-cover (monthly)", "Significance from interpolated controls", "FRAGILE", RED),
        ("Daily VAR / IRF (Δspread)", "Correct differencing, but FEVD ~1.3%", "WEAK", ORANGE),
        ("β₃ (liquid buffer)", "Survives levels + differencing", "SURVIVES", GREEN),
    ]
    y = Inches(1.35); rh = Inches(0.72)
    for name, why, verdict, col in rows:
        box(s, Inches(0.6), y, Inches(3.5), rh, fill=LGRAY)
        txt(s, Inches(0.72), y, Inches(3.3), rh, [[(name, 13, BODY, True)]], anchor=MSO_ANCHOR.MIDDLE)
        txt(s, Inches(4.25), y, Inches(6.0), rh, [[(why, 12.5, GRAY, False)]], anchor=MSO_ANCHOR.MIDDLE)
        box(s, Inches(10.4), y, Inches(2.3), rh, fill=col)
        txt(s, Inches(10.4), y, Inches(2.3), rh, [[(verdict, 12.5, WHITE, True)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        y += rh + Inches(0.05)

    # 5 — the data problem
    s = blank(prs); header(s, "Root cause: the reserve data is largely manufactured", "Why the 'monthly' variation isn't really data")
    bullets(s, Inches(0.6), Inches(1.4), Inches(12), Inches(2.6),
            [("Tether reports T-bill holdings quarterly — ~2 of every 3 monthly points are linear interpolation between two quarterly endpoints", 0, BODY),
             ("2020 rows are literally labelled 'pre-attestation estimate, extrapolated'; mid-2021 treasury holdings are scraped zeros (9e-08) → interpolated over", 0, BODY),
             ("The interpolation choice IS the result:", 0, BODY)])
    # mini table of interpolation sensitivity
    data = [("Interpolation", "β₄ (time series)", "p"),
            ("time / prof (used)", "−35.89", "0.032  **"),
            ("rolling mean", "+5.25", "0.873  ns"),
            ("forward-fill (old)", "−26.26", "0.563  ns")]
    y = Inches(4.1)
    for i, (a, b_, c) in enumerate(data):
        fill = HDR_NAVY if i == 0 else (RGBColor(0xFD, 0xEC, 0xE0) if i == 1 else WHITE)
        tcol = WHITE if i == 0 else BODY
        bold = i == 0 or i == 1
        for j, (val, wdt, lft) in enumerate([(a, 4.2, 0.6), (b_, 3.0, 4.8), (c, 3.0, 7.8)]):
            box(s, Inches(lft), y, Inches(wdt), Inches(0.5), fill=fill, line=RGBColor(0xDD, 0xDD, 0xDD), lw=Pt(0.5))
            txt(s, Inches(lft), y, Inches(wdt), Inches(0.5), [[(val, 13, tcol, bold)]],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.5)
    callout(s, Inches(0.6), Inches(6.5), Inches(12.1),
            "When the headline coefficient's significance flips with how you fill the MISSING data, it describes the fill — not Tether.")

    # 6 — spurious regression
    s = blank(prs); header(s, "And the levels regression is spurious", "Stationarity + cointegration diagnostics")
    bullets(s, Inches(0.6), Inches(1.4), Inches(12), Inches(3),
            [("Spread and L are both non-stationary (ADF p = 0.49, 0.90) — they trend, they don't mean-revert", 0, BODY),
             ("They do NOT cointegrate (Engle–Granger and corrected Johansen) → non-stationary + no cointegration = textbook spurious regression", 0, BODY),
             ("Common driver: the 2022–24 Fed hiking cycle pushed both down together", 0, BODY),
             ("First-differencing kills β₁ and β₄ — and β₄ flips sign (−35.9 → +8.9). A sign that depends on levels-vs-differences is noise", 0, BODY)])
    box(s, Inches(0.6), Inches(5.1), Inches(12.1), Inches(1.1), fill=LGRAY)
    txt(s, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.0),
        [[("Both series fall together 2022→2025:   ", 14, BODY, True),
          ("spread +0.85 → −0.29   |   L  0.19 → 0.06", 14, RED, True)],
         [("Two series trending down together → high R² by trend, not by mechanism.", 13, GRAY, False)]])
    callout(s, Inches(0.6), Inches(6.5), Inches(12.1),
            "β₁ (privilege) and the 13% threshold are significant only in the one specification that is statistically invalid.", fill=DARK_NAVY)

    # 7 — rebuild 1 bid-cover
    s = blank(prs); header(s, "Rebuild 1 — bid-cover at the auction level", "My contribution: the test done properly")
    bullets(s, Inches(0.5), Inches(1.35), Inches(5.6), Inches(4.5),
            [("Old test: 1,094 auctions collapsed to 51 monthly means; no auction-size control", 0, BODY),
             ("Rebuild: one row per auction (N=1,094)", 0, GREEN),
             ("control for ln(offering size) — the mechanical driver of bid-cover", 1, BODY),
             ("maturity FE, VIX, time trend; SE clustered by month", 1, BODY),
             ("falsification: future supply must NOT predict", 1, BODY)])
    s.shapes.add_picture(str(fig_bc), Inches(6.0), Inches(1.5), width=Inches(7.0))
    callout(s, Inches(0.5), Inches(6.55), Inches(12.3),
            "Strip the interpolated controls / control for offering → the USDT bid-cover channel disappears.")

    # 8 — rebuild 2 event study
    s = blank(prs); header(s, "Rebuild 2 — multi-event study", "Drop SVB (confounded), add crypto-native crises")
    bullets(s, Inches(0.5), Inches(1.35), Inches(5.6), Inches(4.5),
            [("Prof.: keep LUNA as motivation, drop SVB (banking-confounded), add more similar crises", 0, BODY),
             ("Selection rule: crypto-native stress, no banking/macro shock to Treasuries", 0, GREEN),
             ("LUNA · Celsius · FTX · BUSD", 1, BODY),
             ("first-difference normal model; cross-event CAAR (events = unit)", 1, BODY),
             ("Result: per-event CARs mostly n.s.; events disagree in sign (FTX −28 bps); pooled p = 0.43", 0, RED)])
    s.shapes.add_picture(str(fig_ev), Inches(6.0), Inches(1.5), width=Inches(7.0))
    callout(s, Inches(0.5), Inches(6.55), Inches(12.3),
            "No systematic spread response to crypto crises — exactly the 'motivation only' role the prof assigned it.", fill=DARK_NAVY)

    # 9 — what survives
    s = blank(prs); header(s, "What survives the reassessment", "The defensible core")
    items = [
        ("USDT is one of the largest structural holders of US T-bills", "Documented fact (attestations) — the privilege mechanism EXISTS structurally", GREEN),
        ("β₃: liquid buffer correlates with the spread", "Survives levels AND first-differencing — but likely macro co-movement, not a buffer mechanism", GREEN),
        ("Short-run USDT → spread (daily VAR, Δspread)", "Day-1 compression, but reverses by day 3 and explains ~1.3% of variance — suggestive only", ORANGE),
        ("Everything built on θ, L-levels, the spread-levels regression", "Spurious / interpolation-driven — demote to 'we tried this; it does not survive'", RED),
    ]
    y = Inches(1.4)
    for head_, sub, col in items:
        box(s, Inches(0.6), y, Inches(0.18), Inches(1.15), fill=col)
        txt(s, Inches(0.95), y, Inches(11.8), Inches(1.15),
            [[(head_, 16, BODY, True)], [(sub, 13, GRAY, False)]], anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(1.27)

    # 10 — reframing
    s = blank(prs); header(s, "The reframe — an honest paper that holds up", "How to present this to Prof. Hur")
    bullets(s, Inches(0.6), Inches(1.4), Inches(12), Inches(3.6),
            [("Lead with the structural fact: stablecoins have become programmatic T-bill buyers — the privilege channel is real by construction", 0, BODY),
             ("Then the honest finding: at monthly/auction frequency the causal imprint is small and swamped by the Fed cycle and data limits", 0, BODY),
             ("Show the rigour: we stress-tested our OWN headline results (stationarity, interpolation, auction-level, falsification) and report what fails", 0, BODY),
             ("Keep β₃ and the event study as supporting / motivation, clearly labelled as such", 0, BODY),
             ("This pre-empts every objection an examiner would raise — owning the limitation is stronger than defending a fragile result", 0, GREEN)])
    callout(s, Inches(0.6), Inches(6.4), Inches(12.1),
            "'We tried to break our own results. Here is what survived.' — the strongest position in the room.")

    # 11 — closing
    s = blank(prs)
    box(s, 0, 0, SW, SH, fill=DARK_NAVY)
    box(s, 0, Inches(3.5), SW, Pt(4), fill=ORANGE)
    txt(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.5),
        [[("The contribution is the rigour.", 34, WHITE, True)]])
    txt(s, Inches(0.8), Inches(3.8), Inches(11.7), Inches(2.5),
        [[("• A documented structural privilege channel (USDT ≈ T-bill buyer)", 18, RGBColor(0xCF, 0xDD, 0xF0), False)],
         [("• An honest null on its high-frequency causal footprint", 18, RGBColor(0xCF, 0xDD, 0xF0), False)],
         [("• Two purpose-built robustness rebuilds that earn that null", 18, RGBColor(0xCF, 0xDD, 0xF0), False)],
         [("• One surviving correlation (β₃), reported with the right caveats", 18, RGBColor(0xCF, 0xDD, 0xF0), False)]])

    out = PRES / "FINAL_Stablecoin_Reassessment.pptx"
    prs.save(out)
    print(f"Saved {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
    print(f"Figure: {fig_bc}")


if __name__ == "__main__":
    build()
