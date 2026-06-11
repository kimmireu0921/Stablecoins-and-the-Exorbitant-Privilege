"""
fix_slide6.py — corrects inaccurate text on slide 6 of the presentation.

Changes:
  - Banner: remove "L ≈ 18–22% cash", replace with "$3.3B (8% of reserves) frozen at SVB"
  - Force A label: "Circle's behavior (HIGH BUFFER)" → "Government backstop (SVB deposit guarantee)"
  - Force A body: remove "large cash buffer" narrative, replace with government intervention
"""

from pptx import Presentation
from pathlib import Path

PATH = Path("presentations/0526_Stablecoin_Exorbitant_Privilege.pptx")
prs = Presentation(str(PATH))
slide = prs.slides[5]  # slide 6

REPLACEMENTS = {
    # Shape 4 — top banner
    "Buffer Condition: HIGH   |   L ≈ 18–22% cash   |   CAR = −18.01 pp***   |   t-stat > 30":
        "Buffer Condition: HIGH   |   $3.3B (8% of reserves) frozen at SVB   |   CAR = −18.01 pp***   |   t-stat > 30",

    # Shape 9 — Force A label
    "Force A: Circle's behavior (HIGH BUFFER)":
        "Force A: Government backstop (SVB deposit guarantee)",

    # Shape 8 — Force A body
    "Circle used its large cash buffer to meet redemptions. Did NOT need to sell T-bills\n"
    "No forced selling pressure from the issuer side. → This is what a flat CAR would look like IF nothing else happened":
        "U.S. regulators guaranteed ALL SVB deposits within 72 hours of collapse.\n"
        "Circle did NOT need to sell T-bills — the government backstop resolved the crisis before forced liquidation. "
        "→ This is what a flat CAR would look like IF nothing else happened",
}


def replace_text(shape, old, new):
    tf = shape.text_frame
    full = tf.text
    if old not in full:
        return False
    # Replace across runs in the first paragraph where possible,
    # otherwise replace paragraph by paragraph
    for para in tf.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
                return True
    # Multi-paragraph replacement: rebuild first paragraph's run
    if tf.paragraphs:
        first_run = tf.paragraphs[0].runs
        if first_run:
            combined = tf.text
            if old in combined:
                first_run[0].text = new
                # Clear remaining runs and paragraphs
                for run in tf.paragraphs[0].runs[1:]:
                    run.text = ""
                return True
    return False


for i, shape in enumerate(slide.shapes):
    if not shape.has_text_frame:
        continue
    current = shape.text_frame.text
    for old, new in REPLACEMENTS.items():
        if old in current:
            # Direct run-level replacement
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if old in run.text:
                        run.text = run.text.replace(old, new)
                        print(f"  Shape {i}: replaced text")
                        break
                # If the old text spans multiple runs, reconstruct
                para_text = "".join(r.text for r in para.runs)
                if old in para_text and para.runs:
                    para.runs[0].text = para_text.replace(old, new)
                    for run in para.runs[1:]:
                        run.text = ""
                    print(f"  Shape {i}: reconstructed paragraph text")

prs.save(str(PATH))
print(f"\nSaved: {PATH}")

# Verify
prs2 = Presentation(str(PATH))
slide2 = prs2.slides[5]
print("\n=== SLIDE 6 AFTER FIX ===")
for shape in slide2.shapes:
    if shape.has_text_frame:
        t = shape.text_frame.text.strip()
        if t:
            print(t[:150])
            print()
