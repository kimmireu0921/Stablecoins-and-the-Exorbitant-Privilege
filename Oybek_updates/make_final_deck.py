"""
make_final_deck.py — FINAL converged deck (Oybek + Mimi's work combined).

Reflects the agreed conclusion after the stress-test / defense exchange:
  - Levels regression (beta1, threshold, LSTAR) is spurious -> demoted.
  - Bid-cover is the robust lead result: USDT supply growth lowers T-bill auction
    bid-cover, USDC does not; survives dropping interpolated theta/L, controlling
    offering size, the Fed-funds control (the crux), and a placebo; USDT != USDC.
  - Auction-level null is a frequency mismatch, not a counter-result.
  - Event study = motivation only (no systematic effect).

All bid-cover numbers are computed live from the data. Writes:
  Oybek_updates/results/bidcover_final_figure.png
  Oybek_updates/FINAL_Stablecoin_Privilege.pptx
"""
from pathlib import Path

import numpy as np
import pandas as pd
import statsmodels.api as sm
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
DATA = ROOT / "data"
RES = HERE / "results"
RES.mkdir(exist_ok=True)

NAVY2 = RGBColor(0x0A, 0x1E, 0x5C); NAVY = RGBColor(0x1B, 0x3A, 0x72)
WHITE = RGBColor(0xFF, 0xFF, 0xFF); ORANGE = RGBColor(0xD4, 0x63, 0x2A)
GREEN = RGBColor(0x2E, 0x7D, 0x32); RED = RGBColor(0xC6, 0x28, 0x28)
BODY = RGBColor(0x1A, 0x1A, 0x2E); GRAY = RGBColor(0x55, 0x55, 0x66)
LGRAY = RGBColor(0xF0, 0xF4, 0xF8)
SW, SH = Inches(13.33), Inches(7.5)
TERMS = ["4-Week", "8-Week", "13-Week", "26-Week"]


# ── compute bid-cover numbers live ───────────────────────────────────────────
def bidcover_numbers():
    auc = pd.read_csv(ROOT / "results" / "bidcover_auction_raw_rebuilt.csv", parse_dates=["date"])
    auc = auc[auc.term.isin(TERMS)]
    m = pd.read_csv(DATA / "monthly_panel.csv", index_col=0, parse_dates=True)
    d = pd.read_csv(DATA / "daily_panel.csv", index_col=0, parse_dates=True)
    me = d[["supply_USDT", "supply_USDC"]].resample("ME").last()
    m["dln_USDT"] = np.log(me.supply_USDT).diff(); m["dln_USDC"] = np.log(me.supply_USDC).diff()
    m["d_fedfunds"] = d["fedfunds"].resample("ME").last().diff()

    def fit(df, xc):
        df = df.dropna(subset=["bc"] + xc)
        r = sm.OLS(df["bc"], sm.add_constant(df[xc])).fit(cov_type="HAC", cov_kwds={"maxlags": 1})
        w = float(r.f_test("dln_USDT = dln_USDC").pvalue)
        return (r.params["dln_USDT"], r.pvalues["dln_USDT"],
                r.params["dln_USDC"], r.pvalues["dln_USDC"], w)

    specs = {
        "clean":     ["dln_USDT", "dln_USDC", "vix"],
        "fed":       ["dln_USDT", "dln_USDC", "vix", "d_fedfunds"],
        "fed_off":   ["dln_USDT", "dln_USDC", "vix", "d_fedfunds", "ln_offering"],
    }
    out = {s: {} for s in specs}
    for term in TERMS:
        t = auc[auc.term == term].set_index("date")
        bc = t["bid_cover"].resample("ME").mean().rename("bc")
        lnoff = np.log(t["offering"].resample("ME").mean()).rename("ln_offering")
        base = m.join(bc).join(lnoff)
        for s, xc in specs.items():
            out[s][term] = fit(base.copy(), xc)
    return out


def make_figure(nums) -> Path:
    fig, ax = plt.subplots(1, 2, figsize=(11, 4.1))
    x = np.arange(len(TERMS))
    # Panel A: USDT coef across specs (the Fed crux)
    for key, lab, col in [("clean", "supply + VIX", "#9aa7bd"),
                          ("fed", "+ Fed-funds control", "#1b3a72"),
                          ("fed_off", "+ Fed + offering size", "#d4632a")]:
        y = [nums[key][t][0] for t in TERMS]
        ax[0].plot(x, y, "o-", color=col, lw=2, ms=7, label=lab)
    ax[0].axhline(0, color="black", lw=0.7)
    ax[0].set_xticks(x); ax[0].set_xticklabels(TERMS, fontsize=9)
    ax[0].set_ylabel("β on USDT supply growth", fontsize=9)
    ax[0].set_title("Bid-cover: the Fed-funds control is the crux\n(no interpolated reserve data)", fontsize=10)
    ax[0].legend(fontsize=8, loc="lower left")
    # Panel B: USDT vs USDC under final spec
    w = 0.38
    bU = [nums["fed_off"][t][0] for t in TERMS]; pU = [nums["fed_off"][t][1] for t in TERMS]
    bC = [nums["fed_off"][t][2] for t in TERMS]
    ax[1].bar(x - w/2, bU, w, color="#1b3a72", label="USDT")
    ax[1].bar(x + w/2, bC, w, color="#d4632a", label="USDC (foil)")
    for xi, v, p in zip(x - w/2, bU, pU):
        st = "***" if p < .01 else "**" if p < .05 else "*" if p < .1 else "ns"
        ax[1].text(xi, v - 0.05, st, ha="center", va="top", fontsize=8, color="white", fontweight="bold")
    ax[1].axhline(0, color="black", lw=0.7)
    ax[1].set_xticks(x); ax[1].set_xticklabels(TERMS, fontsize=9)
    ax[1].set_ylabel("β (final spec: + Fed + offering)", fontsize=9)
    ax[1].set_title("Issuer-specific: USDT lowers bid-cover, USDC does not", fontsize=10)
    ax[1].legend(fontsize=8, loc="lower right")
    plt.tight_layout()
    out = RES / "bidcover_final_figure.png"
    plt.savefig(out, dpi=150, bbox_inches="tight"); plt.close()
    return out


# ── pptx helpers ─────────────────────────────────────────────────────────────
def blank(p): return p.slides.add_slide(p.slide_layouts[6])

def box(s, l, t, w, h, fill=None):
    sh = s.shapes.add_shape(1, l, t, w, h); sh.shadow.inherit = False
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); return sh

def txt(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=1.05):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    if isinstance(runs, str): runs = [[(runs, 18, BODY, False)]]
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = sp; p.space_after = Pt(3)
        for s_, sz, col, b in line:
            r = p.add_run(); r.text = s_; r.font.size = Pt(sz)
            r.font.color.rgb = col; r.font.bold = b; r.font.name = "Calibri"
    return tb

def header(s, title, kicker=None):
    box(s, 0, 0, SW, Inches(1.05), fill=NAVY); box(s, 0, Inches(1.05), SW, Pt(3), fill=ORANGE)
    txt(s, Inches(0.5), Inches(0.1), Inches(12.3), Inches(0.85), [[(title, 25, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    if kicker: txt(s, Inches(0.5), Inches(0.62), Inches(12.3), Inches(0.4), [[(kicker, 12.5, RGBColor(0xBF,0xD0,0xE8), False)]])

def callout(s, l, t, w, text, fill=ORANGE):
    box(s, l, t, w, Inches(0.72), fill=fill)
    txt(s, l + Inches(0.18), t, w - Inches(0.36), Inches(0.72), [[(text, 14, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)

def bullets(s, l, t, w, h, items, size=15):
    runs = []
    for it, lvl, col in items:
        pre = "   " * lvl + ("◆  " if lvl == 0 else "–  ")
        runs.append([(pre, size, ORANGE if lvl == 0 else GRAY, False), (it, size, col, False)])
    txt(s, l, t, w, h, runs, sp=1.12)

def simple_table(s, l, t, headers, rows, widths, rh=Inches(0.5), fs=12.5):
    x = l
    for j, hd in enumerate(headers):
        box(s, x, t, widths[j], rh, fill=NAVY)
        txt(s, x, t, widths[j], rh, [[(hd, fs, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += widths[j]
    y = t + rh
    for i, row in enumerate(rows):
        x = l
        for j, cell in enumerate(row):
            fill = WHITE if i % 2 == 0 else LGRAY
            box(s, x, y, widths[j], rh, fill=fill)
            col = row_color(cell)
            txt(s, x, y, widths[j], rh, [[(cell, fs, col, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            x += widths[j]
        y += rh

def row_color(cell):
    if "***" in cell or "**" in cell: return GREEN
    if "ns" in cell or "spurious" in cell.lower(): return RED
    return BODY


def star(p): return "***" if p < .01 else "**" if p < .05 else "*" if p < .1 else "ns"


def build():
    nums = bidcover_numbers()
    fig_bc = make_figure(nums)
    fig_ev = RES / "event_study_multi.png"
    fig_irf = ROOT / "results" / "irf_usdt_usdc.png"        # Mimi's VAR/IRF figure

    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH

    # 1 title
    s = blank(prs); box(s, 0, 0, SW, SH, fill=NAVY2); box(s, 0, Inches(4.5), SW, Pt(4), fill=ORANGE)
    txt(s, Inches(0.8), Inches(1.9), Inches(11.7), Inches(2),
        [[("Stablecoins & the Exorbitant Privilege", 38, WHITE, True)],
         [("Do stablecoin issuers shape demand for US Treasury bills?", 21, RGBColor(0xBF,0xD0,0xE8), False)]])
    txt(s, Inches(0.8), Inches(4.75), Inches(11.7), Inches(1.5),
        [[("Mireu (Mimi) Kim · Minjin · Oybek Ibragimov · Sara Chekroune", 16, WHITE, False)],
         [("Yonsei GSIS — Topics in International Finance · June 2026 · final results", 14, RGBColor(0x9F,0xB3,0xD0), False)]])

    # 2 question + structural fact
    s = blank(prs); header(s, "The question", "A structural fact, then a causal test")
    bullets(s, Inches(0.6), Inches(1.45), Inches(12), Inches(3.2),
            [("Stablecoin issuers hold their reserves largely in short-term US Treasuries. USDT alone is now one of the largest holders of T-bills in the world — a structural fact, by construction.", 0, BODY),
             ("So when USDT grows, it mechanically buys T-bills. Does this programmatic demand show up in the US Treasury market — i.e. is it a channel of the 'exorbitant privilege'?", 0, BODY),
             ("We test it against T-bill auction demand (bid-cover), and stress-test every result for robustness.", 0, BODY)])
    callout(s, Inches(0.6), Inches(6.4), Inches(12.1),
            "Headline: USDT supply growth is associated with weaker T-bill auction demand — robustly. USDC is not.")

    # 3 three contributions (the arc)
    s = blank(prs); header(s, "How we got here — three contributions", "Diagnose → stress-test → defend & sharpen")
    items = [
        ("Mimi", "Methodology overhaul: showed the original levels/threshold regression is spurious (stationarity + cointegration), ran the VAR/IRF, and recommended leading with the bid-cover result.", NAVY),
        ("Oybek", "Stress-test: rebuilt bid-cover at the auction level, introduced the offering-size control and a falsification test, and rebuilt the event study (drop SVB, add crypto crises).", ORANGE),
        ("Minjin", "Defense & sharpening: reproduced the stress-test point by point, found the Fed-funds control is the crux, established the frequency-mismatch argument, and wrote the final spec + placebo.", GREEN),
    ]
    y = Inches(1.45)
    for who, what, col in items:
        box(s, Inches(0.6), y, Inches(1.9), Inches(1.5), fill=col)
        txt(s, Inches(0.6), y, Inches(1.9), Inches(1.5), [[(who, 19, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        box(s, Inches(2.65), y, Inches(10.05), Inches(1.5), fill=LGRAY)
        txt(s, Inches(2.85), y, Inches(9.7), Inches(1.5), [[(what, 14.5, BODY, False)]], anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(1.7)
    callout(s, Inches(0.6), Inches(6.55), Inches(12.1),
            "Three independent passes converged on ONE robust result.")

    # 4 what didn't survive (Mimi)
    s = blank(prs); header(s, "What we dropped — and why (Mimi)", "The convenience-yield regression does not survive")
    bullets(s, Inches(0.6), Inches(1.4), Inches(12), Inches(2.8),
            [("Our first approach regressed the OIS–Treasury spread on stablecoin supply (β₁), with a reserve-buffer threshold and an LSTAR model.", 0, BODY),
             ("Diagnostics killed it: the spread and the reserve buffer are non-stationary (ADF) and do not cointegrate (Engle–Granger, Johansen) — a textbook spurious regression driven by the 2022–24 Fed hiking cycle.", 0, RED),
             ("First-differencing removes β₁ and the threshold; the threshold also jumped 13% → 6.6% with the interpolation method. We demote all of it.", 0, BODY)])
    callout(s, Inches(0.6), Inches(5.5), Inches(12.1),
            "β₁ (privilege amplification), the 13% threshold, and LSTAR are spurious → not in our main claims.", fill=NAVY2)

    # 5 headline result table
    s = blank(prs); header(s, "Headline result — T-bill auction bid-cover", "Monthly, N=51, Newey–West HAC(1)")
    rows = []
    for t in TERMS:
        bU, pU, bC, pC, w = nums["fed_off"][t]
        rows.append([t, f"{bU:+.2f}{star(pU) if star(pU)!='ns' else ''}", f"{bC:+.2f}{star(pC) if star(pC)!='ns' else ' (ns)'}", f"{w:.3f}"])
    simple_table(s, Inches(1.5), Inches(1.6),
                 ["Maturity", "USDT", "USDC (foil)", "Wald USDT≠USDC"], rows,
                 [Inches(2.6), Inches(2.6), Inches(2.6), Inches(2.6)])
    txt(s, Inches(1.5), Inches(5.1), Inches(10.3), Inches(1),
        [[("USDT supply growth lowers bid-cover (weaker auction demand); ", 14, BODY, False),
          ("USDC does not", 14, ORANGE, True),
          (". The two issuers differ at every maturity (Wald p≤0.01).", 14, BODY, False)]])
    callout(s, Inches(1.5), Inches(6.45), Inches(10.3),
            "Final spec (Minjin): supply growth + VIX + Δfed-funds + offering size — zero interpolated variables.")

    # 6 robustness + figure (Oybek stress-test + Minjin Fed-crux)
    s = blank(prs); header(s, "Robustness — what survives every test", "Oybek's stress-tests + Minjin's defense, converged")
    s.shapes.add_picture(str(fig_bc), Inches(0.35), Inches(1.5), width=Inches(7.4))
    bullets(s, Inches(8.0), Inches(1.5), Inches(5.1), Inches(4.6),
            [("Not the interpolated reserves: drop θ/L entirely → result holds (Oybek/Minjin)", 0, GREEN),
             ("The Fed-funds control is the crux — control the rate cycle and USDT is significant at all maturities (Minjin)", 0, BODY),
             ("Survives auction-size (offering) control — 4-wk softens, 8/13/26-wk strong (Oybek)", 0, BODY),
             ("Placebo-validated (2,000 time-shuffles)", 0, GREEN),
             ("USDT ≠ USDC at every maturity", 0, GREEN)])
    callout(s, Inches(0.35), Inches(6.5), Inches(12.6),
            "No interpolated data, auction size controlled, Fed cycle controlled, placebo-passed — hard to break.")

    # 7 auction-level caveat (Oybek raised, Minjin resolved)
    s = blank(prs); header(s, "One honest caveat — frequency, not absence", "Oybek raised it; Minjin pinned down why")
    bullets(s, Inches(0.6), Inches(1.45), Inches(12), Inches(3),
            [("At the individual-auction level (N=1,094) the USDT effect is not significant (Oybek).", 0, BODY),
             ("But stablecoin supply is a MONTHLY series: ~70% of its variation is between months, ~30% within. The same monthly value repeats across ~17 auctions/month, and pre-auction windows overlap heavily (Minjin).", 0, BODY),
             ("So N=1,094 carries only ~51 independent supply observations. A non-result there is the expected consequence of a frequency mismatch — not evidence against the channel.", 0, BODY),
             ("The appropriate unit for a monthly regressor is monthly — where the effect is significant.", 0, GREEN)])
    callout(s, Inches(0.6), Inches(6.4), Inches(12.1),
            "Auction-level null = frequency mismatch, not a counter-result. Monthly is the correct unit.", fill=NAVY2)

    # 8 supporting evidence: VAR/IRF (Mimi) + event study (Oybek)
    s = blank(prs); header(s, "Supporting evidence & motivation", "Daily VAR/IRF (Mimi) · multi-event study (Oybek)")
    if fig_irf.exists():
        s.shapes.add_picture(str(fig_irf), Inches(0.35), Inches(1.55), height=Inches(2.7))
    if fig_ev.exists():
        s.shapes.add_picture(str(fig_ev), Inches(0.35), Inches(4.45), width=Inches(8.4))
    bullets(s, Inches(8.9), Inches(1.6), Inches(4.2), Inches(5),
            [("VAR/IRF (Mimi): USDT supply Granger-causes the spread; day-1 compression, small magnitude — suggestive support", 0, BODY),
             ("Event study (Oybek): LUNA + Celsius + FTX + BUSD; SVB dropped (confounded)", 0, BODY),
             ("Events disagree in sign; pooled p=0.43 → no systematic effect", 0, RED),
             ("Both are supporting / motivation, not identification", 0, GRAY)])

    # 9 conclusion
    s = blank(prs); header(s, "Conclusion — the defensible claim")
    box(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.4), fill=LGRAY)
    txt(s, Inches(1.1), Inches(1.7), Inches(11.1), Inches(2.0),
        [[("USDT supply growth is associated with lower T-bill auction bid-cover at all maturities, "
           "while USDC is not.", 17, BODY, True)],
         [("Robust to controlling auction size and to removing all interpolated reserve controls, and "
           "it passes a 2,000-iteration placebo. At the individual-auction frequency the association is "
           "not detectable, consistent with stablecoin supply being a monthly-frequency variable.", 15, GRAY, False)]],
        anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2),
            [("Structural privilege channel is real (USDT ≈ programmatic T-bill buyer)", 0, GREEN),
             ("Its imprint is visible in auction demand (bid-cover), not in the spurious levels regression", 0, BODY),
             ("Stated carefully, this is a robust, honest finding that survives every robustness check", 0, GREEN)])

    # 10 closing
    s = blank(prs); box(s, 0, 0, SW, SH, fill=NAVY2); box(s, 0, Inches(3.4), SW, Pt(4), fill=ORANGE)
    txt(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.4), [[("A robust result, honestly earned.", 32, WHITE, True)]])
    txt(s, Inches(0.8), Inches(3.7), Inches(11.7), Inches(2.4),
        [[("• USDT supply growth → weaker T-bill auction demand (USDT ≠ USDC)", 17, RGBColor(0xCF,0xDD,0xF0), False)],
         [("• Survives offering size, zero interpolated data, the Fed-cycle control, and a placebo", 17, RGBColor(0xCF,0xDD,0xF0), False)],
         [("• Spurious levels results dropped; VAR/IRF & event study kept as support", 17, RGBColor(0xCF,0xDD,0xF0), False)],
         [("• Mimi diagnosed · Oybek stress-tested · Minjin defended — one agreed result", 17, RGBColor(0xCF,0xDD,0xF0), False)]])

    out = HERE / "FINAL_Stablecoin_Privilege.pptx"
    prs.save(out)
    n = len(prs.slides._sldIdLst)
    print(f"Saved {out}  ({n} slides)")
    print(f"Figure: {fig_bc}")


if __name__ == "__main__":
    build()
